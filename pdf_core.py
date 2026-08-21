from __future__ import annotations

import ctypes
import io
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import unicodedata
import zipfile
from collections.abc import Callable
from copy import copy
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pypdf import PdfReader, PdfWriter
from pypdf.constants import UserAccessPermissions
from pypdf.generic import (
    ArrayObject,
    BooleanObject,
    ByteStringObject,
    ContentStream,
    DecodedStreamObject,
    DictionaryObject,
    FloatObject,
    NameObject,
    NumberObject,
    RectangleObject,
    TextStringObject,
)

try:
    import pypdfium2 as pdfium

    PDF_RENDER_AVAILABLE = True
except Exception:
    pdfium = None
    PDF_RENDER_AVAILABLE = False

try:
    import pytesseract
    from pytesseract import TesseractNotFoundError

    OCR_AVAILABLE = True
except Exception:
    pytesseract = None

    class TesseractNotFoundError(Exception):
        pass

    OCR_AVAILABLE = False

try:
    import fitz

    PYMUPDF_AVAILABLE = True
except Exception:
    fitz = None
    PYMUPDF_AVAILABLE = False

PDF_SUFFIXES = {".pdf"}
IMAGE_SUFFIXES = {".jpg", ".jpeg", ".png", ".bmp", ".tif", ".tiff", ".webp"}
OFFICE_WORD_SUFFIXES = {".doc", ".docx", ".docm"}
OFFICE_EXCEL_SUFFIXES = {".xls", ".xlsx", ".xlsm", ".xlsb"}
OFFICE_POWERPOINT_SUFFIXES = {".ppt", ".pptx", ".pptm"}
OFFICE_SUFFIXES = OFFICE_WORD_SUFFIXES | OFFICE_EXCEL_SUFFIXES | OFFICE_POWERPOINT_SUFFIXES
OFFICE_DIALOG_FILTER = (
    "Office files (*.doc *.docx *.docm *.xls *.xlsx *.xlsm *.xlsb *.ppt *.pptx *.pptm);;"
    "Word (*.doc *.docx *.docm);;"
    "Excel (*.xls *.xlsx *.xlsm *.xlsb);;"
    "PowerPoint (*.ppt *.pptx *.pptm);;"
    "All files (*.*)"
)
OCR_LANGUAGE_OPTIONS = {
    "auto": "自動偵測（繁／簡／英）",
    "eng": "英文",
    "chi_tra": "繁體中文",
    "chi_sim": "簡體中文",
    "eng+chi_tra": "英文 + 繁中",
    "eng+chi_sim": "英文 + 簡中",
}

_TRAD_SIMP_PAIRS = (
    ("國", "国"),
    ("個", "个"),
    ("這", "这"),
    ("為", "为"),
    ("與", "与"),
    ("從", "从"),
    ("後", "后"),
    ("發", "发"),
    ("對", "对"),
    ("會", "会"),
    ("學", "学"),
    ("經", "经"),
    ("業", "业"),
    ("東", "东"),
    ("車", "车"),
    ("門", "门"),
    ("開", "开"),
    ("關", "关"),
    ("麼", "么"),
    ("並", "并"),
    ("萬", "万"),
    ("億", "亿"),
    ("餘", "余"),
    ("於", "于"),
    ("過", "过"),
    ("還", "还"),
    ("說", "说"),
    ("請", "请"),
    ("語", "语"),
    ("證", "证"),
    ("據", "据"),
    ("總", "总"),
    ("報", "报"),
    ("務", "务"),
    ("審", "审"),
    ("計", "计"),
    ("財", "财"),
    ("帳", "账"),
    ("虧", "亏"),
    ("損", "损"),
    ("權", "权"),
    ("債", "债"),
    ("資", "资"),
    ("產", "产"),
    ("現", "现"),
    ("廣", "广"),
    ("體", "体"),
    ("銀", "银"),
    ("錄", "录"),
    ("際", "际"),
    ("點", "点"),
    ("樣", "样"),
    ("將", "将"),
    ("來", "来"),
    ("們", "们"),
    ("條", "条"),
    ("額", "额"),
    ("匯", "汇"),
    ("兌", "兑"),
    ("團", "团"),
    ("聯", "联"),
    ("營", "营"),
    ("處", "处"),
    ("號", "号"),
    ("師", "师"),
    ("書", "书"),
    ("頁", "页"),
    ("係", "系"),
)
_TRAD_ONLY_CHARS = {traditional for traditional, _simplified in _TRAD_SIMP_PAIRS}
_SIMP_ONLY_CHARS = {simplified for _traditional, simplified in _TRAD_SIMP_PAIRS}
LAST_OCR_LANGUAGE = ""


@dataclass(frozen=True)
class PageItem:
    pdf_path: Path
    page_index: int
    label: str
    rotation: int = 0


@dataclass(frozen=True)
class TextBlock:
    text: str
    x: float
    y: float
    width: float
    height: float
    font_size: float
    font_name: str = ""
    color_rgb: tuple[float, float, float] = (0.0, 0.0, 0.0)
    bbox: tuple[float, float, float, float] = (0.0, 0.0, 0.0, 0.0)
    font_flags: int = 0
    page_font_name: str = ""
    font_file: str = ""


@dataclass
class FormField:
    name: str
    page_index: int
    field_type: str
    value: str
    choices: tuple[str, ...] = ()


@dataclass(frozen=True)
class SignatureInfo:
    field_name: str
    page_index: int
    signer: str
    signed_at: str
    has_contents: bool
    intact: bool
    detail: str


FORM_WIDGET_TYPE_NAMES = {
    0: "unknown",
    1: "button",
    2: "checkbox",
    3: "combobox",
    4: "listbox",
    5: "radio",
    6: "signature",
    7: "text",
}


def block_right(block: TextBlock) -> float:
    return block.x + block.width


def block_bottom(block: TextBlock) -> float:
    return block.y - block.height


def _merge_block_bbox(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> tuple[float, float, float, float]:
    if left == (0.0, 0.0, 0.0, 0.0):
        return right
    if right == (0.0, 0.0, 0.0, 0.0):
        return left
    return (
        min(left[0], right[0]),
        min(left[1], right[1]),
        max(left[2], right[2]),
        max(left[3], right[3]),
    )


def text_blocks_on_same_line(left: TextBlock, right: TextBlock) -> bool:
    tolerance = max(left.font_size, right.font_size) * 0.6
    return abs(left.y - right.y) <= tolerance


def merge_text_blocks(blocks: list[TextBlock]) -> list[TextBlock]:
    if not blocks:
        return []
    sorted_blocks = sorted(blocks, key=lambda block: (-block.y, block.x))
    lines: list[list[TextBlock]] = []
    for block in sorted_blocks:
        for line in lines:
            if text_blocks_on_same_line(line[0], block):
                line.append(block)
                break
        else:
            lines.append([block])

    merged: list[TextBlock] = []
    for line in lines:
        line = sorted(line, key=lambda block: block.x)
        current = line[0]
        for block in line[1:]:
            gap = block.x - block_right(current)
            max_gap = max(current.font_size, block.font_size) * 1.8
            if gap <= max_gap:
                separator = " " if gap > current.font_size * 0.25 else ""
                left = min(current.x, block.x)
                right = max(block_right(current), block_right(block))
                top = max(current.y, block.y)
                bottom = min(block_bottom(current), block_bottom(block))
                current = TextBlock(
                    text=f"{current.text}{separator}{block.text}",
                    x=left,
                    y=top,
                    width=right - left,
                    height=top - bottom,
                    font_size=max(current.font_size, block.font_size),
                    font_name=current.font_name or block.font_name,
                    color_rgb=current.color_rgb,
                    bbox=_merge_block_bbox(current.bbox, block.bbox),
                    font_flags=current.font_flags or block.font_flags,
                    page_font_name=current.page_font_name or block.page_font_name,
                    font_file=current.font_file or block.font_file,
                )
            else:
                merged.append(current)
                current = block
        merged.append(current)
    return sorted(merged, key=lambda block: (-block.y, block.x))


def parse_pages(spec: str, page_count: int) -> list[int]:
    spec = (spec or "").strip()
    if not spec:
        raise ValueError("請輸入頁碼，例如 1,3,5-8。")

    pages: list[int] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_text, end_text = part.split("-", 1)
            start, end = int(start_text), int(end_text)
            if start > end:
                raise ValueError(f"頁碼範圍不正確：{part}")
            pages.extend(range(start - 1, end))
        else:
            pages.append(int(part) - 1)

    deduped: list[int] = []
    for page in pages:
        if page < 0 or page >= page_count:
            raise ValueError(f"頁碼超出範圍：{page + 1}，目前共有 {page_count} 頁。")
        if page not in deduped:
            deduped.append(page)
    return deduped


def safe_output_name(name: str) -> str:
    return re.sub(r"[^A-Za-z0-9._ -]+", "_", name).strip() or "output.pdf"


_WINDOWS_FORBIDDEN_NAME = re.compile(r'[<>:"/\\|?*\x00-\x1f]')


def suggested_pdf_name_for_source(source: Path) -> str:
    """Keep the original file name (including CJK), only swapping the suffix to .pdf."""

    stem = _WINDOWS_FORBIDDEN_NAME.sub("_", Path(source).stem).strip(" .") or "output"
    return f"{stem}.pdf"


def suggested_pdf_path_for_source(source: Path) -> Path:
    source = Path(source)
    return source.with_name(suggested_pdf_name_for_source(source))


def suggested_images_zip_name_for_source(source: Path) -> str:
    """Keep the original stem (including CJK) and append -images.zip."""

    stem = _WINDOWS_FORBIDDEN_NAME.sub("_", Path(source).stem).strip(" .") or "output"
    if stem.lower().endswith("-images"):
        return f"{stem}.zip"
    return f"{stem}-images.zip"


def suggested_images_zip_path_for_source(source: Path) -> Path:
    source = Path(source)
    return source.with_name(suggested_images_zip_name_for_source(source))


def _report_office_progress(
    progress: Callable[[int, int, str], None] | None,
    current: int,
    total: int,
    text: str,
) -> None:
    if progress is None:
        return
    try:
        progress(current, total, text)
    except Exception:
        return


def open_reader(path: Path, password: str = "") -> PdfReader:
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        if not password:
            raise ValueError(f"{path.name} 已加密，請輸入密碼。")
        result = reader.decrypt(password)
        if result == 0:
            raise ValueError(f"{path.name} 密碼不正確。")
    return reader


def write_pdf(writer: PdfWriter, path: Path) -> None:
    with path.open("wb") as stream:
        writer.write(stream)


def page_item_label(item: PageItem) -> str:
    if item.rotation:
        return f"{item.label} (旋轉 {item.rotation}度)"
    return item.label


def page_object_for_item(reader: PdfReader, item: PageItem):
    page = copy(reader.pages[item.page_index])
    if item.rotation:
        page.rotate(item.rotation)
    return page


def write_page_items_merged(page_items: list[PageItem], indexes: list[int], target: Path, password: str = "") -> None:
    if not indexes:
        raise ValueError("請先選取要擷取的頁面。")
    writer = PdfWriter()
    reader_cache: dict[Path, PdfReader] = {}
    for index in indexes:
        item = page_items[index]
        reader = reader_cache.get(item.pdf_path)
        if reader is None:
            reader = open_reader(item.pdf_path, password)
            reader_cache[item.pdf_path] = reader
        writer.add_page(page_object_for_item(reader, item))
    write_pdf(writer, target)


def write_page_items_separately(page_items: list[PageItem], indexes: list[int], folder: Path, password: str = "") -> int:
    if not indexes:
        raise ValueError("請先選取要擷取的頁面。")
    folder.mkdir(parents=True, exist_ok=True)
    reader_cache: dict[Path, PdfReader] = {}
    for output_index, page_index in enumerate(indexes, start=1):
        item = page_items[page_index]
        reader = reader_cache.get(item.pdf_path)
        if reader is None:
            reader = open_reader(item.pdf_path, password)
            reader_cache[item.pdf_path] = reader
        writer = PdfWriter()
        writer.add_page(page_object_for_item(reader, item))
        filename = safe_output_name(f"{output_index:03d}-{item.pdf_path.stem}-page-{item.page_index + 1}.pdf")
        write_pdf(writer, folder / filename)
    return len(indexes)


def merge_pdf_files(paths: list[Path], target: Path, password: str = "") -> None:
    if len(paths) < 2:
        raise ValueError("合併 PDF 至少需要兩個檔案。")
    writer = PdfWriter()
    for path in paths:
        reader = open_reader(path, password)
        for page in reader.pages:
            writer.add_page(page)
    write_pdf(writer, target)


def split_pdf_to_zip(source: Path, target_zip: Path, password: str = "") -> int:
    reader = open_reader(source, password)
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        generated: list[Path] = []
        for index, page in enumerate(reader.pages, start=1):
            writer = PdfWriter()
            writer.add_page(page)
            part = temp_path / f"page-{index:04d}.pdf"
            write_pdf(writer, part)
            generated.append(part)
        with zipfile.ZipFile(target_zip, "w", zipfile.ZIP_DEFLATED) as archive:
            for item in generated:
                archive.write(item, item.name)
    return len(reader.pages)


def extract_pdf_pages(source: Path, target: Path, pages_spec: str, password: str = "") -> None:
    reader = open_reader(source, password)
    pages = parse_pages(pages_spec, len(reader.pages))
    extract_pdf_page_indexes(source, target, pages, password)


def extract_pdf_page_indexes(
    source: Path,
    target: Path,
    page_indexes: list[int],
    password: str = "",
) -> int:
    reader = open_reader(source, password)
    page_count = len(reader.pages)
    indexes: list[int] = []
    seen: set[int] = set()
    for raw in page_indexes:
        index = int(raw)
        if index < 0 or index >= page_count or index in seen:
            continue
        seen.add(index)
        indexes.append(index)
    if not indexes:
        raise ValueError("沒有可擷取的頁面。")
    writer = PdfWriter()
    for index in indexes:
        writer.add_page(reader.pages[index])
    write_pdf(writer, target)
    return len(indexes)


def apply_edits_and_extract_pages(
    source: Path,
    target: Path,
    page_indexes: list[int],
    *,
    overlays: list[dict] | None = None,
    erase_marks: list[EraseMark] | None = None,
    markups: list[tuple[int, MarkupAnnotation]] | None = None,
    password: str = "",
    remove_content: bool = True,
) -> int:
    """Apply pending overlays/markups, then write only the edited pages."""

    overlays = list(overlays or [])
    erase_marks = list(erase_marks or [])
    markups = list(markups or [])
    wanted: list[int] = []
    seen: set[int] = set()
    for raw in page_indexes:
        index = int(raw)
        if index in seen:
            continue
        seen.add(index)
        wanted.append(index)
    if not wanted:
        raise ValueError("沒有可擷取的修改頁。")
    if not overlays and not erase_marks and not markups:
        return extract_pdf_page_indexes(source, target, wanted, password)

    temps: list[Path] = []
    current = Path(source)
    current_password = password
    try:
        if erase_marks or overlays:
            handle = tempfile.NamedTemporaryFile(suffix=".pdf", prefix="edit_pages_", delete=False)
            handle.close()
            temp = Path(handle.name)
            temps.append(temp)
            apply_erase_then_text_overlays(
                current,
                temp,
                erase_marks,
                overlays,
                current_password,
                remove_content,
            )
            current = temp
            current_password = ""
        if markups:
            handle = tempfile.NamedTemporaryFile(suffix=".pdf", prefix="markup_pages_", delete=False)
            handle.close()
            temp = Path(handle.name)
            temps.append(temp)
            apply_markup_annotations(current, temp, markups, current_password)
            current = temp
            current_password = ""
        return extract_pdf_page_indexes(current, target, wanted, current_password)
    finally:
        for path in temps:
            try:
                path.unlink(missing_ok=True)
            except OSError:
                pass


def delete_pdf_pages(source: Path, target: Path, pages_spec: str, password: str = "") -> None:
    reader = open_reader(source, password)
    delete_pages = set(parse_pages(pages_spec, len(reader.pages)))
    writer = PdfWriter()
    for index, page in enumerate(reader.pages):
        if index not in delete_pages:
            writer.add_page(page)
    if not writer.pages:
        raise ValueError("不能刪除全部頁面。")
    write_pdf(writer, target)


def rotate_pdf_pages(source: Path, target: Path, angle: int, pages_spec: str = "", password: str = "") -> None:
    if angle not in {90, 180, 270}:
        raise ValueError("旋轉角度只支援 90、180、270。")
    reader = open_reader(source, password)
    rotate_pages = set(parse_pages(pages_spec, len(reader.pages))) if pages_spec.strip() else set(range(len(reader.pages)))
    writer = PdfWriter()
    for index, page in enumerate(reader.pages):
        if index in rotate_pages:
            page.rotate(angle)
        writer.add_page(page)
    write_pdf(writer, target)


def encrypt_pdf(source: Path, target: Path, new_password: str, password: str = "") -> None:
    if not new_password:
        raise ValueError("請輸入新密碼。")
    reader = open_reader(source, password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(new_password)
    write_pdf(writer, target)


def decrypt_pdf(source: Path, target: Path, password: str = "") -> None:
    reader = open_reader(source, password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    write_pdf(writer, target)


def compress_pdf(source: Path, target: Path, password: str = "") -> None:
    reader = open_reader(source, password)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    write_pdf(writer, target)


def extract_pdf_text(source: Path, target: Path, password: str = "") -> None:
    reader = open_reader(source, password)
    pieces = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pieces.append(f"--- Page {index} ---\n{text.strip()}\n")
    target.write_text("\n".join(pieces), encoding="utf-8")


def _page_indexes_or_all(source: Path, password: str = "", pages_spec: str = "") -> list[int]:
    reader = open_reader(source, password)
    if (pages_spec or "").strip():
        return parse_pages(pages_spec, len(reader.pages))
    return list(range(len(reader.pages)))


def _pdf_page_size(source: Path, password: str = "", page_index: int = 0) -> tuple[float, float]:
    if PYMUPDF_AVAILABLE and fitz is not None:
        document = fitz.open(str(source))
        try:
            if document.is_encrypted and password:
                document.authenticate(password)
            if document.page_count <= 0:
                return 595.0, 842.0
            page = document[min(max(page_index, 0), document.page_count - 1)]
            return float(page.rect.width), float(page.rect.height)
        finally:
            document.close()
    reader = open_reader(source, password)
    box = reader.pages[min(max(page_index, 0), len(reader.pages) - 1)].mediabox
    return float(box.width), float(box.height)


def _pdf_extractable_char_count(source: Path, password: str = "", pages_spec: str = "") -> int:
    reader = open_reader(source, password)
    indexes = _page_indexes_or_all(source, password, pages_spec)
    total = 0
    for index in indexes:
        total += len((reader.pages[index].extract_text() or "").strip())
    return total


def _is_cjk_char(char: str) -> bool:
    code = ord(char)
    return (
        0x2E80 <= code <= 0x2EFF
        or 0x2F00 <= code <= 0x2FDF
        or 0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0xF900 <= code <= 0xFAFF
    )


def count_cjk_chars(text: str) -> int:
    return sum(1 for char in text or "" if _is_cjk_char(char))


_CJK_RADICAL_SUPPLEMENT = str.maketrans(
    {
        "\u2e9d": "\u6c11",  # ⺠ → 民
        "\u2ed1": "\u9577",  # ⻑ → 長
    }
)


def normalize_cjk_text(text: str) -> str:
    """Fix Tesseract / pypdf CJK radicals and spaces between Chinese characters."""
    if not text:
        return text
    normalized = unicodedata.normalize("NFKC", text).translate(_CJK_RADICAL_SUPPLEMENT)
    return re.sub(
        r"(?<=[\u2e80-\u2fdf\u3400-\u9fff\uf900-\ufaff])[ \t]+(?=[\u2e80-\u2fdf\u3400-\u9fff\uf900-\ufaff])",
        "",
        normalized,
    )


def _pdf_plain_text(source: Path, password: str = "", pages_spec: str = "") -> str:
    indexes = _page_indexes_or_all(source, password, pages_spec)
    if PYMUPDF_AVAILABLE and fitz is not None:
        document = None
        try:
            document = fitz.open(str(source))
            if document.is_encrypted:
                if not password or document.authenticate(password) == 0:
                    raise ValueError(f"{source.name} 已加密，請輸入密碼。")
            return "\n".join((document[index].get_text("text") or "") for index in indexes)
        except ValueError:
            raise
        except Exception:
            pass
        finally:
            if document is not None:
                document.close()
    reader = open_reader(source, password)
    return "\n".join((reader.pages[index].extract_text() or "") for index in indexes)


def _pdf_cjk_char_count(source: Path, password: str = "", pages_spec: str = "") -> int:
    return count_cjk_chars(_pdf_plain_text(source, password, pages_spec))


def pdf_has_usable_text_layer(source: Path, password: str = "", pages_spec: str = "") -> bool:
    try:
        chars = _pdf_extractable_char_count(source, password, pages_spec)
        cjk = _pdf_cjk_char_count(source, password, pages_spec)
    except Exception:
        return True
    return chars >= 8 or cjk >= 8


def _docx_plain_text(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read("word/document.xml").decode("utf-8", errors="replace")
    except Exception:
        return ""
    return "".join(re.findall(r"<w:t[^>]*>([^<]*)</w:t>", xml))


def _docx_keeps_cjk(target: Path, source_cjk: int) -> bool:
    if source_cjk < 24:
        return True
    if not target.is_file():
        return False
    kept = count_cjk_chars(_docx_plain_text(target))
    return kept >= max(12, int(source_cjk * 0.35))


def _tesseract_language(language: str) -> str:
    mapping = {
        "eng+chi_tra": "chi_tra+eng",
        "eng+chi_sim": "chi_sim+eng",
        "eng+chi_tra+chi_sim": "chi_tra+chi_sim+eng",
        "auto": "chi_tra+chi_sim+eng",
    }
    return mapping.get(language or "", language or "eng")


def _tesseract_config(language: str) -> str:
    lang = _tesseract_language(language)
    if "chi_" in lang:
        return "--oem 1 --psm 3 -c preserve_interword_spaces=1"
    return "--psm 3 -c preserve_interword_spaces=1"


def _tesseract_column_config(tess_config: str) -> str:
    config = tess_config or "--psm 6 -c preserve_interword_spaces=1"
    if re.search(r"--psm \d+", config):
        return re.sub(r"--psm \d+", "--psm 6", config)
    return f"{config} --psm 6".strip()


def _column_tess_lang(tess_lang: str, prefer_cjk: bool) -> str:
    parts = [part for part in (tess_lang or "eng").split("+") if part]
    if not parts:
        return "eng"
    chinese = [part for part in parts if part.startswith("chi_")]
    rest = [part for part in parts if not part.startswith("chi_")]
    if not chinese:
        return "+".join(parts)
    if prefer_cjk:
        return "+".join(chinese + rest)
    return "+".join(rest + chinese) if rest else "+".join(chinese)


def _ocr_render_dpi(language: str, dpi: int) -> int:
    requested = language or ""
    if requested in {"auto", "eng+chi_tra+chi_sim"} or "chi_" in requested:
        if dpi < 300:
            return 300
    return dpi


def detect_ocr_language(text: str) -> str:
    """Pick a Tesseract language pack from a short OCR sample."""
    sample = text or ""
    trad = sum(1 for char in sample if char in _TRAD_ONLY_CHARS)
    simp = sum(1 for char in sample if char in _SIMP_ONLY_CHARS)
    cjk = count_cjk_chars(sample)
    latin = sum(1 for char in sample if ("A" <= char <= "Z") or ("a" <= char <= "z"))
    if cjk < 8 and latin >= 12:
        return "eng"
    if simp > trad and simp >= 3:
        return "eng+chi_sim"
    if trad > simp and trad >= 3:
        return "eng+chi_tra"
    if cjk >= 8:
        return "eng+chi_tra+chi_sim"
    if latin >= 8:
        return "eng"
    return "eng+chi_tra+chi_sim"


def ocr_language_label(language: str) -> str:
    return OCR_LANGUAGE_OPTIONS.get(language, language or "自動偵測")


def last_ocr_language() -> str:
    return LAST_OCR_LANGUAGE


def ocr_language_short_label(language: str) -> str:
    return {
        "eng": "英文",
        "chi_tra": "繁中",
        "chi_sim": "簡中",
        "eng+chi_tra": "英文+繁中",
        "eng+chi_sim": "英文+簡中",
        "eng+chi_tra+chi_sim": "中英混合",
        "auto": "自動偵測",
    }.get(language, language or "自動偵測")


def _set_last_ocr_language(language: str) -> None:
    global LAST_OCR_LANGUAGE
    LAST_OCR_LANGUAGE = language


def _resolved_ocr_settings(language: str, sample_image: Image.Image | None = None) -> tuple[str, str, str]:
    resolved = resolve_ocr_language(language, sample_image)
    _set_last_ocr_language(resolved)
    return resolved, _tesseract_language(resolved), _tesseract_config(resolved)


def _ocr_image_to_string(image: Image.Image, language: str) -> str:
    prepared = _prepare_ocr_image(image)
    try:
        tess_lang = _tesseract_language(language)
        return pytesseract.image_to_string(
            prepared,
            lang=tess_lang,
            config=_tesseract_config(language),
        )
    finally:
        if prepared is not image:
            prepared.close()


def _probe_ocr_text(image: Image.Image) -> str:
    for language in ("eng+chi_tra+chi_sim", "eng+chi_sim", "eng+chi_tra", "eng"):
        try:
            text = _ocr_image_to_string(image, language)
        except Exception:
            continue
        if (text or "").strip():
            return text
    return ""


def resolve_ocr_language(language: str, sample_image: Image.Image | None = None) -> str:
    requested = (language or "auto").strip() or "auto"
    if requested != "auto":
        return requested
    if sample_image is None:
        return "eng+chi_tra+chi_sim"
    return detect_ocr_language(_probe_ocr_text(sample_image))


def _prepare_ocr_image(image: Image.Image) -> Image.Image:
    from PIL import ImageFilter, ImageOps

    width, height = image.size
    longest = max(width, height)
    work = image
    resample = getattr(getattr(Image, "Resampling", Image), "LANCZOS", Image.LANCZOS)
    if longest < 1800:
        scale = 1800 / float(longest)
        work = image.resize((max(1, int(width * scale)), max(1, int(height * scale))), resample)
    elif longest > 3200:
        scale = 3200 / float(longest)
        work = image.resize((max(1, int(width * scale)), max(1, int(height * scale))), resample)
    gray = work.convert("L")
    if work is not image:
        work.close()
    gray = ImageOps.autocontrast(gray, cutoff=1)
    return gray.filter(ImageFilter.UnsharpMask(radius=1.2, percent=140, threshold=2))


def _write_unlocked_pdf_copy(source: Path, password: str = "") -> Path:
    reader = open_reader(source, password)
    if not getattr(reader, "is_encrypted", False) and not password:
        return source
    handle = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    handle.close()
    target = Path(handle.name)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    write_pdf(writer, target)
    return target


def _docx_available() -> bool:
    try:
        import docx  # noqa: F401

        return True
    except Exception:
        return False


def _xlsx_available() -> bool:
    try:
        import openpyxl  # noqa: F401

        return True
    except Exception:
        return False


def _pdf2docx_available() -> bool:
    try:
        from pdf2docx import Converter  # noqa: F401

        return True
    except Exception:
        return False


def _pdf_to_docx_via_word(source: Path, target: Path) -> tuple[bool, str]:
    if sys.platform != "win32":
        return False, "not-windows"
    try:
        import pythoncom
        import win32com.client
    except ImportError:
        return False, "no-com"
    initialized = False
    word = None
    document = None
    try:
        pythoncom.CoInitialize()
        initialized = True
        word = _dispatch_com_app(win32com.client, ("Word.Application", "Word.Application.16"))
        try:
            word.Visible = False
        except Exception:
            pass
        try:
            word.DisplayAlerts = 0
        except Exception:
            pass
        document = _try_com_call(
            lambda: word.Documents.Open(str(source.resolve()), False, True, False),
            lambda: word.Documents.Open(str(source.resolve())),
        )
        target.parent.mkdir(parents=True, exist_ok=True)
        _try_com_call(
            lambda: document.SaveAs2(str(target.resolve()), 16),
            lambda: document.SaveAs(str(target.resolve()), 16),
        )
        if target.is_file() and target.stat().st_size > 0:
            return True, "word"
        return False, "empty-output"
    except Exception as exc:
        return False, str(exc) or type(exc).__name__
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        if initialized:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _pdf_to_office_via_libreoffice(source: Path, target: Path, convert_to: str) -> tuple[bool, str]:
    executable = find_libreoffice_executable()
    if executable is None:
        return False, "no-libreoffice"
    with tempfile.TemporaryDirectory() as temp_dir:
        out_dir = Path(temp_dir)
        command = [
            str(executable),
            "--headless",
            "--norestore",
            "--nolockcheck",
            "--convert-to",
            convert_to,
            "--outdir",
            str(out_dir),
            str(source.resolve()),
        ]
        try:
            subprocess.run(command, capture_output=True, timeout=180, check=False)
        except Exception as exc:
            return False, str(exc)
        produced = out_dir / f"{source.stem}.{convert_to}"
        if not produced.exists():
            matches = list(out_dir.glob(f"*.{convert_to}"))
            produced = matches[0] if matches else produced
        if not produced.exists() or produced.stat().st_size <= 0:
            return False, "empty-output"
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(produced, target)
        return True, "libreoffice"


def _docx_body_font_name(text: str) -> str:
    sample = text or ""
    if count_cjk_chars(sample) < 8:
        return "Calibri"
    trad = sum(1 for char in sample if char in _TRAD_ONLY_CHARS)
    simp = sum(1 for char in sample if char in _SIMP_ONLY_CHARS)
    if trad > simp:
        return "微軟正黑體"
    return "微软雅黑"


def _docx_body_size(sizes: list[float]) -> float:
    if not sizes:
        return 11.0
    ordered = sorted(float(size) for size in sizes if size)
    if not ordered:
        return 11.0
    median = ordered[len(ordered) // 2]
    return max(10.0, min(12.0, round(median)))


def _docx_display_size(raw_size: float, body_size: float) -> float:
    size = float(raw_size or body_size)
    if size >= body_size * 1.75:
        return min(18.0, body_size + 5)
    if size >= body_size * 1.3:
        return min(14.0, body_size + 2)
    if size <= body_size * 0.8:
        return max(9.0, body_size - 1)
    return body_size


def _docx_set_rfonts(element, font_name: str) -> None:
    from docx.oxml.ns import qn

    r_pr = element.get_or_add_rPr()
    r_fonts = r_pr.find(qn("w:rFonts"))
    if r_fonts is None:
        from docx.oxml import OxmlElement

        r_fonts = OxmlElement("w:rFonts")
        r_pr.append(r_fonts)
    r_fonts.set(qn("w:ascii"), font_name)
    r_fonts.set(qn("w:hAnsi"), font_name)
    r_fonts.set(qn("w:eastAsia"), font_name)
    r_fonts.set(qn("w:cs"), font_name)


def _configure_docx_normal_style(document, font_name: str, body_size: float) -> None:
    from docx.shared import Pt

    style = document.styles["Normal"]
    style.font.name = font_name
    style.font.size = Pt(body_size)
    _docx_set_rfonts(style.element, font_name)


def _configure_docx_section_from_pdf(document, width_pt: float, height_pt: float) -> None:
    from docx.enum.section import WD_ORIENT
    from docx.shared import Inches, Pt

    width_pt = max(200.0, float(width_pt or 595.0))
    height_pt = max(200.0, float(height_pt or 842.0))
    section = document.sections[0]
    landscape = width_pt > height_pt
    section.orientation = WD_ORIENT.LANDSCAPE if landscape else WD_ORIENT.PORTRAIT
    section.page_width = Pt(width_pt)
    section.page_height = Pt(height_pt)
    margin = Pt(28) if landscape else Inches(0.85)
    section.left_margin = margin
    section.right_margin = margin
    section.top_margin = margin
    section.bottom_margin = margin


def _style_docx_run(run, font_name: str, size_pt: float, bold: bool = False) -> None:
    from docx.shared import Pt

    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.italic = False
    _docx_set_rfonts(run._element, font_name)


def _add_docx_paragraph(document, text: str, font_name: str, size_pt: float, bold: bool = False, align=None):
    paragraph = document.add_paragraph()
    run = paragraph.add_run(text)
    _style_docx_run(run, font_name, size_pt, bold)
    _set_docx_paragraph_spacing(paragraph, after_pt=3.0)
    if align is not None:
        paragraph.alignment = align
    return paragraph


def _docx_cell_font_name(text: str, fallback: str = "Calibri") -> str:
    sample = text or ""
    if count_cjk_chars(sample) <= 0:
        return "Calibri" if any(char.isascii() and char.isalpha() for char in sample) else fallback
    trad = sum(1 for char in sample if char in _TRAD_ONLY_CHARS)
    simp = sum(1 for char in sample if char in _SIMP_ONLY_CHARS)
    if simp > trad:
        return "微软雅黑"
    return "微軟正黑體"


def _set_docx_paragraph_spacing(paragraph, after_pt: float = 0.0) -> None:
    from docx.shared import Pt

    paragraph.paragraph_format.space_before = Pt(0)
    paragraph.paragraph_format.space_after = Pt(after_pt)
    paragraph.paragraph_format.line_spacing = 1.08


def _set_docx_cell_text(cell, text: str, font_name: str, size_pt: float) -> None:
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT

    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
    cleaned = normalize_cjk_text(text or "").strip()
    cell.text = cleaned
    paragraph = cell.paragraphs[0]
    _set_docx_paragraph_spacing(paragraph)
    if paragraph.runs:
        _style_docx_run(paragraph.runs[0], font_name, size_pt)


def _set_docx_table_no_borders(table) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    borders = OxmlElement("w:tblBorders")
    for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
        element = OxmlElement(f"w:{edge}")
        element.set(qn("w:val"), "nil")
        element.set(qn("w:sz"), "0")
        element.set(qn("w:space"), "0")
        element.set(qn("w:color"), "auto")
        borders.append(element)
    table._tbl.tblPr.append(borders)


def _style_docx_layout_table(table, document, bordered: bool = False) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Emu

    section = document.sections[0]
    usable = int(section.page_width - section.left_margin - section.right_margin)
    cols = max(len(table.columns), 1)
    col_emu = max(1, usable // cols)
    table.autofit = False
    table.allow_autofit = False
    tbl_pr = table._tbl.tblPr
    tbl_w = tbl_pr.find(qn("w:tblW"))
    if tbl_w is None:
        tbl_w = OxmlElement("w:tblW")
        tbl_pr.append(tbl_w)
    tbl_w.set(qn("w:w"), str(int(Emu(usable).twips)))
    tbl_w.set(qn("w:type"), "dxa")
    layout = tbl_pr.find(qn("w:tblLayout"))
    if layout is None:
        layout = OxmlElement("w:tblLayout")
        tbl_pr.append(layout)
    layout.set(qn("w:type"), "fixed")
    grid = table._tbl.find(qn("w:tblGrid"))
    if grid is not None:
        for grid_col in grid.findall(qn("w:gridCol")):
            grid_col.set(qn("w:w"), str(int(Emu(col_emu).twips)))
    if bordered:
        try:
            table.style = "Table Grid"
        except Exception:
            pass
    else:
        _set_docx_table_no_borders(table)
    for row in table.rows:
        for cell in row.cells:
            cell.width = Emu(col_emu)
            tc_pr = cell._tc.get_or_add_tcPr()
            tc_w = tc_pr.find(qn("w:tcW"))
            if tc_w is None:
                tc_w = OxmlElement("w:tcW")
                tc_pr.append(tc_w)
            tc_w.set(qn("w:w"), str(int(Emu(col_emu).twips)))
            tc_w.set(qn("w:type"), "dxa")


def _style_bilingual_docx_table(table, document) -> None:
    _style_docx_layout_table(table, document, bordered=False)


def _add_docx_layout_table(
    document,
    rows: list[list[str]],
    default_font: str,
    size_pt: float,
    bordered: bool = False,
) -> None:
    cleaned_rows: list[list[str]] = []
    for row in rows or []:
        cells = [normalize_cjk_text(str(cell or "")).strip() for cell in row]
        if any(cells):
            cleaned_rows.append(cells)
    if not cleaned_rows:
        return
    width = max(len(row) for row in cleaned_rows)
    cleaned_rows = [row + [""] * (width - len(row)) for row in cleaned_rows]
    table = document.add_table(rows=len(cleaned_rows), cols=width)
    _style_docx_layout_table(table, document, bordered=bordered)
    for index, row in enumerate(cleaned_rows):
        for column, text in enumerate(row):
            _set_docx_cell_text(
                table.cell(index, column),
                text,
                _docx_cell_font_name(text, default_font),
                size_pt,
            )


def _page_rows_are_two_column(rows: list[list[str]]) -> bool:
    if len(rows) < 3:
        return False
    wide = 0
    paired = 0
    for row in rows:
        filled = [str(cell).strip() for cell in row if str(cell).strip()]
        if len(filled) >= 3:
            wide += 1
        if len(row) >= 2 and str(row[0]).strip() and str(row[1]).strip() and len(filled) == 2:
            paired += 1
    if wide >= 2:
        return False
    return paired >= max(3, int(len(rows) * 0.3))


def _rows_are_side_by_side_layout(rows: list[list[str]]) -> bool:
    if not _page_rows_are_two_column(rows):
        return False
    left = " ".join(str(row[0]).strip() for row in rows if row)
    right = " ".join(str(row[1]).strip() for row in rows if len(row) > 1)
    if min(len(left), len(right)) < 12:
        return False
    right_alnum = [char for char in right if char.isalnum() or _is_cjk_char(char)]
    if right_alnum and sum(char.isdigit() for char in right_alnum) >= max(3, int(len(right_alnum) * 0.55)):
        return False
    left_cjk = count_cjk_chars(left)
    right_cjk = count_cjk_chars(right)
    left_latin = sum(1 for char in left if char.isascii() and char.isalpha())
    right_latin = sum(1 for char in right if char.isascii() and char.isalpha())
    bilingual = (left_latin >= 16 and right_cjk >= 6 and left_cjk <= left_latin) or (
        right_latin >= 16 and left_cjk >= 6 and right_cjk <= right_latin
    )
    balanced = min(len(left), len(right)) >= 40
    return bilingual or balanced


def _ocr_rows_use_two_column_table(rows: list[list[str]]) -> bool:
    if _page_rows_are_two_column(rows):
        return True
    paired = 0
    for row in rows:
        if len(row) >= 2 and str(row[0]).strip() and str(row[1]).strip():
            paired += 1
    return paired >= 2


def _add_docx_ocr_page(document, rows: list[list[str]], default_font: str, size_pt: float = 11.0) -> None:
    cleaned_rows: list[list[str]] = []
    for row in rows or []:
        cells = [normalize_cjk_text(str(cell or "")).strip() for cell in row]
        if any(cells):
            cleaned_rows.append(cells)
    if not cleaned_rows:
        _add_docx_paragraph(document, "（此頁沒有可辨識文字）", default_font, size_pt)
        return
    if _ocr_rows_use_two_column_table(cleaned_rows):
        _add_docx_layout_table(document, cleaned_rows, default_font, size_pt, bordered=False)
        return
    for row in cleaned_rows:
        text = "  ".join(cell for cell in row if cell)
        if text:
            _add_docx_paragraph(document, text, _docx_cell_font_name(text, default_font), size_pt)


def _block_bbox(block: TextBlock) -> tuple[float, float, float, float]:
    if block.bbox != (0.0, 0.0, 0.0, 0.0):
        return block.bbox
    return (block.x, block.y - block.height, block.x + block.width, block.y)


def _text_blocks_to_words(blocks: list[TextBlock]) -> list[tuple]:
    words: list[tuple] = []
    for block in blocks:
        text = normalize_cjk_text((block.text or "").strip())
        if not text:
            continue
        x0, y0, x1, y1 = _block_bbox(block)
        words.append((float(x0), float(y0), float(x1), float(y1), text, 0, 0, 0))
    return words


def _fitz_page_words(page) -> list[tuple]:
    words: list[tuple] = []
    try:
        raw = page.get_text("words") or []
    except Exception:
        return words
    for item in raw:
        x0, y0, x1, y1, text, *_rest = item
        cleaned = normalize_cjk_text(str(text or "")).strip()
        if cleaned:
            words.append((float(x0), float(y0), float(x1), float(y1), cleaned, 0, 0, 0))
    return words


def _page_width_from_blocks(blocks: list[TextBlock], fallback: float = 500.0) -> float:
    rights = [_block_bbox(block)[2] for block in blocks]
    if not rights:
        return fallback
    return max(max(rights) * 1.08, fallback * 0.5)


def _block_is_centered(block: TextBlock, page_width: float) -> bool:
    if page_width <= 0:
        return False
    x0, _y0, x1, _y1 = _block_bbox(block)
    mid = (x0 + x1) / 2
    return x0 > page_width * 0.18 and abs(mid - page_width / 2) <= page_width * 0.14


def _add_docx_formatted_blocks(
    document,
    blocks: list[TextBlock],
    page_width: float,
    default_font: str,
    body_size: float,
) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    ordered = sorted(blocks, key=lambda block: (_block_bbox(block)[1], _block_bbox(block)[0]))
    for block in ordered:
        text = normalize_cjk_text((block.text or "").strip())
        if not text:
            continue
        bold = bool(block.font_flags & 16) or "bold" in (block.font_name or "").lower()
        size = _docx_display_size(float(block.font_size or body_size), body_size)
        font_name = _docx_cell_font_name(text, default_font)
        align = WD_ALIGN_PARAGRAPH.CENTER if _block_is_centered(block, page_width) else None
        _add_docx_paragraph(document, text, font_name, size, bold, align=align)


def _add_docx_page_layout(
    document,
    page_width: float,
    blocks: list[TextBlock],
    words: list[tuple],
    tables: list[list[list[str]]],
    default_font: str,
    body_size: float,
) -> None:
    col_rows = _ocr_words_to_rows(words, page_width) if words else []
    if _rows_are_side_by_side_layout(col_rows):
        _add_docx_layout_table(document, col_rows, default_font, body_size, bordered=False)
        return
    useful_tables = [rows for rows in tables if _table_looks_like_data(rows)]
    if useful_tables:
        for rows in useful_tables:
            bordered = _table_filled_columns(rows) >= 3
            _add_docx_layout_table(document, rows, default_font, body_size, bordered=bordered)
        return
    if col_rows and _is_useful_table(col_rows) and _table_filled_columns(col_rows) >= 3:
        _add_docx_layout_table(document, col_rows, default_font, body_size, bordered=True)
        return
    if not blocks:
        _add_docx_paragraph(document, "（此頁沒有可抽取的文字）", default_font, body_size)
        return
    _add_docx_formatted_blocks(document, blocks, page_width, default_font, body_size)


def _write_docx_from_text_blocks(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    progress=None,
) -> None:
    from docx import Document

    indexes = _page_indexes_or_all(source, password, pages_spec)
    layouts: list[tuple[float, float, list[TextBlock], list[tuple], list[list[list[str]]]]] = []
    all_text: list[str] = []
    all_sizes: list[float] = []
    if PYMUPDF_AVAILABLE and fitz is not None:
        document_pdf = fitz.open(str(source))
        try:
            if document_pdf.is_encrypted:
                if not password or document_pdf.authenticate(password) == 0:
                    raise ValueError(f"{source.name} 已加密，請輸入密碼。")
            total_pages = len(indexes)
            for offset, page_index in enumerate(indexes):
                if page_index < 0 or page_index >= document_pdf.page_count:
                    raise ValueError("頁碼超出範圍。")
                if progress:
                    progress(
                        offset + 1,
                        total_pages,
                        f"正在檢查第 {page_index + 1} / {total_pages} 頁版面…",
                    )
                page = document_pdf[page_index]
                blocks = _text_blocks_from_fitz_page(page)
                if not blocks:
                    text = normalize_cjk_text((page.get_text("text") or "").strip())
                    for line in text.splitlines():
                        if line.strip():
                            blocks.append(TextBlock(line.strip(), 0, 0, 0, 0, 11.0))
                words = _fitz_page_words(page) or _text_blocks_to_words(blocks)
                tables = _pymupdf_page_tables(page)
                layouts.append((float(page.rect.width), float(page.rect.height), blocks, words, tables))
                for block in blocks:
                    all_text.append(block.text or "")
                    if block.font_size:
                        all_sizes.append(float(block.font_size))
        finally:
            document_pdf.close()
    else:
        reader = open_reader(source, password)
        for page_index in indexes:
            blocks = extract_page_text_blocks(source, page_index, password)
            if not blocks:
                text = normalize_cjk_text((reader.pages[page_index].extract_text() or "").strip())
                blocks = [
                    TextBlock(line.strip(), 0, 0, 0, 0, 11.0)
                    for line in text.splitlines()
                    if line.strip()
                ]
            box = reader.pages[page_index].mediabox
            page_width = float(box.width)
            page_height = float(box.height)
            words = _text_blocks_to_words(blocks)
            layouts.append((page_width, page_height, blocks, words, []))
            for block in blocks:
                all_text.append(block.text or "")
                if block.font_size:
                    all_sizes.append(float(block.font_size))
    if not any(layout[2] or layout[3] or layout[4] for layout in layouts):
        layouts = [
            (
                595.0,
                842.0,
                [TextBlock("（此頁沒有可抽取的文字）", 0, 0, 0, 0, 11.0)],
                [],
                [],
            )
        ]
    font_name = _docx_body_font_name("\n".join(all_text))
    body_size = _docx_body_size(all_sizes)
    document = Document()
    _configure_docx_normal_style(document, font_name, body_size)
    first_width, first_height = layouts[0][0], layouts[0][1]
    _configure_docx_section_from_pdf(document, first_width, first_height)
    total = len(layouts)
    for offset, (page_width, _page_height, blocks, words, tables) in enumerate(layouts):
        if progress:
            progress(offset + 1, total, f"正在重建第 {offset + 1} 頁版面…")
        if offset:
            document.add_page_break()
        _add_docx_page_layout(document, page_width, blocks, words, tables, font_name, body_size)
    if progress:
        progress(total, total, f"已重建 {total} 頁版面")
    target.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(target))


def _write_docx_from_ocr(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    language: str = "auto",
    dpi: int = 200,
    progress=None,
) -> int:
    from docx import Document

    pages, tess_lang, tess_config = _ocr_iter_pages(source, password, pages_spec, language, dpi, progress)
    page_rows: list[list[list[str]]] = []
    all_text: list[str] = []
    count = 0
    for _page_index, image in pages:
        count += 1
        if progress:
            progress(count - 1, 0, f"掃描件 OCR 第 {_page_index + 1} 頁，請稍候…")
        try:
            rows = _ocr_image_to_rows(image, tess_lang, tess_config)
        finally:
            image.close()
        page_rows.append(rows)
        all_text.extend(cell for row in rows for cell in row)
    if progress:
        progress(count, count, f"OCR 完成 {count} 頁")
    font_name = _docx_body_font_name("\n".join(all_text))
    document = Document()
    _configure_docx_normal_style(document, font_name, 11.0)
    width_pt, height_pt = _pdf_page_size(source, password)
    _configure_docx_section_from_pdf(document, width_pt, height_pt)
    for offset, rows in enumerate(page_rows):
        if offset:
            document.add_page_break()
        _add_docx_ocr_page(document, rows, font_name)
    target.parent.mkdir(parents=True, exist_ok=True)
    document.save(str(target))
    return count


def _convert_pdf_to_docx_pdf2docx(source: Path, target: Path) -> bool:
    if not _pdf2docx_available():
        return False
    try:
        from pdf2docx import Converter

        converter = Converter(str(source))
        try:
            converter.convert(str(target))
        finally:
            converter.close()
        return target.is_file() and target.stat().st_size > 0
    except Exception:
        target.unlink(missing_ok=True)
        return False


def pdf_to_docx(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    language: str = "auto",
    dpi: int = 200,
    progress=None,
) -> tuple[int, str]:
    if not _docx_available():
        raise ValueError("PDF 轉 Word 需要安裝 python-docx。")
    indexes = _page_indexes_or_all(source, password, pages_spec)
    if not indexes:
        raise ValueError("沒有可轉換的頁面。")
    if progress:
        progress(0, 0, "正在檢查 PDF 文字層…")
    unlocked = _write_unlocked_pdf_copy(source, password)
    try:
        char_count = _pdf_extractable_char_count(unlocked, "", pages_spec)
        source_cjk = _pdf_cjk_char_count(unlocked, "", pages_spec)
        plain_len = len(_pdf_plain_text(unlocked, "", pages_spec).strip())
        has_text_layer = char_count >= 8 or source_cjk >= 8 or plain_len >= 8
        cjk_heavy = source_cjk >= 40
        limited_pages = bool((pages_spec or "").strip())

        def accept_layout(method: str) -> bool:
            if not target.is_file() or target.stat().st_size <= 0:
                return False
            if _docx_keeps_cjk(target, source_cjk):
                return True
            target.unlink(missing_ok=True)
            return False

        # Word/LibreOffice often drop subset CJK fonts. Prefer Unicode extractors first.
        if cjk_heavy:
            if not limited_pages and _convert_pdf_to_docx_pdf2docx(unlocked, target) and accept_layout("pdf2docx"):
                return len(indexes), "pdf2docx"
            if has_text_layer:
                _write_docx_from_text_blocks(source, target, password, pages_spec, progress)
                return len(indexes), "text"
        elif char_count >= 20 and not limited_pages:
            ok, _reason = _pdf_to_docx_via_word(unlocked, target)
            if ok and accept_layout("word"):
                return len(indexes), "word"
            ok, _reason = _pdf_to_office_via_libreoffice(unlocked, target, "docx")
            if ok and accept_layout("libreoffice"):
                return len(indexes), "libreoffice"
            if _convert_pdf_to_docx_pdf2docx(unlocked, target) and accept_layout("pdf2docx"):
                return len(indexes), "pdf2docx"
        if has_text_layer:
            _write_docx_from_text_blocks(source, target, password, pages_spec, progress)
            return len(indexes), "text"
        page_count = _write_docx_from_ocr(source, target, password, pages_spec, language, dpi, progress)
        return page_count, "ocr"
    finally:
        if unlocked != source:
            unlocked.unlink(missing_ok=True)


def _normalize_extracted_rows(raw_rows) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in raw_rows or []:
        cells = []
        for cell in row:
            text = "" if cell is None else str(cell).replace("\n", " ").strip()
            cells.append(normalize_cjk_text(text))
        if any(cell for cell in cells):
            rows.append(cells)
    return _pad_table_rows(rows)


def _pad_table_rows(rows: list[list[str]]) -> list[list[str]]:
    width = max((len(row) for row in rows), default=0)
    return [row + [""] * (width - len(row)) for row in rows]


def _table_filled_columns(rows: list[list[str]]) -> int:
    if not rows:
        return 0
    width = max(len(row) for row in rows)
    filled = 0
    for column in range(width):
        if any(column < len(row) and str(row[column]).strip() for row in rows):
            filled += 1
    return filled


def _is_useful_table(rows: list[list[str]]) -> bool:
    if not rows:
        return False
    nonempty_rows = [row for row in rows if any(str(cell).strip() for cell in row)]
    filled_cols = _table_filled_columns(rows)
    cells = sum(1 for row in nonempty_rows for cell in row if str(cell).strip())
    if filled_cols < 2 or cells < 6:
        return False
    if filled_cols >= 3 and len(nonempty_rows) >= 2:
        return True
    return len(nonempty_rows) >= 3


def _table_looks_like_data(rows: list[list[str]]) -> bool:
    if not _is_useful_table(rows):
        return False
    if _table_looks_fragmented(rows):
        return False
    cols = _table_filled_columns(rows)
    nonempty = [row for row in rows if any(str(cell).strip() for cell in row)]
    filled = sum(1 for row in nonempty for cell in row if str(cell).strip())
    total = max(1, len(nonempty) * max(cols, 1))
    if cols >= 6 and filled < total * 0.45:
        return False
    return True


def _table_looks_fragmented(rows: list[list[str]]) -> bool:
    cells = [str(cell).strip() for row in rows for cell in row if str(cell).strip()]
    if not cells:
        return True
    cols = _table_filled_columns(rows)
    lengths = sorted(len(cell) for cell in cells)
    median = lengths[len(lengths) // 2]
    short = sum(1 for length in lengths if length <= 8)
    if cols >= 6 and median <= 12:
        return True
    if cols >= 5 and short >= max(6, int(len(cells) * 0.55)):
        return True
    return False


def _join_word_cluster(words: list[tuple]) -> str:
    texts = [str(word[4]).strip() for word in words if str(word[4]).strip()]
    if not texts:
        return ""
    joined = texts[0]
    for index, text in enumerate(texts[1:], start=1):
        previous = texts[index - 1]
        if previous and text and _is_cjk_char(previous[-1]) and _is_cjk_char(text[0]):
            joined += text
        else:
            joined += " " + text
    return normalize_cjk_text(joined).strip()


def _group_words_into_lines(words: list[tuple], y_tol: float | None = None) -> list[list[tuple]]:
    if not words:
        return []
    heights = sorted(max(1.0, float(word[3]) - float(word[1])) for word in words)
    median_height = heights[len(heights) // 2]
    tolerance = y_tol if y_tol is not None else max(3.0, median_height * 0.45)
    ordered = sorted(words, key=lambda word: ((float(word[1]) + float(word[3])) / 2, float(word[0])))
    lines: list[list[tuple]] = []
    line_ys: list[float] = []
    for word in ordered:
        y_mid = (float(word[1]) + float(word[3])) / 2
        if lines and abs(y_mid - line_ys[-1]) <= tolerance:
            lines[-1].append(word)
            count = len(lines[-1])
            line_ys[-1] = (line_ys[-1] * (count - 1) + y_mid) / count
        else:
            lines.append([word])
            line_ys.append(y_mid)
    for line in lines:
        line.sort(key=lambda word: float(word[0]))
    return lines


def _column_gap_threshold(lines: list[list[tuple]], page_width: float) -> float:
    gaps: list[float] = []
    for line in lines:
        for previous, current in zip(line, line[1:]):
            gap = float(current[0]) - float(previous[2])
            if gap > 0:
                gaps.append(gap)
    floor = max(10.0, page_width * 0.03)
    ceiling = max(floor, page_width * 0.14)
    if not gaps:
        return max(floor, page_width * 0.025)
    small = [gap for gap in gaps if gap < page_width * 0.04]
    if small:
        typical = small[len(small) // 2]
        return max(floor, min(ceiling, typical * 4.0))
    return max(floor, page_width * 0.04)


def _split_line_by_x_gaps(line_words: list[tuple], min_gap: float) -> list[str]:
    if not line_words:
        return []
    clusters = [[line_words[0]]]
    for previous, current in zip(line_words, line_words[1:]):
        gap = float(current[0]) - float(previous[2])
        if gap >= min_gap:
            clusters.append([current])
        else:
            clusters[-1].append(current)
    return [_join_word_cluster(cluster) for cluster in clusters]


def _words_to_table_rows(words: list[tuple], page_width: float) -> list[list[str]]:
    if not words:
        return []
    lines = _group_words_into_lines(words)
    threshold = _column_gap_threshold(lines, page_width)
    rows = [_split_line_by_x_gaps(line, threshold) for line in lines]
    rows = [row for row in rows if any(str(cell).strip() for cell in row)]
    return _pad_table_rows(rows)


def _text_lines_to_rows(lines: list[str]) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in lines:
        cleaned = normalize_cjk_text(line).strip()
        if not cleaned:
            continue
        if cleaned.startswith("--- Page "):
            rows.append([cleaned])
            continue
        if "\t" in cleaned:
            parts = [normalize_cjk_text(part).strip() for part in cleaned.split("\t")]
        else:
            parts = [part.strip() for part in re.split(r" {2,}", cleaned)]
        rows.append(parts or [cleaned])
    return _pad_table_rows(rows)


def _page_looks_scanned(page) -> bool:
    text = page.get_text("text") or ""
    stripped = text.strip()
    if len(stripped) >= 40 or count_cjk_chars(text) >= 16:
        return False
    if len(stripped) < 8 and count_cjk_chars(text) < 4:
        return True
    try:
        images = page.get_images() or []
    except Exception:
        images = []
    return bool(images) and len(stripped) < 24


def _pymupdf_page_tables(page) -> list[list[list[str]]]:
    # Lined tables only. Text-grid strategies split words across cells and can
    # freeze the UI for many seconds on pages with dense drawings.
    try:
        finder = page.find_tables()
    except Exception:
        return []
    useful: list[list[list[str]]] = []
    for table in getattr(finder, "tables", None) or []:
        try:
            rows = _normalize_extracted_rows(table.extract())
        except Exception:
            continue
        if _is_useful_table(rows) and not _table_looks_fragmented(rows):
            useful.append(rows)
    return useful


def _page_words_to_table(page) -> list[list[str]] | None:
    try:
        words = page.get_text("words") or []
    except Exception:
        return None
    rows = _words_to_table_rows(words, float(page.rect.width))
    if rows and _is_useful_table(rows) and not _table_looks_fragmented(rows):
        return rows
    return None


def _page_text_to_rows(page) -> list[list[str]]:
    text = normalize_cjk_text(page.get_text("text") or "")
    return _text_lines_to_rows(text.splitlines())


def _is_noisy_ocr_token(text: str, confidence: float) -> bool:
    token = (text or "").strip()
    if not token:
        return True
    if confidence < 40:
        return True
    cjk = count_cjk_chars(token)
    latin = sum(1 for char in token if char.isascii() and char.isalpha())
    symbols = sum(
        1
        for char in token
        if not (
            char.isalnum()
            or _is_cjk_char(char)
            or char in " .,;:%()[]-+/&'’\"“”《》【】、。|｜"
        )
    )
    if symbols >= 2 and symbols >= max(2, len(token) * 0.35):
        return True
    if cjk == 0 and latin >= 6:
        vowels = sum(1 for char in token.lower() if char in "aeiou")
        if vowels == 0:
            return True
    return False


def _is_noisy_ocr_line(text: str) -> bool:
    line = (text or "").strip()
    if not line:
        return True
    cjk = count_cjk_chars(line)
    latin = sum(1 for char in line if char.isascii() and char.isalpha())
    symbols = sum(
        1
        for char in line
        if not (char.isalnum() or _is_cjk_char(char) or char.isspace() or char in ".,;:%()[]-+/&'|｜《》【】、。")
    )
    if symbols >= 4 and symbols >= cjk + max(latin, 1):
        return True
    tokens = [part for part in re.split(r"\s+", line) if part]
    if not tokens:
        return True
    noisy = sum(1 for token in tokens if _is_noisy_ocr_token(token, 100))
    if cjk < 4 and noisy >= max(2, int(len(tokens) * 0.5)):
        return True
    return False


def _detect_page_gutter_x(image: Image.Image) -> float | None:
    """Find a vertical whitespace valley near page centre (typical EN | 中 layout)."""
    gray = image.convert("L") if image.mode != "L" else image
    width, height = gray.size
    if width < 240 or height < 160:
        return None
    target_w = min(width, 480)
    scale = target_w / float(width)
    target_h = max(48, int(height * scale))
    small = gray.resize((target_w, target_h), getattr(getattr(Image, "Resampling", Image), "BILINEAR", Image.BILINEAR))
    try:
        pixels = small.load()
        y0 = int(target_h * 0.12)
        y1 = max(y0 + 1, int(target_h * 0.88))
        band_h = y1 - y0
        window = max(3, target_w // 70)
        x0 = int(target_w * 0.28)
        x1 = int(target_w * 0.72)
        best: tuple[float, int] | None = None
        for thresh in (100, 120, 140):
            ink = [0.0] * target_w
            for x in range(target_w):
                count = 0
                for y in range(y0, y1):
                    if pixels[x, y] < thresh:
                        count += 1
                ink[x] = count / float(band_h)
            smooth = []
            for x in range(target_w):
                lo = max(0, x - window)
                hi = min(target_w, x + window + 1)
                smooth.append(sum(ink[lo:hi]) / (hi - lo))
            band = smooth[x0:x1]
            if not band:
                continue
            min_val = min(band)
            min_x = x0 + band.index(min_val)
            left_peak = max(smooth[:x0], default=0.0)
            right_peak = max(smooth[x1:], default=0.0)
            peak = min(left_peak, right_peak)
            if peak < 0.003:
                continue
            dip = (peak - min_val) / peak
            if dip < 0.32:
                continue
            run = 0
            for value in band:
                if value <= min_val + max(0.002, (peak - min_val) * 0.2):
                    run += 1
            if run < max(3, int(target_w * 0.012)):
                continue
            if best is None or dip > best[0]:
                best = (dip, min_x)
        if best is None:
            return None
        return best[1] / scale
    finally:
        if small is not gray and small is not image:
            small.close()


def _ocr_column_split_x(words: list[tuple], page_width: float) -> float | None:
    if len(words) < 10 or page_width <= 0:
        return None
    mids = sorted((float(word[0]) + float(word[2])) / 2 for word in words)
    best_gap = 0.0
    split_at = None
    center = page_width * 0.5
    for left_mid, right_mid in zip(mids, mids[1:]):
        gap_mid = (left_mid + right_mid) / 2
        if abs(gap_mid - center) > page_width * 0.22:
            continue
        gap = right_mid - left_mid
        if gap > best_gap:
            best_gap = gap
            split_at = gap_mid
    if split_at is None or best_gap < page_width * 0.06:
        return None
    left_count = sum(1 for word in words if (float(word[0]) + float(word[2])) / 2 < split_at)
    right_count = len(words) - left_count
    if left_count < 6 or right_count < 6:
        return None
    return split_at


def _ocr_column_line_entries(words: list[tuple]) -> list[tuple[float, str]]:
    entries: list[tuple[float, str]] = []
    for line in _group_words_into_lines(words):
        text = normalize_cjk_text(_join_word_cluster(line))
        if not text or _is_noisy_ocr_line(text):
            continue
        y_mid = sum((float(word[1]) + float(word[3])) / 2 for word in line) / len(line)
        entries.append((y_mid, text))
    return entries


def _ocr_words_to_rows(
    words: list[tuple],
    page_width: float,
    split_at: float | None = None,
) -> list[list[str]]:
    if not words:
        return []
    gutter = split_at if split_at is not None else _ocr_column_split_x(words, page_width)
    if gutter is not None:
        left_words = [word for word in words if (float(word[0]) + float(word[2])) / 2 < gutter]
        right_words = [word for word in words if (float(word[0]) + float(word[2])) / 2 >= gutter]
        if len(left_words) >= 3 and len(right_words) >= 3:
            left_entries = _ocr_column_line_entries(left_words)
            right_entries = _ocr_column_line_entries(right_words)
            heights = [max(1.0, float(word[3]) - float(word[1])) for word in words]
            heights.sort()
            y_tol = max(10.0, heights[len(heights) // 2] * 1.05)
            rows: list[list[str]] = []
            left_index = 0
            right_index = 0
            while left_index < len(left_entries) or right_index < len(right_entries):
                if left_index >= len(left_entries):
                    rows.append(["", right_entries[right_index][1]])
                    right_index += 1
                    continue
                if right_index >= len(right_entries):
                    rows.append([left_entries[left_index][1], ""])
                    left_index += 1
                    continue
                left_y, left_text = left_entries[left_index]
                right_y, right_text = right_entries[right_index]
                if abs(left_y - right_y) <= y_tol:
                    rows.append([left_text, right_text])
                    left_index += 1
                    right_index += 1
                elif left_y < right_y:
                    rows.append([left_text, ""])
                    left_index += 1
                else:
                    rows.append(["", right_text])
                    right_index += 1
            return rows
    lines = _group_words_into_lines(words)
    threshold = max(_column_gap_threshold(lines, page_width), page_width * 0.1)
    rows = []
    for line in lines:
        parts = [normalize_cjk_text(part) for part in _split_line_by_x_gaps(line, threshold)]
        parts = [part for part in parts if part and not _is_noisy_ocr_line(part)]
        if parts:
            rows.append(parts)
    return _pad_table_rows(rows)


def _ocr_rows_to_plain_text(rows: list[list[str]]) -> str:
    lines: list[str] = []
    for row in rows or []:
        cells = [normalize_cjk_text(str(cell or "")).strip() for cell in row]
        if len(cells) >= 2 and cells[0] and cells[1]:
            lines.append(f"{cells[0]}  |  {cells[1]}")
        else:
            filled = [cell for cell in cells if cell]
            if filled:
                lines.append(filled[0] if len(filled) == 1 else "  |  ".join(filled))
    return "\n".join(lines)


def _ocr_words_to_plain_text(words: list[tuple], page_width: float) -> str:
    return _ocr_rows_to_plain_text(_ocr_words_to_rows(words, page_width))


def _tesseract_data_to_words(data: dict) -> list[tuple]:
    texts = data.get("text") or []
    words: list[tuple] = []
    for index, raw in enumerate(texts):
        text = str(raw or "").strip()
        try:
            confidence = float((data.get("conf") or ["-1"])[index])
        except Exception:
            confidence = 0.0
        if not text or _is_noisy_ocr_token(text, confidence):
            continue
        try:
            left = float(data["left"][index])
            top = float(data["top"][index])
            width = float(data["width"][index])
            height = float(data["height"][index])
        except Exception:
            continue
        words.append((left, top, left + width, top + height, text, 0, 0, index))
    return words


def _ocr_image_to_words(image: Image.Image, tess_lang: str, tess_config: str) -> list[tuple]:
    output_type = getattr(getattr(pytesseract, "Output", None), "DICT", "dict")
    try:
        data = pytesseract.image_to_data(
            image,
            lang=tess_lang,
            config=tess_config,
            output_type=output_type,
        )
    except Exception as exc:
        detail = str(exc)
        if "traineddata" in detail or "Failed loading language" in detail:
            raise ValueError(
                f"Tesseract 找不到語言包「{tess_lang}」。請確認 tessdata 內有 chi_sim / chi_tra。"
            ) from exc
        return []
    if isinstance(data, dict):
        return _tesseract_data_to_words(data)
    return []


def _offset_ocr_words(words: list[tuple], dx: float, dy: float) -> list[tuple]:
    shifted: list[tuple] = []
    for word in words:
        left, top, right, bottom, text, *rest = word
        shifted.append((left + dx, top + dy, right + dx, bottom + dy, text, *rest))
    return shifted


def _ocr_bilingual_column_words(
    image: Image.Image,
    gutter: float,
    tess_lang: str,
    tess_config: str,
) -> list[tuple]:
    width, height = image.size
    pad = max(8, int(width * 0.012))
    gutter_x = max(1, min(width - 1, int(gutter)))
    col_config = _tesseract_column_config(tess_config)
    left_box = (0, 0, min(width, gutter_x + pad), height)
    right_origin = max(0, gutter_x - pad)
    right_box = (right_origin, 0, width, height)
    left_img = image.crop(left_box)
    right_img = image.crop(right_box)
    try:
        left_words = _ocr_image_to_words(left_img, _column_tess_lang(tess_lang, prefer_cjk=False), col_config)
        right_words = _ocr_image_to_words(right_img, _column_tess_lang(tess_lang, prefer_cjk=True), col_config)
    finally:
        left_img.close()
        right_img.close()
    left_kept = [
        word
        for word in left_words
        if (float(word[0]) + float(word[2])) / 2 < gutter
    ]
    right_kept = [
        word
        for word in _offset_ocr_words(right_words, right_origin, 0)
        if (float(word[0]) + float(word[2])) / 2 >= gutter
    ]
    if len(left_kept) < 4 or len(right_kept) < 4:
        return []
    return left_kept + right_kept


def _ocr_image_to_rows(image: Image.Image, tess_lang: str, tess_config: str) -> list[list[str]]:
    prepared = _prepare_ocr_image(image)
    try:
        raw_gutter = _detect_page_gutter_x(image)
        if raw_gutter is not None and image.width:
            gutter = raw_gutter * (prepared.width / float(image.width))
        else:
            gutter = _detect_page_gutter_x(prepared)
        words: list[tuple] = []
        if gutter is not None:
            words = _ocr_bilingual_column_words(prepared, gutter, tess_lang, tess_config)
        if not words:
            words = _ocr_image_to_words(prepared, tess_lang, tess_config)
        if words:
            rows = _ocr_words_to_rows(words, float(prepared.width), split_at=gutter)
            if rows:
                return rows
        try:
            text = pytesseract.image_to_string(prepared, lang=tess_lang, config=tess_config)
        except Exception as exc:
            detail = str(exc)
            if "traineddata" in detail or "Failed loading language" in detail:
                raise ValueError(
                    f"Tesseract 找不到語言包「{tess_lang}」。請確認 tessdata 內有 chi_sim / chi_tra。"
                ) from exc
            text = ""
        rows = []
        for line in (text or "").splitlines():
            cleaned = normalize_cjk_text(line).strip()
            if cleaned and not _is_noisy_ocr_line(cleaned):
                rows.append([cleaned])
        return rows
    finally:
        if prepared is not image:
            prepared.close()


def _ocr_indexes_to_tables(
    source: Path,
    password: str,
    page_indexes: list[int],
    language: str,
    dpi: int,
    progress=None,
) -> list[tuple[int, list[list[str]]]]:
    if not page_indexes:
        return []
    pages_spec = ",".join(str(index + 1) for index in page_indexes)
    pages, tess_lang, tess_config = _ocr_iter_pages(source, password, pages_spec, language, dpi, progress)
    sheets: list[tuple[int, list[list[str]]]] = []
    count = 0
    for page_index, image in pages:
        count += 1
        if progress:
            progress(count - 1, 0, f"掃描件 OCR 第 {page_index + 1} 頁，請稍候…")
        try:
            rows = _ocr_image_to_rows(image, tess_lang, tess_config)
        finally:
            image.close()
        if rows:
            sheets.append((page_index, rows))
    return sheets


def _preferred_xlsx_method(methods: list[str]) -> str:
    for name in ("tables", "columns", "ocr", "text"):
        if name in methods:
            return name
    return "text"


def _collect_xlsx_sheets(
    source: Path,
    password: str = "",
    pages_spec: str = "",
    language: str = "auto",
    dpi: int = 200,
    progress=None,
) -> tuple[list[tuple[int, list[list[str]]]], str]:
    indexes = _page_indexes_or_all(source, password, pages_spec)
    sheets: list[tuple[int, list[list[str]]]] = []
    methods: list[str] = []
    ocr_indexes: list[int] = []
    if PYMUPDF_AVAILABLE and fitz is not None:
        document = fitz.open(str(source))
        try:
            if document.is_encrypted:
                if not password or document.authenticate(password) == 0:
                    raise ValueError(f"{source.name} 已加密，請輸入密碼。")
            total = len(indexes)
            for offset, page_index in enumerate(indexes):
                if progress:
                    progress(
                        offset + 1,
                        total,
                        f"正在檢查第 {page_index + 1} / {total} 頁表格…",
                    )
                page = document[page_index]
                page_tables = _pymupdf_page_tables(page)
                if page_tables:
                    sheets.extend((page_index, rows) for rows in page_tables)
                    methods.append("tables")
                    continue
                clustered = _page_words_to_table(page)
                if clustered:
                    sheets.append((page_index, clustered))
                    methods.append("columns")
                    continue
                text_rows = _page_text_to_rows(page)
                if text_rows and not _page_looks_scanned(page):
                    sheets.append((page_index, text_rows))
                    methods.append("text")
                    continue
                ocr_indexes.append(page_index)
        finally:
            document.close()
    else:
        reader = open_reader(source, password)
        total = len(indexes)
        for offset, page_index in enumerate(indexes):
            if progress:
                progress(
                    offset + 1,
                    total,
                    f"正在檢查第 {page_index + 1} / {total} 頁文字…",
                )
            text = normalize_cjk_text(reader.pages[page_index].extract_text() or "")
            rows = _text_lines_to_rows(text.splitlines())
            if rows and (len(text.strip()) >= 8 or count_cjk_chars(text) >= 4):
                sheets.append((page_index, rows))
                methods.append("text")
            else:
                ocr_indexes.append(page_index)
    if ocr_indexes:
        try:
            ocr_sheets = _ocr_indexes_to_tables(source, password, ocr_indexes, language, dpi, progress)
        except Exception:
            if not sheets:
                raise
            ocr_sheets = []
        if ocr_sheets:
            sheets.extend(ocr_sheets)
            methods.append("ocr")
    sheets.sort(key=lambda item: item[0])
    return sheets, _preferred_xlsx_method(methods)


def _extract_pdf_tables(
    source: Path,
    password: str = "",
    pages_spec: str = "",
) -> list[tuple[int, list[list[str]]]]:
    sheets, _method = _collect_xlsx_sheets(source, password, pages_spec)
    return [(page_index, rows) for page_index, rows in sheets if _is_useful_table(rows)]


def _xlsx_sheet_name(page_index: int, table_number: int, used_names: set[str]) -> str:
    base = f"第{page_index + 1}頁" if table_number == 1 else f"第{page_index + 1}頁_表{table_number}"
    name = base[:31]
    suffix = 2
    while name in used_names:
        name = f"{base[:28]}_{suffix}"[:31]
        suffix += 1
    used_names.add(name)
    return name


def _excel_cell_value(value: str):
    text = (value or "").strip()
    if not text:
        return ""
    negative = text.startswith("(") and text.endswith(")")
    body = text[1:-1].strip() if negative else text
    percent = body.endswith("%")
    if percent:
        body = body[:-1].strip()
    compact = body.replace(",", "").replace("，", "").replace(" ", "")
    if not re.fullmatch(r"-?\d+(?:\.\d+)?", compact):
        return text
    if re.fullmatch(r"0\d+", compact):
        return text
    number = float(compact)
    if negative:
        number = -number
    if percent:
        return number / 100.0
    if "." not in compact and number == int(number):
        return int(number)
    return number


def _write_xlsx_tables(tables: list[tuple[int, list[list[str]]]], target: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    first = True
    used_names: set[str] = set()
    table_numbers: dict[int, int] = {}
    for page_index, rows in tables:
        table_numbers[page_index] = table_numbers.get(page_index, 0) + 1
        name = _xlsx_sheet_name(page_index, table_numbers[page_index], used_names)
        sheet = workbook.active if first else workbook.create_sheet(title=name)
        if first:
            sheet.title = name
            first = False
        for row_index, row in enumerate(rows, start=1):
            for column_index, value in enumerate(row, start=1):
                sheet.cell(row_index, column_index, _excel_cell_value(value))
        for column_index in range(1, (max((len(row) for row in rows), default=0) + 1)):
            longest = 8
            for row in rows:
                if column_index - 1 < len(row):
                    longest = max(longest, min(42, len(str(row[column_index - 1] or ""))))
            sheet.column_dimensions[get_column_letter(column_index)].width = longest + 2
    target.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(str(target))


def _write_xlsx_from_text(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    lines: list[str] | None = None,
) -> None:
    if lines is None:
        indexes = _page_indexes_or_all(source, password, pages_spec)
        reader = open_reader(source, password)
        collected: list[str] = []
        for page_index in indexes:
            collected.append(f"--- Page {page_index + 1} ---")
            collected.extend((reader.pages[page_index].extract_text() or "").splitlines())
        rows = _text_lines_to_rows(collected)
    else:
        rows = _text_lines_to_rows(lines)
    _write_xlsx_tables([(0, rows or [["（沒有可抽出的文字）"]])], target)


def pdf_to_xlsx(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    language: str = "auto",
    dpi: int = 200,
    progress=None,
) -> tuple[int, str]:
    if not _xlsx_available():
        raise ValueError("PDF 轉 Excel 需要安裝 openpyxl。")
    indexes = _page_indexes_or_all(source, password, pages_spec)
    if not indexes:
        raise ValueError("沒有可轉換的頁面。")
    if progress:
        progress(0, len(indexes), "正在檢查 PDF 表格／文字層…")
    unlocked = _write_unlocked_pdf_copy(source, password)
    try:
        sheets, method = _collect_xlsx_sheets(unlocked, "", pages_spec, language, dpi, progress)
        if not sheets:
            raise ValueError("無法從這份 PDF 抽出表格或文字。掃描件請確認已安裝 Tesseract OCR。")
        _write_xlsx_tables(sheets, target)
        return len(indexes), method
    finally:
        if unlocked != source:
            unlocked.unlink(missing_ok=True)


def write_pdf_info(source: Path, target: Path, password: str = "") -> None:
    reader = open_reader(source, password)
    metadata = reader.metadata or {}
    lines = [
        "Victor PDF Tools Box - PDF Info",
        f"File: {source.name}",
        f"Pages: {len(reader.pages)}",
        f"Encrypted: {reader.is_encrypted}",
        "",
        "Metadata:",
    ]
    for key, value in metadata.items():
        lines.append(f"{key}: {value}")
    target.write_text("\n".join(lines), encoding="utf-8")


def ensure_pymupdf_available() -> None:
    if not PYMUPDF_AVAILABLE or fitz is None:
        raise ValueError("無痕文字替換需要 PyMuPDF（pymupdf）。")


def normalize_font_key(font_name: str) -> str:
    normalized = re.sub(r"^[A-Z0-9]{6}\+", "", font_name or "", flags=re.IGNORECASE)
    return re.sub(r"[^a-z0-9]+", "", normalized.lower())


def int_color_to_rgb(color: int) -> tuple[float, float, float]:
    if color == 0:
        return (0.0, 0.0, 0.0)
    return (
        ((color >> 16) & 255) / 255.0,
        ((color >> 8) & 255) / 255.0,
        (color & 255) / 255.0,
    )


WINDOWS_FONT_CANDIDATES: dict[str, tuple[str, ...]] = {
    "arial": ("arial.ttf", "arialbd.ttf", "ariali.ttf", "arialbi.ttf"),
    "helvetica": ("arial.ttf", "arialbd.ttf", "ariali.ttf", "arialbi.ttf"),
    "timesnewroman": ("times.ttf", "timesbd.ttf", "timesi.ttf", "timesbi.ttf"),
    "timesroman": ("times.ttf", "timesbd.ttf", "timesi.ttf", "timesbi.ttf"),
    "courier": ("cour.ttf", "courbd.ttf", "couri.ttf", "courbi.ttf"),
    "couriernew": ("cour.ttf", "courbd.ttf", "couri.ttf", "courbi.ttf"),
    "microsoftyahei": ("msyh.ttc", "msyhbd.ttc"),
    "yahei": ("msyh.ttc", "msyhbd.ttc"),
    "mingliu": ("mingliu.ttc", "mingliub.ttc"),
    "pmingliu": ("mingliu.ttc", "mingliub.ttc"),
    "simsun": ("simsun.ttc", "simsunb.ttf"),
    "nsimsun": ("simsun.ttc", "simsunb.ttf"),
    "simhei": ("simhei.ttf",),
    "calibri": ("calibri.ttf", "calibrib.ttf", "calibrii.ttf", "calibriz.ttf"),
    "tahoma": ("tahoma.ttf", "tahomabd.ttf"),
    "verdana": ("verdana.ttf", "verdanab.ttf", "verdanai.ttf", "verdanaz.ttf"),
}


def _font_file_variant_index(bold: bool, italic: bool) -> int:
    if bold and italic:
        return 3
    if bold:
        return 1
    if italic:
        return 2
    return 0


def resolve_system_font_file(font_name: str, font_flags: int = 0) -> str:
    key = normalize_font_key(font_name)
    candidates = None
    for pattern, files in WINDOWS_FONT_CANDIDATES.items():
        if pattern in key or key in pattern:
            candidates = files
            break
    if candidates is None:
        if any(
            token in key
            for token in (
                "cjk",
                "noto",
                "sourcehan",
                "sourcehansans",
                "song",
                "ming",
                "kai",
                "hei",
                "yahei",
                "jhenghei",
                "gothic",
                "mincho",
                "pingfang",
                "simsun",
            )
        ):
            candidates = ("msjh.ttc", "msyh.ttc", "msjhbd.ttc", "msyhbd.ttc", "mingliu.ttc", "simsun.ttc")
        else:
            candidates = WINDOWS_FONT_CANDIDATES["arial"]
    bold = bool(font_flags & 16)
    italic = bool(font_flags & 2)
    index = _font_file_variant_index(bold, italic)
    fonts_dir = Path("C:/Windows/Fonts")
    for offset in (0, 1, 2, 3):
        candidate_index = min(index + offset, len(candidates) - 1)
        path = fonts_dir / candidates[candidate_index]
        if path.exists():
            return str(path)
    fallback = fonts_dir / candidates[0]
    if fallback.exists():
        return str(fallback)
    return ""


def match_page_font_name(span_font: str, page_fonts: list[tuple]) -> str:
    span_key = normalize_font_key(span_font)
    if not span_key:
        return ""
    best_name = ""
    best_score = 0
    for _xref, _ext, _kind, basefont, name, _encoding in page_fonts:
        for candidate in (basefont, name):
            candidate_key = normalize_font_key(str(candidate))
            if not candidate_key:
                continue
            if candidate_key == span_key:
                return str(name)
            if candidate_key in span_key or span_key in candidate_key:
                score = min(len(candidate_key), len(span_key))
                if score > best_score:
                    best_score = score
                    best_name = str(name)
    return best_name


def _pymupdf_text_flags() -> int:
    flags = int(getattr(fitz, "TEXTFLAGS_DICT", 0) or 0)
    clip = int(getattr(fitz, "TEXT_MEDIABOX_CLIP", 0) or 0)
    if flags and clip:
        return flags & ~clip
    return 0


def _span_text(span: dict) -> str:
    text = (span.get("text") or "").replace("\u00a0", " ")
    if text.strip():
        return text.strip()
    chars = span.get("chars") or []
    assembled = "".join(str(item.get("c") or "") for item in chars)
    return assembled.replace("\u00a0", " ").strip()


def span_to_text_block(span: dict, page_fonts: list[tuple]) -> TextBlock:
    text = _span_text(span)
    bbox = tuple(float(value) for value in span.get("bbox", (0, 0, 0, 0)))
    origin = span.get("origin") or (bbox[0], bbox[3])
    font_name = str(span.get("font") or "")
    font_flags = int(span.get("flags") or 0)
    font_size = float(span.get("size") or 12)
    page_font_name = match_page_font_name(font_name, page_fonts)
    font_file = resolve_system_font_file(font_name, font_flags)
    width = max(bbox[2] - bbox[0], font_size * 2.0)
    height = max(bbox[3] - bbox[1], font_size * 1.2)
    return TextBlock(
        text=text,
        x=float(origin[0]),
        y=float(origin[1]),
        width=width,
        height=height,
        font_size=font_size,
        font_name=font_name,
        color_rgb=int_color_to_rgb(int(span.get("color") or 0)),
        bbox=bbox,
        font_flags=font_flags,
        page_font_name=page_font_name,
        font_file=font_file,
    )


def _text_blocks_from_pymupdf_dict(page_data: dict, page_fonts: list[tuple]) -> list[TextBlock]:
    blocks: list[TextBlock] = []
    for block in page_data.get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                item = span_to_text_block(span, page_fonts)
                if item.text:
                    blocks.append(item)
    return blocks


def _text_blocks_from_fitz_page(page) -> list[TextBlock]:
    page_fonts = page.get_fonts()
    flags = _pymupdf_text_flags()
    blocks = _text_blocks_from_pymupdf_dict(page.get_text("dict", flags=flags), page_fonts)
    if not blocks:
        blocks = _text_blocks_from_pymupdf_dict(page.get_text("rawdict", flags=flags), page_fonts)
    if not blocks:
        for item in page.get_text("words", flags=flags):
            x0, y0, x1, y1, word, *_rest = item
            text = (word or "").strip()
            if not text:
                continue
            height = max(float(y1) - float(y0), 8.0)
            width = max(float(x1) - float(x0), 8.0)
            blocks.append(
                TextBlock(
                    text,
                    float(x0),
                    float(y1),
                    width,
                    height,
                    height,
                    bbox=(float(x0), float(y0), float(x1), float(y1)),
                )
            )
    if not blocks:
        plain = (page.get_text("text") or "").strip()
        if plain:
            rect = page.rect
            blocks.append(
                TextBlock(
                    plain,
                    float(rect.x0),
                    float(rect.y1),
                    float(rect.width),
                    float(rect.height),
                    12.0,
                    bbox=(float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)),
                )
            )
    return merge_text_blocks(blocks)


def extract_page_text_blocks_pymupdf(source: Path, page_index: int, password: str = "") -> list[TextBlock]:
    ensure_pymupdf_available()
    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        if page_index < 0 or page_index >= document.page_count:
            raise ValueError("頁碼超出範圍。")
        return _text_blocks_from_fitz_page(document[page_index])
    finally:
        document.close()


def extract_page_text_blocks_pypdf(source: Path, page_index: int, password: str = "") -> list[TextBlock]:
    reader = open_reader(source, password)
    if page_index < 0 or page_index >= len(reader.pages):
        raise ValueError("頁碼超出範圍。")

    blocks: list[TextBlock] = []

    def visitor_text(text, _cm, tm, font_dict, font_size):
        normalized = (text or "").strip()
        if not normalized:
            return
        x = float(tm[4])
        y = float(tm[5])
        size = float(font_size or 12)
        width = max(len(normalized) * size * 0.55, size * 2.0)
        height = max(size * 1.5, 12.0)
        font_name = ""
        if font_dict:
            font_name = str(font_dict.get("/BaseFont", ""))
        blocks.append(TextBlock(normalized, x, y, width, height, size, font_name))

    reader.pages[page_index].extract_text(visitor_text=visitor_text)
    return merge_text_blocks(blocks)


def extract_page_text_blocks(source: Path, page_index: int, password: str = "") -> list[TextBlock]:
    if PYMUPDF_AVAILABLE:
        try:
            blocks = extract_page_text_blocks_pymupdf(source, page_index, password)
            if blocks:
                return blocks
        except ValueError:
            raise
        except Exception:
            pass
    return extract_page_text_blocks_pypdf(source, page_index, password)


def text_block_redaction_rect(block: TextBlock, replacement: str = "") -> "fitz.Rect":
    ensure_pymupdf_available()
    padding = max(block.font_size * 0.15, 1.5)
    if block.bbox != (0.0, 0.0, 0.0, 0.0):
        rect = fitz.Rect(block.bbox)
    else:
        rect = fitz.Rect(
            block.x,
            block.y - block.height,
            block.x + block.width,
            block.y + max(block.font_size * 0.35, 4.0),
        )
    rect = rect + (-padding, -padding, padding, padding)
    if replacement:
        extra_width = max(len(replacement) - len(block.text.strip()), 0) * block.font_size * 0.55
        if extra_width > 0:
            rect.x1 += extra_width + padding
    return rect


def _is_pdf_internal_font_id(font_name: str) -> bool:
    key = (font_name or "").strip()
    if not key:
        return True
    if "+" in key:
        return True
    if len(key) <= 4 and key[:1].isalpha() and key[1:].isdigit():
        return True
    return False


def _font_file_looks_cjk(font_file: str) -> bool:
    name = Path(font_file).name.lower()
    return any(token in name for token in ("msyh", "msjh", "mingliu", "simsun", "simhei", "noto", "cjk", "sourcehan"))


def _insert_text_with_block_style(page, block: TextBlock, replacement: str, rect: "fitz.Rect") -> None:
    fontsize = max(block.font_size, 6.0)
    color = block.color_rgb
    kwargs: dict = {"fontsize": fontsize, "color": color, "align": fitz.TEXT_ALIGN_LEFT}
    bold = bool(block.font_flags & 16) or "bold" in (block.font_name or "").lower()
    need_cjk = overlay_needs_embedded_font(replacement) or overlay_needs_embedded_font(block.text)
    fontfile = ""
    if need_cjk:
        if block.font_file and _font_file_looks_cjk(block.font_file):
            fontfile = block.font_file
        else:
            fontfile = resolve_cjk_font_file(bold) or ""
        if fontfile:
            kwargs["fontfile"] = fontfile
            kwargs["fontname"] = "edit-cjk"
        else:
            kwargs["fontname"] = "china-s"
    elif block.font_file:
        kwargs["fontfile"] = block.font_file
        kwargs["fontname"] = "edit-latin"
    elif block.page_font_name and not _is_pdf_internal_font_id(block.page_font_name):
        kwargs["fontname"] = block.page_font_name
    else:
        kwargs["fontname"] = "helv"

    insert_kwargs = {key: value for key, value in kwargs.items() if key != "align"}
    if "\n" in replacement:
        overflow = page.insert_textbox(rect, replacement, **kwargs)
        if overflow < 0:
            page.insert_text(
                fitz.Point(rect.x0, min(rect.y1 - 2, rect.y0 + fontsize)),
                replacement,
                **insert_kwargs,
            )
        return
    page.insert_text(fitz.Point(rect.x0 + 0.5, rect.y0 + fontsize * 0.85), replacement, **insert_kwargs)


def replace_text_block_seamless(
    source: Path,
    target: Path,
    page_index: int,
    block: TextBlock,
    replacement: str,
    password: str = "",
) -> None:
    """Remove the original text run and rewrite with matched font/size/color."""

    if not replacement.strip():
        raise ValueError("請輸入替換文字。")
    ensure_pymupdf_available()

    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        if page_index < 0 or page_index >= document.page_count:
            raise ValueError("頁碼超出範圍。")

        page = document[page_index]
        replacement_text = replacement.strip()
        rect = text_block_redaction_rect(block, replacement_text)

        needle = block.text.strip()
        if needle:
            matches = page.search_for(needle)
            if matches:
                rect = matches[0]
                if replacement_text and len(replacement_text) > len(needle):
                    extra = (len(replacement_text) - len(needle)) * block.font_size * 0.55
                    rect.x1 += extra + max(block.font_size * 0.15, 1.5)

        page.add_redact_annot(rect, fill=(1, 1, 1))
        page.apply_redactions()

        _insert_text_with_block_style(page, block, replacement_text, rect)
        if hasattr(document, "subset_fonts"):
            document.subset_fonts()
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()


def replace_text_block_overlay(
    source: Path,
    target: Path,
    page_index: int,
    block: TextBlock,
    replacement: str,
    password: str = "",
) -> None:
    if not replacement.strip():
        raise ValueError("請輸入替換文字。")
    bold = "bold" in (block.font_name or "").lower()
    add_text_overlay_annotation(
        source=source,
        target=target,
        page_index=page_index,
        x=block.x,
        y=block.y,
        text=replacement,
        font_size=max(int(round(block.font_size)), 8),
        cover_original=True,
        cover_width=max(block.width + block.font_size * 0.8, len(replacement) * block.font_size * 0.65),
        cover_height=max(block.height, block.font_size * 1.8),
        password=password,
        font_key=font_key_for_pdf_font(block.font_name),
        bold=bold,
        color_rgb=block.color_rgb,
    )


def redact_text_block_overlay(
    source: Path,
    target: Path,
    page_index: int,
    block: TextBlock,
    password: str = "",
) -> None:
    reader = open_reader(source, password)
    if page_index < 0 or page_index >= len(reader.pages):
        raise ValueError("頁碼超出範圍。")

    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    page = writer.pages[page_index]
    add_annotation_to_page(writer, page, redaction_annotation_for_block(block))
    write_pdf(writer, target)


@dataclass(frozen=True)
class EraseMark:
    page_index: int
    kind: str
    color_rgb: tuple[float, float, float]
    points: tuple[tuple[float, float], ...]
    radius: float = 16.0


def apply_erase_marks(
    source: Path,
    target: Path,
    marks: list[EraseMark],
    password: str = "",
    remove_content: bool = True,
) -> int:
    if not marks:
        raise ValueError("請先用橡皮刷或範圍遮擋標記要處理的位置。")
    ensure_pymupdf_available()
    document = None
    try:
        document = fitz.open(str(source))
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        grouped: dict[int, list[EraseMark]] = {}
        for mark in marks:
            if mark.page_index < 0 or mark.page_index >= document.page_count:
                raise ValueError("頁碼超出範圍。")
            grouped.setdefault(mark.page_index, []).append(mark)
        applied = 0
        for page_index, page_marks in grouped.items():
            page = document[page_index]
            height = float(page.rect.height)

            def to_fitz_rect(x0: float, y0: float, x1: float, y1: float) -> "fitz.Rect":
                return fitz.Rect(
                    min(x0, x1),
                    height - max(y0, y1),
                    max(x0, x1),
                    height - min(y0, y1),
                )

            if remove_content:
                for mark in page_marks:
                    fill = mark.color_rgb
                    if mark.kind == "rect" and len(mark.points) >= 2:
                        (x0, y0), (x1, y1) = mark.points[0], mark.points[1]
                        page.add_redact_annot(to_fitz_rect(x0, y0, x1, y1), fill=fill)
                        applied += 1
                    else:
                        radius = max(float(mark.radius), 2.0)
                        for pdf_x, pdf_y in mark.points:
                            center_x, center_y = pdf_x, height - pdf_y
                            page.add_redact_annot(
                                fitz.Rect(
                                    center_x - radius,
                                    center_y - radius,
                                    center_x + radius,
                                    center_y + radius,
                                ),
                                fill=fill,
                            )
                            applied += 1
                page.apply_redactions()
            for mark in page_marks:
                fill = mark.color_rgb
                if mark.kind == "rect" and len(mark.points) >= 2:
                    (x0, y0), (x1, y1) = mark.points[0], mark.points[1]
                    page.draw_rect(to_fitz_rect(x0, y0, x1, y1), color=fill, fill=fill, width=0)
                    if not remove_content:
                        applied += 1
                else:
                    radius = max(float(mark.radius), 2.0)
                    for pdf_x, pdf_y in mark.points:
                        page.draw_circle(
                            fitz.Point(pdf_x, height - pdf_y),
                            radius,
                            color=fill,
                            fill=fill,
                            width=0,
                        )
                        if not remove_content:
                            applied += 1
        temp_target = target.with_name(f"{target.stem}.erase-tmp{target.suffix}")
        if hasattr(document, "subset_fonts"):
            document.subset_fonts()
        document.save(str(temp_target), garbage=4, deflate=True)
        document.close()
        document = None
        temp_target.replace(target)
        return applied
    finally:
        if document is not None:
            document.close()


def apply_erase_then_text_overlays(
    source: Path,
    target: Path,
    marks: list[EraseMark],
    overlays: list[dict],
    password: str = "",
    remove_content: bool = True,
) -> int:
    """Erase first, then place text boxes, writing once to target."""

    if not marks and not overlays:
        raise ValueError("請先用橡皮刷遮擋，或加入文字方塊。")
    source = Path(source)
    target = Path(target)
    if marks and overlays:
        handle = tempfile.NamedTemporaryFile(suffix=".pdf", prefix="erase_text_", delete=False)
        handle.close()
        temp = Path(handle.name)
        try:
            applied = apply_erase_marks(source, temp, marks, password, remove_content)
            add_text_overlay_annotations(temp, target, overlays, password="")
            return applied
        finally:
            try:
                temp.unlink(missing_ok=True)
            except OSError:
                pass
    if marks:
        return apply_erase_marks(source, target, marks, password, remove_content)
    add_text_overlay_annotations(source, target, overlays, password)
    return 0


def redaction_annotation_for_block(block: TextBlock) -> DictionaryObject:
    padding = max(block.font_size * 0.2, 2.0)
    return DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/Square"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(block.x - padding),
                    FloatObject(block.y - block.height - padding),
                    FloatObject(block.x + block.width + padding),
                    FloatObject(block.y + padding),
                ]
            ),
            NameObject("/C"): ArrayObject([FloatObject(0), FloatObject(0), FloatObject(0)]),
            NameObject("/IC"): ArrayObject([FloatObject(0), FloatObject(0), FloatObject(0)]),
            NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(0)]),
            NameObject("/F"): NumberObject(4),
        }
    )


def redact_matching_text_blocks_overlay(
    source: Path,
    target: Path,
    query: str,
    password: str = "",
    case_sensitive: bool = False,
    whole_word: bool = False,
) -> int:
    if not (query or "").strip():
        raise ValueError("請輸入要批量遮蔽的搜尋文字。")

    reader = open_reader(source, password)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    match_count = 0
    for page_index, page in enumerate(writer.pages):
        blocks = extract_page_text_blocks(source, page_index, password)
        for block in blocks:
            if text_matches_query(block.text, query, case_sensitive, whole_word):
                add_annotation_to_page(writer, page, redaction_annotation_for_block(block))
                match_count += 1
    if match_count == 0:
        raise ValueError(f"找不到要遮蔽的文字：{query.strip()}")
    write_pdf(writer, target)
    return match_count


def text_matches_query(text: str, query: str, case_sensitive: bool = False, whole_word: bool = False) -> bool:
    needle = (query or "").strip()
    if not needle:
        return False
    haystack = text or ""
    if whole_word:
        flags = 0 if case_sensitive else re.IGNORECASE
        return re.search(rf"(?<!\w){re.escape(needle)}(?!\w)", haystack, flags) is not None
    if case_sensitive:
        return needle in haystack
    return needle.lower() in haystack.lower()


def validate_content_stream_replacement(original: str, replacement: str) -> tuple[str, str]:
    old_text = re.sub(r"\s+", " ", original or "").strip()
    new_text = replacement.strip()
    if not old_text:
        raise ValueError("沒有可替換的原文字。")
    if not new_text:
        raise ValueError("請輸入替換文字。")
    if not re.fullmatch(r"[\x20-\x7E]+", old_text) or not re.fullmatch(r"[\x20-\x7E]+", new_text):
        raise ValueError("直接改內容流目前只支援簡單英文、數字和半形符號。")
    if len(old_text) != len(new_text):
        raise ValueError("直接改內容流目前要求新舊文字長度相同，以避免版面走位。")
    return old_text, new_text


def replace_text_in_show_operand(operand, old_text: str, new_text: str) -> tuple[object, bool]:
    if isinstance(operand, TextStringObject) and str(operand) == old_text:
        return TextStringObject(new_text), True
    if isinstance(operand, ByteStringObject):
        try:
            decoded = bytes(operand).decode("latin-1")
        except Exception:
            return operand, False
        if decoded == old_text:
            return ByteStringObject(new_text.encode("latin-1")), True
    return operand, False


def replace_text_in_array_operand(operand, old_text: str, new_text: str) -> tuple[object, bool]:
    if not isinstance(operand, ArrayObject):
        return operand, False
    string_parts = [part for part in operand if isinstance(part, (TextStringObject, ByteStringObject))]
    joined = "".join(str(part) if isinstance(part, TextStringObject) else bytes(part).decode("latin-1") for part in string_parts)
    if joined != old_text:
        return operand, False

    replaced = ArrayObject()
    offset = 0
    for part in operand:
        if isinstance(part, TextStringObject):
            length = len(str(part))
            replaced.append(TextStringObject(new_text[offset : offset + length]))
            offset += length
        elif isinstance(part, ByteStringObject):
            length = len(bytes(part))
            replaced.append(ByteStringObject(new_text[offset : offset + length].encode("latin-1")))
            offset += length
        else:
            replaced.append(part)
    return replaced, True


def replace_text_block_content_stream(
    source: Path,
    target: Path,
    page_index: int,
    block: TextBlock,
    replacement: str,
    password: str = "",
) -> None:
    old_text, new_text = validate_content_stream_replacement(block.text, replacement)
    replace_text_in_content_stream(source, target, page_index, old_text, new_text, password)


def replace_text_in_content_stream(
    source: Path,
    target: Path,
    page_index: int,
    old_text: str,
    new_text: str,
    password: str = "",
) -> None:
    reader = open_reader(source, password)
    if page_index < 0 or page_index >= len(reader.pages):
        raise ValueError("頁碼超出範圍。")

    writer = PdfWriter()
    replacements = 0
    for index, page in enumerate(reader.pages):
        page_copy = copy(page)
        if index == page_index:
            contents = page_copy.get_contents()
            if contents is None:
                raise ValueError("此頁沒有可修改的內容流。")
            content = ContentStream(contents, reader)
            new_operations = []
            for operands, operator in content.operations:
                operands = list(operands)
                replaced = False
                if operator in {b"Tj", b"'"} and operands:
                    operands[0], replaced = replace_text_in_show_operand(operands[0], old_text, new_text)
                elif operator == b'"' and len(operands) >= 3:
                    operands[2], replaced = replace_text_in_show_operand(operands[2], old_text, new_text)
                elif operator == b"TJ" and operands:
                    operands[0], replaced = replace_text_in_array_operand(operands[0], old_text, new_text)
                if replaced:
                    replacements += 1
                new_operations.append((operands, operator))
            if replacements != 1:
                raise ValueError("找不到唯一可直接修改的文字；請改用覆蓋方式。")
            content.operations = new_operations
            stream = DecodedStreamObject()
            stream.set_data(content.get_data())
            page_copy[NameObject("/Contents")] = stream
        writer.add_page(page_copy)
    write_pdf(writer, target)


def redact_text_block_secure(
    source: Path,
    target: Path,
    page_index: int,
    block: TextBlock,
    password: str = "",
) -> None:
    old_text = re.sub(r"\s+", " ", block.text or "").strip()
    if not old_text:
        raise ValueError("沒有可遮蔽的原文字。")
    if not re.fullmatch(r"[\x20-\x7E]+", old_text):
        raise ValueError("安全遮蔽目前只支援簡單英文、數字和半形符號。")
    with tempfile.TemporaryDirectory() as temp_dir:
        stripped_pdf = Path(temp_dir) / "redaction-content-removed.pdf"
        replace_text_in_content_stream(
            source,
            stripped_pdf,
            page_index,
            old_text,
            " " * len(old_text),
            password,
        )
        redact_text_block_overlay(stripped_pdf, target, page_index, block)


def font_key_for_pdf_font(font_name: str) -> str:
    normalized = (font_name or "").lower()
    if "times" in normalized or "roman" in normalized:
        return "times"
    if "courier" in normalized or "mono" in normalized:
        return "courier"
    return "helvetica"


def _tesseract_exe_candidates() -> list[Path]:
    names = ["tesseract.exe"] if sys.platform == "win32" else ["tesseract"]
    folders = []
    if sys.platform == "win32":
        folders.extend(
            [
                Path(os.environ.get("ProgramFiles", r"C:\Program Files")) / "Tesseract-OCR",
                Path(os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)")) / "Tesseract-OCR",
                Path.home() / "AppData" / "Local" / "Programs" / "Tesseract-OCR",
                Path.home() / "scoop" / "apps" / "tesseract" / "current",
            ]
        )
    found: list[Path] = []
    for folder in folders:
        for name in names:
            candidate = folder / name
            if candidate not in found:
                found.append(candidate)
    return found


def configure_tesseract() -> Path | None:
    if not OCR_AVAILABLE or pytesseract is None:
        return None
    try:
        pytesseract.get_tesseract_version()
        cmd = getattr(getattr(pytesseract, "pytesseract", None), "tesseract_cmd", None)
        return Path(cmd) if cmd else Path("tesseract")
    except TesseractNotFoundError:
        pass
    except Exception:
        pass
    for candidate in _tesseract_exe_candidates():
        if not candidate.is_file():
            continue
        pytesseract.pytesseract.tesseract_cmd = str(candidate)
        tessdata = candidate.parent / "tessdata"
        if tessdata.is_dir():
            os.environ["TESSDATA_PREFIX"] = str(tessdata)
        try:
            pytesseract.get_tesseract_version()
            return candidate
        except Exception:
            continue
    return None


def ensure_ocr_available() -> None:
    if not OCR_AVAILABLE or pytesseract is None:
        raise ValueError(
            "這份 PDF 沒有可抽取的文字層（掃描件），需要本機 OCR。\n"
            "請安裝 pytesseract，以及 Tesseract OCR（含 chi_sim / chi_tra 語言包）。\n"
            "Windows 安裝包：https://github.com/UB-Mannheim/tesseract/wiki"
        )
    if configure_tesseract() is None:
        raise ValueError(
            "找不到 Tesseract OCR 執行檔。掃描 PDF 必須用 OCR 才能轉 Word / Excel。\n"
            "本機若已安裝，常見路徑是 C:\\Program Files\\Tesseract-OCR\\tesseract.exe。\n"
            "尚未安裝請用 UB Mannheim 安裝包，並勾選 chi_sim / chi_tra。"
        )


def selected_page_indexes(source: Path, password: str = "", pages_spec: str = "") -> list[int]:
    reader = open_reader(source, password)
    if (pages_spec or "").strip():
        return parse_pages(pages_spec, len(reader.pages))
    return list(range(len(reader.pages)))


def iter_pdf_page_images(
    source: Path,
    password: str = "",
    pages_spec: str = "",
    dpi: int = 200,
):
    if not PDF_RENDER_AVAILABLE or pdfium is None:
        raise ValueError("PDF OCR 需要 pypdfium2 預覽元件。")
    if dpi < 72 or dpi > 600:
        raise ValueError("解析度請介於 72 至 600 DPI。")

    page_indexes = selected_page_indexes(source, password, pages_spec)
    if not page_indexes:
        raise ValueError("沒有可處理的頁面。")

    scale = dpi / 72.0
    document = pdfium.PdfDocument(str(source), password=password or None)
    try:
        for page_index in page_indexes:
            page = document.get_page(page_index)
            try:
                image = page.render(scale=scale).to_pil().convert("RGB")
            finally:
                page.close()
            yield page_index, image
    finally:
        document.close()


def render_pdf_page_images(
    source: Path,
    password: str = "",
    pages_spec: str = "",
    dpi: int = 200,
) -> list[tuple[int, Image.Image]]:
    return list(iter_pdf_page_images(source, password, pages_spec, dpi))


def _ocr_image_with_lang(image: Image.Image, tess_lang: str, tess_config: str) -> str:
    return _ocr_rows_to_plain_text(_ocr_image_to_rows(image, tess_lang, tess_config))


def _ocr_iter_pages(
    source: Path,
    password: str,
    pages_spec: str,
    language: str,
    dpi: int,
    progress=None,
):
    ensure_ocr_available()
    render_dpi = _ocr_render_dpi(language, dpi)
    iterator = iter(iter_pdf_page_images(source, password, pages_spec, render_dpi))
    first = next(iterator, None)
    if first is None:
        raise ValueError("沒有可處理的頁面。")
    _first_index, first_image = first
    requested = (language or "auto").strip() or "auto"
    if progress:
        progress(0, 0, "掃描件：正在偵測 OCR 語言…" if requested == "auto" else "掃描件：開始 OCR…")
    _resolved, tess_lang, tess_config = _resolved_ocr_settings(language, first_image)

    def pages():
        yield first
        yield from iterator

    return pages(), tess_lang, tess_config


def ocr_pdf_to_text(
    source: Path,
    target: Path,
    password: str = "",
    language: str = "auto",
    pages_spec: str = "",
    dpi: int = 200,
    progress=None,
) -> int:
    pages, tess_lang, tess_config = _ocr_iter_pages(source, password, pages_spec, language, dpi, progress)
    pieces: list[str] = []
    count = 0
    for page_index, image in pages:
        count += 1
        if progress:
            progress(count - 1, 0, f"掃描件 OCR 第 {page_index + 1} 頁，請稍候…")
        try:
            text = _ocr_image_with_lang(image, tess_lang, tess_config)
        finally:
            image.close()
        pieces.append(f"--- Page {page_index + 1} ---\n{normalize_cjk_text(text).strip()}\n")
    if progress:
        progress(count, count, f"OCR 完成 {count} 頁")
    target.write_text("\n".join(pieces), encoding="utf-8")
    return count


def ocr_pdf_to_searchable_pdf(
    source: Path,
    target: Path,
    password: str = "",
    language: str = "auto",
    pages_spec: str = "",
    dpi: int = 200,
    progress=None,
) -> int:
    pages, tess_lang, tess_config = _ocr_iter_pages(source, password, pages_spec, language, dpi, progress)
    writer = PdfWriter()
    count = 0
    for _page_index, image in pages:
        count += 1
        if progress:
            progress(count - 1, 0, f"掃描件 OCR 第 {_page_index + 1} 頁，請稍候…")
        prepared = _prepare_ocr_image(image)
        try:
            pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                prepared,
                extension="pdf",
                lang=tess_lang,
                config=tess_config,
            )
        finally:
            if prepared is not image:
                prepared.close()
            image.close()
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)
    if progress:
        progress(count, count, f"OCR 完成 {count} 頁")
    write_pdf(writer, target)
    return count


IMAGE_EXPORT_FORMATS = {"png", "jpg", "jpeg", "webp"}


def _image_extension(image_format: str) -> str:
    normalized = (image_format or "png").lower()
    if normalized in {"jpg", "jpeg"}:
        return "jpg"
    if normalized == "webp":
        return "webp"
    return "png"


def _save_rendered_image(image: Image.Image, path: Path, image_format: str) -> None:
    normalized = (image_format or "png").lower()
    if normalized in {"jpg", "jpeg"}:
        image.save(path, format="JPEG", quality=92)
    elif normalized == "webp":
        image.save(path, format="WEBP", quality=90)
    else:
        image.save(path, format="PNG")


def pdf_to_images(
    source: Path,
    target: Path,
    password: str = "",
    image_format: str = "png",
    dpi: int = 200,
    pages_spec: str = "",
) -> int:
    if not PDF_RENDER_AVAILABLE or pdfium is None:
        raise ValueError("PDF 轉圖片需要 pypdfium2 預覽元件。")
    if dpi < 72 or dpi > 600:
        raise ValueError("解析度請介於 72 至 600 DPI。")

    page_indexes = selected_page_indexes(source, password, pages_spec)
    if not page_indexes:
        raise ValueError("沒有可輸出的頁面。")

    scale = dpi / 72.0
    extension = _image_extension(image_format)
    stem = safe_output_name(source.stem)

    def write_images(folder: Path) -> list[Path]:
        folder.mkdir(parents=True, exist_ok=True)
        generated: list[Path] = []
        document = pdfium.PdfDocument(str(source), password=password or None)
        try:
            for output_index, page_index in enumerate(page_indexes, start=1):
                page = document.get_page(page_index)
                try:
                    image = page.render(scale=scale).to_pil().convert("RGB")
                finally:
                    page.close()
                filename = safe_output_name(f"{stem}-page-{page_index + 1:04d}.{extension}")
                output_path = folder / filename
                _save_rendered_image(image, output_path, image_format)
                generated.append(output_path)
        finally:
            document.close()
        return generated

    if target.suffix.lower() == ".zip":
        with tempfile.TemporaryDirectory() as temp_dir:
            generated = write_images(Path(temp_dir))
            with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as archive:
                for item in generated:
                    archive.write(item, item.name)
        return len(page_indexes)

    generated = write_images(target)
    return len(generated)


def images_to_pdf(
    paths: list[Path],
    target: Path,
    resolution: float = 150.0,
    lossless: bool = False,
) -> None:
    if not paths:
        raise ValueError("請選擇圖片。")
    dpi = float(resolution) or 150.0
    if lossless and PYMUPDF_AVAILABLE:
        document = fitz.open()
        try:
            for image_path in paths:
                with Image.open(image_path) as image:
                    width_px, height_px = image.size
                page = document.new_page(
                    width=width_px * 72.0 / dpi,
                    height=height_px * 72.0 / dpi,
                )
                page.insert_image(page.rect, filename=str(image_path))
            document.save(str(target), deflate=True)
        finally:
            document.close()
        return
    converted: list[Image.Image] = []
    for image_path in paths:
        image = Image.open(image_path)
        if image.mode in {"RGBA", "P"}:
            image = image.convert("RGB")
        converted.append(image)
    first, rest = converted[0], converted[1:]
    first.save(
        target,
        save_all=True,
        append_images=rest,
        resolution=dpi,
        quality=95,
    )
    for image in converted:
        image.close()


def office_app_for_path(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in OFFICE_WORD_SUFFIXES:
        return "word"
    if suffix in OFFICE_EXCEL_SUFFIXES:
        return "excel"
    if suffix in OFFICE_POWERPOINT_SUFFIXES:
        return "powerpoint"
    raise ValueError(f"不支援的 Office 格式：{path.suffix or path.name}")


def find_libreoffice_executable() -> Path | None:
    for name in ("soffice", "soffice.exe", "libreoffice"):
        found = shutil.which(name)
        if found:
            return Path(found)
    home = Path.home()
    for candidate in (
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files\LibreOffice 24\program\soffice.exe"),
        Path(r"C:\Program Files\LibreOffice 25\program\soffice.exe"),
        home / r"AppData\Local\Programs\LibreOffice\program\soffice.exe",
        Path("/usr/bin/soffice"),
        Path("/usr/bin/libreoffice"),
    ):
        if candidate.exists():
            return candidate
    return None


def _dispatch_com_app(win32com_client, prog_ids: tuple[str, ...]):
    last_error: Exception | None = None
    for prog_id in prog_ids:
        for factory in (win32com_client.DispatchEx, win32com_client.Dispatch):
            try:
                return factory(prog_id)
            except Exception as exc:
                last_error = exc
    raise last_error or RuntimeError("無法啟動 Office 應用程式。")


def _libreoffice_pdf_filter(source: Path) -> str:
    """PDF export filter that keeps original images and embeds fonts."""

    exporter = {
        "word": "writer_pdf_Export",
        "excel": "calc_pdf_Export",
        "powerpoint": "impress_pdf_Export",
    }[office_app_for_path(source)]
    options = (
        '{"UseLosslessCompression":{"type":"boolean","value":"true"},'
        '"ReduceImageResolution":{"type":"boolean","value":"false"},'
        '"MaxImageResolution":{"type":"long","value":"600"},'
        '"Quality":{"type":"long","value":"100"},'
        '"SelectPdfVersion":{"type":"long","value":"0"},'
        '"EmbedStandardFonts":{"type":"boolean","value":"true"}}'
    )
    return f"pdf:{exporter}:{options}"


def _try_com_call(action, *fallbacks):
    errors: list[Exception] = []
    for attempt in (action, *fallbacks):
        try:
            return attempt()
        except Exception as exc:
            errors.append(exc)
    raise errors[-1]


def _hide_com_window(app) -> None:
    """Keep Office automation in the background; PowerPoint often forbids Visible=False."""

    for value in (False, 0):
        try:
            app.Visible = value
            break
        except Exception:
            continue
    else:
        try:
            app.Visible = True
        except Exception:
            pass
        try:
            app.WindowState = 2
        except Exception:
            pass
    for value in (0, 1):
        try:
            app.DisplayAlerts = value
            break
        except Exception:
            continue
    hwnd = 0
    try:
        hwnd = int(app.HWND)
    except Exception:
        hwnd = 0
    if hwnd:
        try:
            ctypes.windll.user32.ShowWindow(hwnd, 0)
        except Exception:
            pass


def _prepare_word_font_embedding(document) -> None:
    for name, value in (
        ("EmbedTrueTypeFonts", True),
        ("DoNotEmbedSystemFonts", False),
        ("SaveSubsetFonts", False),
    ):
        try:
            setattr(document, name, value)
        except Exception:
            continue


def _export_word_pdf(word, source_s: str, target_s: str, progress=None) -> None:
    _report_office_progress(progress, 0, 1, f"正在轉換 {Path(source_s).name}")
    document = None
    try:
        document = _try_com_call(
            lambda: word.Documents.Open(source_s, False, True, False),
            lambda: word.Documents.Open(source_s),
        )
        _prepare_word_font_embedding(document)
        _try_com_call(
            lambda: document.ExportAsFixedFormat(
                target_s,
                17,
                False,
                0,
                0,
                1,
                1,
                0,
                True,
                True,
                1,
                True,
                False,
                False,
            ),
            lambda: document.ExportAsFixedFormat(target_s, 17, False, 0),
            lambda: document.SaveAs(target_s, 17),
        )
        _report_office_progress(progress, 1, 1, "完成")
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass


def _export_excel_pdf(excel, source_s: str, target_s: str, progress=None) -> None:
    _report_office_progress(progress, 0, 1, f"正在轉換 {Path(source_s).name}")
    workbook = None
    try:
        workbook = _try_com_call(
            lambda: excel.Workbooks.Open(source_s, UpdateLinks=0, ReadOnly=True),
            lambda: excel.Workbooks.Open(source_s),
        )
        _try_com_call(
            lambda: workbook.ExportAsFixedFormat(0, target_s, 0, True, False, None, None, False),
            lambda: workbook.ExportAsFixedFormat(0, target_s, 0),
            lambda: workbook.ExportAsFixedFormat(0, target_s),
        )
        _report_office_progress(progress, 1, 1, "完成")
    finally:
        if workbook is not None:
            try:
                workbook.Close(False)
            except Exception:
                pass


def _export_powerpoint_pdf(powerpoint, source_s: str, target_s: str, progress=None) -> None:
    _hide_com_window(powerpoint)
    _report_office_progress(progress, 0, 1, f"正在開啟 {Path(source_s).name}")
    presentation = None
    try:
        presentation = _try_com_call(
            lambda: powerpoint.Presentations.Open(source_s, True, False, False),
            lambda: powerpoint.Presentations.Open(source_s, False, False, False),
            lambda: powerpoint.Presentations.Open(source_s, WithWindow=False),
            lambda: powerpoint.Presentations.Open(source_s),
        )
        _hide_com_window(powerpoint)
        try:
            _export_powerpoint_slides_as_pdf(presentation, Path(target_s), progress=progress)
        except Exception:
            _try_com_call(
                lambda: presentation.SaveAs(target_s, 32, True),
                lambda: presentation.SaveAs(target_s, 32),
                lambda: presentation.ExportAsFixedFormat(target_s, 2, 2),
            )
    finally:
        if presentation is not None:
            try:
                presentation.Saved = True
            except Exception:
                pass
            try:
                presentation.Close()
            except Exception:
                pass


def _export_powerpoint_slides_as_pdf(presentation, target: Path, dpi: int = 300, progress=None) -> None:
    """Rasterize each slide at print DPI so bevel, shadow and vectors match the deck."""

    count = int(presentation.Slides.Count)
    if count <= 0:
        raise RuntimeError("簡報沒有投影片。")
    width_pt = float(presentation.PageSetup.SlideWidth)
    height_pt = float(presentation.PageSetup.SlideHeight)
    px_w = max(int(round(width_pt / 72.0 * dpi)), 1)
    px_h = max(int(round(height_pt / 72.0 * dpi)), 1)
    total = count + 1
    with tempfile.TemporaryDirectory() as temp_dir:
        slide_dir = Path(temp_dir)
        pngs: list[Path] = []
        for index in range(1, count + 1):
            _report_office_progress(progress, index - 1, total, f"正在匯出第 {index} / {count} 頁")
            png = slide_dir / f"slide-{index:04d}.png"
            presentation.Slides(index).Export(str(png), "PNG", px_w, px_h)
            if not png.is_file() or png.stat().st_size == 0:
                raise RuntimeError(f"第 {index} 頁投影片匯出失敗。")
            pngs.append(png)
        _report_office_progress(progress, count, total, "正在寫入 PDF")
        images_to_pdf(pngs, target, resolution=float(dpi), lossless=True)
        _report_office_progress(progress, total, total, "完成")


def _export_office_via_com(source: Path, target: Path, progress=None) -> None:
    import win32com.client

    app_kind = office_app_for_path(source)
    source_s = str(source.resolve())
    target_s = str(target.resolve())
    if app_kind == "word":
        word = _dispatch_com_app(win32com.client, ("Word.Application", "Word.Application.16"))
        _hide_com_window(word)
        try:
            _export_word_pdf(word, source_s, target_s, progress=progress)
        finally:
            try:
                word.Quit()
            except Exception:
                pass
        return
    if app_kind == "excel":
        excel = _dispatch_com_app(win32com.client, ("Excel.Application", "Excel.Application.16"))
        _hide_com_window(excel)
        try:
            _export_excel_pdf(excel, source_s, target_s, progress=progress)
        finally:
            try:
                excel.Quit()
            except Exception:
                pass
        return
    powerpoint = _dispatch_com_app(win32com.client, ("PowerPoint.Application", "PowerPoint.Application.16"))
    try:
        _export_powerpoint_pdf(powerpoint, source_s, target_s, progress=progress)
    finally:
        try:
            powerpoint.Quit()
        except Exception:
            pass


def _convert_office_with_com(
    source: Path,
    target: Path,
    progress=None,
    pump: Callable[[], None] | None = None,
) -> tuple[bool, str]:
    try:
        import pythoncom  # noqa: F401
        import win32com.client  # noqa: F401
    except ImportError:
        return False, "未安裝 Windows COM 元件。"

    result: dict[str, object] = {"ok": False, "error": ""}

    def worker() -> None:
        try:
            import pythoncom

            pythoncom.CoInitializeEx(pythoncom.COINIT_APARTMENTTHREADED)
            try:
                _export_office_via_com(source, target, progress=progress)
                if target.is_file() and target.stat().st_size > 0:
                    result["ok"] = True
                else:
                    result["error"] = "Microsoft Office 沒有產生 PDF。"
            finally:
                pythoncom.CoUninitialize()
        except Exception as exc:
            result["error"] = str(exc) or type(exc).__name__

    thread = threading.Thread(target=worker, name="office-com-sta", daemon=True)
    thread.start()
    timeout = 360 if office_app_for_path(source) == "powerpoint" else 180
    if pump is None:
        thread.join(timeout=timeout)
    else:
        waited = 0.0
        slice_s = 0.05
        while thread.is_alive():
            pump()
            thread.join(timeout=slice_s)
            waited += slice_s
            if waited >= timeout:
                break
    if thread.is_alive():
        return False, "Office 轉換逾時。請確認 Word / Excel / PowerPoint 沒有跳出對話框。"
    if result["ok"]:
        return True, ""
    return False, str(result["error"] or "Microsoft Office 無法轉換此檔。")


def _convert_office_with_libreoffice(source: Path, target: Path, progress=None) -> None:
    executable = find_libreoffice_executable()
    if executable is None:
        raise ValueError("未找到 LibreOffice。")
    _report_office_progress(progress, 0, 1, f"正在用 LibreOffice 轉換 {source.name}")
    with tempfile.TemporaryDirectory() as temp_dir:
        out_dir = Path(temp_dir)
        command = [
            str(executable),
            "--headless",
            "--norestore",
            "--nolockcheck",
            "--convert-to",
            _libreoffice_pdf_filter(source),
            "--outdir",
            str(out_dir),
            str(source.resolve()),
        ]
        completed = subprocess.run(command, capture_output=True, timeout=180, check=False)
        produced = out_dir / f"{source.stem}.pdf"
        if not produced.exists():
            fallback = [
                str(executable),
                "--headless",
                "--norestore",
                "--nolockcheck",
                "--convert-to",
                "pdf",
                "--outdir",
                str(out_dir),
                str(source.resolve()),
            ]
            completed = subprocess.run(fallback, capture_output=True, timeout=180, check=False)
            produced = out_dir / f"{source.stem}.pdf"
        if not produced.exists():
            pdfs = list(out_dir.glob("*.pdf"))
            if not pdfs:
                detail = (completed.stderr or completed.stdout or b"").decode("utf-8", errors="replace").strip()
                raise ValueError(detail or "LibreOffice 沒有產生 PDF。")
            produced = pdfs[0]
        shutil.copy2(produced, target)
        _report_office_progress(progress, 1, 1, "完成")


def convert_office_file_to_pdf(
    source: Path,
    target: Path,
    progress=None,
    pump: Callable[[], None] | None = None,
) -> None:
    if not source.is_file():
        raise ValueError(f"找不到檔案：{source.name}")
    office_app_for_path(source)
    target = Path(target)
    target.parent.mkdir(parents=True, exist_ok=True)
    _report_office_progress(progress, 0, 1, f"正在準備 {source.name}")
    with tempfile.TemporaryDirectory() as temp_dir:
        local_source = Path(temp_dir) / source.name
        shutil.copy2(source, local_source)
        com_ok, com_error = _convert_office_with_com(
            local_source, target, progress=progress, pump=pump
        )
        if com_ok:
            return
        try:
            _convert_office_with_libreoffice(local_source, target, progress=progress)
        except Exception as libre_error:
            details = []
            if com_error:
                details.append(f"Microsoft Office：{com_error}")
            details.append(f"LibreOffice：{libre_error}")
            raise ValueError(
                "無法轉換 Office 檔。請確認 Word / Excel / PowerPoint 可開啟此檔，"
                "或安裝 LibreOffice 後再試。\n" + "\n".join(details)
            ) from libre_error
    if not target.is_file() or target.stat().st_size == 0:
        raise ValueError(f"轉換失敗，沒有產生 PDF：{source.name}")


def office_files_to_pdf(
    paths: list[Path],
    target: Path,
    progress=None,
    pump: Callable[[], None] | None = None,
) -> int:
    office_paths = [path for path in paths if path.suffix.lower() in OFFICE_SUFFIXES]
    if not office_paths:
        raise ValueError("請加入 Word、Excel 或 PowerPoint 檔案。")
    target = Path(target)
    target.parent.mkdir(parents=True, exist_ok=True)
    if len(office_paths) == 1:
        convert_office_file_to_pdf(office_paths[0], target, progress=progress, pump=pump)
        return 1
    with tempfile.TemporaryDirectory() as temp_dir:
        converted: list[Path] = []
        file_count = len(office_paths)
        for index, source in enumerate(office_paths, start=1):
            prefix = f"({index}/{file_count}) {source.name}："

            def nested(current: int, total: int, text: str, prefix=prefix) -> None:
                _report_office_progress(progress, current, total, prefix + text)

            pdf_path = Path(temp_dir) / f"{index:03d}-{suggested_pdf_name_for_source(source)}"
            convert_office_file_to_pdf(source, pdf_path, progress=nested, pump=pump)
            converted.append(pdf_path)
        _report_office_progress(progress, 0, 1, "正在合併 PDF")
        merge_pdf_files(converted, target)
        _report_office_progress(progress, 1, 1, "完成")
    return len(office_paths)


ANNOTATION_FONT_OPTIONS = {
    "cjk": ("Helv", "Helv-Bold"),
    "helvetica": ("Helv", "Helv-Bold"),
    "times": ("Times-Roman", "Times-Bold"),
    "courier": ("Courier", "Courier-Bold"),
}

ANNOTATION_COLOR_PRESETS: dict[str, tuple[float, float, float]] = {
    "black": (0.0, 0.0, 0.0),
    "red": (0.8, 0.0, 0.0),
    "blue": (0.0, 0.0, 0.7),
    "gray": (0.4, 0.4, 0.4),
}

ANNOTATION_FILL_PRESETS: dict[str, tuple[float, float, float] | None] = {
    "white": (1.0, 1.0, 1.0),
    "none": None,
    "yellow": (1.0, 0.96, 0.72),
    "blue": (0.86, 0.93, 1.0),
    "green": (0.86, 0.96, 0.86),
}


def resolve_annotation_fill(
    color_rgb: tuple[float, float, float],
    fill_rgb: tuple[float, float, float] | None = None,
    fill_none: bool = False,
) -> tuple[float, float, float] | None:
    if fill_none:
        return None
    if fill_rgb is not None:
        return fill_rgb
    return tuple(channel * 0.22 + 0.78 for channel in color_rgb)


def build_annotation_da(
    font_key: str,
    font_size: int,
    bold: bool,
    color_rgb: tuple[float, float, float],
) -> str:
    regular, bold_name = ANNOTATION_FONT_OPTIONS.get(font_key, ANNOTATION_FONT_OPTIONS["helvetica"])
    font_name = bold_name if bold else regular
    red, green, blue = color_rgb
    return f"/{font_name} {font_size} Tf {red} {green} {blue} rg"


def text_contains_cjk(text: str) -> bool:
    return any(
        "\u3400" <= char <= "\u9fff" or "\uf900" <= char <= "\ufaff" or "\u3000" <= char <= "\u303f"
        for char in text or ""
    )


def overlay_needs_embedded_font(text: str, font_key: str = "") -> bool:
    key = (font_key or "").strip().lower()
    return key in {"cjk", "yahei", "mingliu"} or text_contains_cjk(text)


def resolve_cjk_font_file(bold: bool = False) -> str | None:
    fonts_dir = Path("C:/Windows/Fonts")
    names: list[str] = []
    if bold:
        names.extend(["msjhbd.ttc", "msyhbd.ttc"])
    names.extend(["msjh.ttc", "msyh.ttc", "mingliu.ttc", "simsun.ttc"])
    for name in names:
        path = fonts_dir / name
        if path.is_file():
            return str(path)
    return None


def _embed_overlay_text(
    target: Path,
    page_index: int,
    x: float,
    y: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    bold: bool,
    color_rgb: tuple[float, float, float],
    password: str = "",
) -> None:
    ensure_pymupdf_available()
    document = None
    try:
        document = fitz.open(str(target))
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{target.name} 已加密，請輸入密碼。")
        page = document[page_index]
        page_height = float(page.rect.height)
        rect = fitz.Rect(x + 2, page_height - (y + height) + 1, x + width - 2, page_height - y - 1)
        kwargs: dict = {
            "color": color_rgb,
            "align": fitz.TEXT_ALIGN_LEFT,
        }
        fontfile = resolve_cjk_font_file(bold)
        if fontfile:
            kwargs["fontfile"] = fontfile
            kwargs["fontname"] = "overlay-cjk"
        else:
            kwargs["fontname"] = "china-t"
        fontsize = float(max(font_size, 8))
        unused = page.insert_textbox(rect, text, fontsize=fontsize, **kwargs)
        if unused < 0:
            for size in range(int(fontsize) - 1, 7, -1):
                unused = page.insert_textbox(rect, text, fontsize=float(size), **kwargs)
                if unused >= 0:
                    fontsize = float(size)
                    break
        if unused < 0:
            point = fitz.Point(rect.x0, min(rect.y1 - 1, rect.y0 + fontsize))
            insert_kwargs = {key: value for key, value in kwargs.items() if key != "align"}
            page.insert_text(point, text, fontsize=max(fontsize * 0.75, 8), **insert_kwargs)
        if hasattr(document, "subset_fonts"):
            document.subset_fonts()
        temp_target = target.with_name(f"{target.stem}.overlay-tmp{target.suffix}")
        document.save(str(temp_target), garbage=4, deflate=True)
        document.close()
        document = None
        temp_target.replace(target)
    finally:
        if document is not None:
            document.close()


def rgb_to_hex(color_rgb: tuple[float, float, float]) -> str:
    red, green, blue = color_rgb
    return "#{:02x}{:02x}{:02x}".format(
        int(max(0.0, min(red, 1.0)) * 255),
        int(max(0.0, min(green, 1.0)) * 255),
        int(max(0.0, min(blue, 1.0)) * 255),
    )


def add_annotation_to_page(writer: PdfWriter, page, annotation: DictionaryObject) -> None:
    annotation_ref = writer._add_object(annotation)
    if "/Annots" not in page:
        page[NameObject("/Annots")] = ArrayObject()
    page["/Annots"].append(annotation_ref)


def add_text_overlay_annotation(
    source: Path,
    target: Path,
    page_index: int,
    x: float,
    y: float,
    text: str,
    font_size: int,
    cover_original: bool,
    cover_width: float,
    cover_height: float,
    password: str = "",
    font_key: str = "helvetica",
    bold: bool = False,
    color_rgb: tuple[float, float, float] = (0.0, 0.0, 0.0),
    shape: str = "box",
    pointer: tuple[float, float] | None = None,
    fill_rgb: tuple[float, float, float] | None = None,
    fill_none: bool = False,
) -> None:
    reader = open_reader(source, password)
    if page_index < 0 or page_index >= len(reader.pages):
        raise ValueError("頁碼超出範圍。")
    if not text.strip():
        raise ValueError("請輸入要加入的文字。")

    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    page = writer.pages[page_index]
    text_height = max(float(font_size) * 1.8, 24.0)
    rect_width = max(float(cover_width), 80.0)
    rect_height = max(float(cover_height), text_height)

    shape = (shape or "box").strip().lower()
    embed_text = overlay_needs_embedded_font(text, font_key)
    if shape in CALLOUT_KINDS:
        pointer_x, pointer_y = pointer if pointer is not None else (x + rect_width / 2.0, y - 16.0)
        markup = MarkupAnnotation(
            kind=shape,
            x0=pointer_x,
            y0=pointer_y,
            x1=x,
            y1=y,
            color_rgb=color_rgb,
            contents="" if embed_text else text,
            box_width=rect_width,
            box_height=rect_height,
            fill_rgb=fill_rgb,
            fill_none=fill_none,
        )
        da = build_annotation_da(font_key, font_size, bold, color_rgb)
        for annotation in iter_markup_pdf_annotations(markup):
            if annotation.get("/Subtype") == "/FreeText":
                if embed_text:
                    annotation[NameObject("/Contents")] = TextStringObject("")
                annotation[NameObject("/DA")] = TextStringObject(da)
            add_annotation_to_page(writer, page, annotation)
        write_pdf(writer, target)
        if embed_text:
            _embed_overlay_text(
                target,
                page_index,
                x,
                y,
                rect_width,
                rect_height,
                text,
                font_size,
                bold,
                color_rgb,
                password,
            )
        return

    if shape == "rect":
        fill = resolve_annotation_fill(color_rgb, fill_rgb, fill_none)
        frame = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/Square"),
                NameObject("/Rect"): ArrayObject(
                    [
                        FloatObject(x),
                        FloatObject(y),
                        FloatObject(x + rect_width),
                        FloatObject(y + rect_height),
                    ]
                ),
                NameObject("/C"): _color_array(color_rgb),
                NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), FloatObject(1.5)]),
                NameObject("/F"): NumberObject(4),
            }
        )
        if fill is not None:
            frame[NameObject("/IC")] = _color_array(fill)
        add_annotation_to_page(writer, page, frame)
        cover_original = False

    if cover_original:
        cover = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/Square"),
                NameObject("/Rect"): ArrayObject(
                    [
                        FloatObject(x),
                        FloatObject(y - 4),
                        FloatObject(x + rect_width),
                        FloatObject(y + rect_height),
                    ]
                ),
                NameObject("/C"): ArrayObject([FloatObject(1), FloatObject(1), FloatObject(1)]),
                NameObject("/IC"): ArrayObject([FloatObject(1), FloatObject(1), FloatObject(1)]),
                NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(0)]),
                NameObject("/F"): NumberObject(4),
            }
        )
        add_annotation_to_page(writer, page, cover)

    if not embed_text:
        free_text = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/FreeText"),
                NameObject("/Rect"): ArrayObject(
                    [
                        FloatObject(x),
                        FloatObject(y),
                        FloatObject(x + rect_width),
                        FloatObject(y + rect_height),
                    ]
                ),
                NameObject("/Contents"): TextStringObject(text),
                NameObject("/DA"): TextStringObject(build_annotation_da(font_key, font_size, bold, color_rgb)),
                NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(0)]),
                NameObject("/F"): NumberObject(4),
            }
        )
        add_annotation_to_page(writer, page, free_text)
    write_pdf(writer, target)
    if embed_text:
        _embed_overlay_text(
            target,
            page_index,
            x,
            y,
            rect_width,
            rect_height,
            text,
            font_size,
            bold,
            color_rgb,
            password,
        )


def add_text_overlay_annotations(
    source: Path,
    target: Path,
    overlays: list[dict],
    password: str = "",
) -> None:
    if not overlays:
        raise ValueError("請先加入至少一筆文字標註。")
    source = Path(source)
    target = Path(target)
    current = source
    temps: list[Path] = []
    try:
        for index, overlay in enumerate(overlays):
            is_last = index == len(overlays) - 1
            if is_last:
                out = target
            else:
                handle = tempfile.NamedTemporaryFile(suffix=".pdf", prefix=f"ann{index}_", delete=False)
                handle.close()
                out = Path(handle.name)
                temps.append(out)
            shape = str(overlay.get("shape") or "box")
            pointer = None
            if shape in CALLOUT_KINDS:
                pointer = (float(overlay["pointer_x"]), float(overlay["pointer_y"]))
            add_text_overlay_annotation(
                source=current,
                target=out,
                page_index=int(overlay["page_index"]),
                x=float(overlay["pdf_x"]),
                y=float(overlay["pdf_y"]),
                text=str(overlay["text"]),
                font_size=int(overlay["font_size"]),
                cover_original=bool(overlay.get("cover")),
                cover_width=float(overlay["rect_width"]),
                cover_height=float(overlay["rect_height"]),
                password=password if current == source else "",
                font_key=str(overlay.get("font_key") or "helvetica"),
                bold=bool(overlay.get("bold")),
                color_rgb=tuple(overlay.get("color_rgb") or (0.0, 0.0, 0.0)),
                shape=shape,
                pointer=pointer,
                fill_rgb=tuple(overlay["fill_rgb"]) if overlay.get("fill_rgb") is not None else None,
                fill_none=bool(overlay.get("fill_none")),
            )
            current = out
    finally:
        for path in temps:
            try:
                path.unlink(missing_ok=True)
            except OSError:
                pass


def add_page_numbers(source: Path, target: Path, template: str = "Page {page} of {total}", password: str = "") -> None:
    reader = open_reader(source, password)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    total = len(writer.pages)
    for index, page in enumerate(writer.pages, start=1):
        width = float(page.mediabox.width)
        text = template.format(page=index, total=total)
        rect_width = max(160.0, len(text) * 7.0)
        x = max((width - rect_width) / 2, 24.0)
        annotation = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/FreeText"),
                NameObject("/Rect"): ArrayObject(
                    [FloatObject(x), FloatObject(20), FloatObject(x + rect_width), FloatObject(42)]
                ),
                NameObject("/Contents"): TextStringObject(text),
                NameObject("/DA"): TextStringObject("/Helv 10 Tf 0.25 0.25 0.25 rg"),
                NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(0)]),
                NameObject("/F"): NumberObject(4),
            }
        )
        add_annotation_to_page(writer, page, annotation)
    write_pdf(writer, target)


WATERMARK_POSITIONS = {
    "center": "正中",
    "top-left": "左上",
    "top-center": "上中",
    "top-right": "右上",
    "middle-left": "左中",
    "middle-right": "右中",
    "bottom-left": "左下",
    "bottom-center": "下中",
    "bottom-right": "右下",
    "custom": "自訂座標（左 % / 上 %）",
}

WATERMARK_ROTATIONS = (
    (45, "斜向 45°（Word 預設）"),
    (0, "水平"),
    (30, "斜向 30°"),
    (60, "斜向 60°"),
    (-45, "反向斜向 -45°"),
    (90, "直向 90°"),
)

WATERMARK_COLORS = {
    "gray": ((0.55, 0.55, 0.55), "灰色"),
    "light-gray": ((0.75, 0.75, 0.75), "淺灰"),
    "red": ((0.72, 0.16, 0.16), "紅色"),
    "blue": ((0.16, 0.28, 0.68), "藍色"),
    "black": ((0.18, 0.18, 0.18), "黑色"),
}

_WATERMARK_ANCHORS = {
    "center": ("center", "middle"),
    "top-left": ("left", "top"),
    "top-center": ("center", "top"),
    "top-right": ("right", "top"),
    "middle-left": ("left", "middle"),
    "middle-right": ("right", "middle"),
    "bottom-left": ("left", "bottom"),
    "bottom-center": ("center", "bottom"),
    "bottom-right": ("right", "bottom"),
}


def _watermark_insert_kwargs(text: str) -> tuple[dict, Callable[[float], float]]:
    kwargs: dict = {}
    if overlay_needs_embedded_font(text):
        fontfile = resolve_cjk_font_file(False)
        if fontfile:
            kwargs["fontfile"] = fontfile
            kwargs["fontname"] = "wm-cjk"
            font = fitz.Font(fontfile=fontfile)

            def width_at(size: float, _font=font) -> float:
                return float(_font.text_length(text, fontsize=size))

            return kwargs, width_at
        kwargs["fontname"] = "china-t"

        def width_at(size: float) -> float:
            return float(fitz.get_text_length(text, fontname="china-t", fontsize=size))

        return kwargs, width_at
    kwargs["fontname"] = "helv"

    def width_at(size: float) -> float:
        return float(fitz.get_text_length(text, fontname="helv", fontsize=size))

    return kwargs, width_at


def _watermark_pivot(
    page_rect,
    position: str,
    text_width: float,
    font_size: float,
    rotation: float,
    x_percent: float | None = None,
    y_percent: float | None = None,
):
    width = float(page_rect.width)
    height = float(page_rect.height)
    if position == "custom":
        xp = 50.0 if x_percent is None else float(x_percent)
        yp = 50.0 if y_percent is None else float(y_percent)
        xp = min(max(xp, 0.0), 100.0)
        yp = min(max(yp, 0.0), 100.0)
        return fitz.Point(width * xp / 100.0, height * yp / 100.0)

    if position not in _WATERMARK_ANCHORS:
        raise ValueError("不支援的浮水印位置。")

    rad = math.radians(rotation)
    box_w = abs(text_width * math.cos(rad)) + abs(font_size * math.sin(rad))
    box_h = abs(text_width * math.sin(rad)) + abs(font_size * math.cos(rad))
    half_w = box_w / 2.0
    half_h = box_h / 2.0
    margin = 48.0
    horiz, vert = _WATERMARK_ANCHORS[position]
    if horiz == "left":
        x = margin + half_w
    elif horiz == "right":
        x = width - margin - half_w
    else:
        x = width / 2.0
    if vert == "top":
        y = margin + half_h
    elif vert == "bottom":
        y = height - margin - half_h
    else:
        y = height / 2.0
    return fitz.Point(min(max(x, margin), width - margin), min(max(y, margin), height - margin))


def add_watermark(
    source: Path,
    target: Path,
    text: str,
    password: str = "",
    position: str = "center",
    rotation: int = 45,
    font_size: int = 48,
    opacity: float = 0.25,
    color_rgb: tuple[float, float, float] = (0.55, 0.55, 0.55),
    pages_spec: str = "",
    x_percent: float | None = None,
    y_percent: float | None = None,
) -> int:
    """Draw a Word-style text watermark on selected pages."""

    stamp = (text or "").strip()
    if not stamp:
        raise ValueError("請輸入水印 / 印章文字。")
    if position not in WATERMARK_POSITIONS:
        raise ValueError("不支援的浮水印位置。")
    ensure_pymupdf_available()
    fontsize = float(min(max(int(font_size), 6), 144))
    alpha = min(max(float(opacity), 0.05), 1.0)
    angle = float(rotation)
    red, green, blue = color_rgb
    color = (
        min(max(float(red), 0.0), 1.0),
        min(max(float(green), 0.0), 1.0),
        min(max(float(blue), 0.0), 1.0),
    )
    font_kwargs, width_at = _watermark_insert_kwargs(stamp)
    text_width = max(width_at(fontsize), 8.0)

    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        page_indexes = (
            parse_pages(pages_spec, document.page_count)
            if (pages_spec or "").strip()
            else list(range(document.page_count))
        )
        if not page_indexes:
            raise ValueError("沒有可加入浮水印的頁面。")
        rotate = abs(angle) % 360.0 > 0.01
        matrix = fitz.Matrix(1, 1).prerotate(angle) if rotate else None
        for page_index in page_indexes:
            page = document[page_index]
            page_size = fontsize
            page_width = text_width
            max_width = min(float(page.rect.width), float(page.rect.height)) * 0.85
            if page_width > max_width:
                page_size = max(6.0, fontsize * (max_width / page_width))
                page_width = max(width_at(page_size), 8.0)
            pivot = _watermark_pivot(
                page.rect,
                position,
                page_width,
                page_size,
                angle,
                x_percent=x_percent,
                y_percent=y_percent,
            )
            start = fitz.Point(pivot.x - page_width / 2.0, pivot.y + page_size * 0.35)
            insert_kwargs = {
                "fontsize": page_size,
                "color": color,
                "fill_opacity": alpha,
                "overlay": True,
                **font_kwargs,
            }
            if matrix is not None:
                insert_kwargs["morph"] = (pivot, matrix)
            page.insert_text(start, stamp, **insert_kwargs)
        if overlay_needs_embedded_font(stamp) and hasattr(document, "subset_fonts"):
            document.subset_fonts()
        document.save(str(target), garbage=4, deflate=True)
        return len(page_indexes)
    finally:
        document.close()


def rendered_page_is_blank(page, threshold: int) -> bool:
    if not PDF_RENDER_AVAILABLE or pdfium is None:
        raise ValueError("刪除空白頁需要 pypdfium2 預覽元件。")
    bitmap = page.render(scale=0.15)
    image = bitmap.to_pil().convert("L")
    non_white = 0
    for pixel in image.getdata():
        if pixel < 245:
            non_white += 1
            if non_white > threshold:
                return False
    return True


def remove_blank_pages(source: Path, target: Path, threshold: int = 25, password: str = "") -> int:
    reader = open_reader(source, password)
    document = pdfium.PdfDocument(str(source), password=password or None) if pdfium is not None else None
    writer = PdfWriter()
    removed = 0
    try:
        for index, source_page in enumerate(reader.pages):
            rendered_page = document.get_page(index)
            try:
                is_blank = rendered_page_is_blank(rendered_page, threshold)
            finally:
                rendered_page.close()
            if is_blank:
                removed += 1
            else:
                writer.add_page(source_page)
    finally:
        if document is not None:
            document.close()
    if not writer.pages:
        raise ValueError("所有頁面都被判定為空白，已取消輸出。")
    write_pdf(writer, target)
    return removed


def clean_metadata(source: Path, target: Path, password: str = "") -> None:
    reader = open_reader(source, password)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    writer.add_metadata({"/Producer": "Victor PDF Tools Box"})
    write_pdf(writer, target)


@dataclass(frozen=True)
class BookmarkItem:
    title: str
    page_index: int
    level: int = 0


def extract_outline(source: Path, password: str = "") -> list[BookmarkItem]:
    """Return a flattened list of bookmarks with indentation levels."""

    reader = open_reader(source, password)
    page_count = len(reader.pages)
    items: list[BookmarkItem] = []

    def walk(entries, level: int) -> None:
        for entry in entries:
            if isinstance(entry, list):
                walk(entry, level + 1)
                continue
            title = str(getattr(entry, "title", "") or "").strip()
            try:
                page_index = reader.get_destination_page_number(entry)
            except Exception:
                page_index = None
            if page_index is None or page_index < 0:
                page_index = 0
            if page_count and page_index >= page_count:
                page_index = page_count - 1
            items.append(BookmarkItem(title or "(未命名書籤)", int(page_index), max(0, level)))

    try:
        outline = reader.outline
    except Exception:
        outline = []
    walk(outline, 0)
    return items


def apply_outline(source: Path, target: Path, items: list[BookmarkItem], password: str = "") -> None:
    """Rebuild the document outline from a flat list of bookmarks."""

    reader = open_reader(source, password)
    page_count = len(reader.pages)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)

    parents: dict[int, object] = {}
    for item in items:
        title = (item.title or "").strip() or "(未命名書籤)"
        page_index = item.page_index
        if page_index < 0:
            page_index = 0
        if page_count and page_index >= page_count:
            page_index = page_count - 1
        level = max(0, item.level)
        parent = parents.get(level - 1) if level > 0 else None
        ref = writer.add_outline_item(title, page_index, parent=parent)
        parents[level] = ref
        for deeper in [key for key in parents if key > level]:
            del parents[deeper]

    write_pdf(writer, target)


# --- markup annotations (highlight / shapes / sticky notes) ------------------

MARKUP_KINDS = (
    "highlight",
    "underline",
    "strikeout",
    "rect",
    "ellipse",
    "line",
    "arrow",
    "note",
)

MARKUP_SUBTYPES = {
    "highlight": "/Highlight",
    "underline": "/Underline",
    "strikeout": "/StrikeOut",
    "rect": "/Square",
    "ellipse": "/Circle",
    "line": "/Line",
    "arrow": "/Line",
    "note": "/Text",
    "callout": "/FreeText",
    "speech": "/Polygon",
    "cloud": "/Polygon",
}

CALLOUT_KINDS = frozenset({"callout", "speech", "cloud", "comment"})

MARKUP_COLOR_PRESETS: dict[str, tuple[float, float, float]] = {
    "yellow": (1.0, 0.92, 0.23),
    "green": (0.45, 0.86, 0.45),
    "blue": (0.40, 0.74, 1.0),
    "pink": (1.0, 0.58, 0.79),
    "red": (0.88, 0.16, 0.16),
    "black": (0.0, 0.0, 0.0),
}


@dataclass(frozen=True)
class MarkupAnnotation:
    """A markup placed on a page, in PDF user space (origin bottom-left).

    For rectangle-like markups (highlight/underline/strikeout/rect/ellipse) the
    two points are opposite corners. For line/arrow they are the endpoints. For
    a sticky note only (x0, y0) is used as the anchor point.
    """

    kind: str
    x0: float
    y0: float
    x1: float
    y1: float
    color_rgb: tuple[float, float, float] = (1.0, 0.92, 0.23)
    contents: str = ""
    box_width: float = 0.0
    box_height: float = 0.0
    fill_rgb: tuple[float, float, float] | None = None
    fill_none: bool = False


@dataclass(frozen=True)
class CalloutLayout:
    pointer: tuple[float, float]
    left: float
    bottom: float
    right: float
    top: float
    attach: tuple[float, float]
    vertices: tuple[tuple[float, float], ...]
    text: str
    path: tuple[tuple[float, float], ...] = ()


def _normalized_rect(x0: float, y0: float, x1: float, y1: float) -> tuple[float, float, float, float]:
    return (min(x0, x1), min(y0, y1), max(x0, x1), max(y0, y1))


def _color_array(color_rgb: tuple[float, float, float]) -> ArrayObject:
    red, green, blue = color_rgb
    return ArrayObject([FloatObject(red), FloatObject(green), FloatObject(blue)])


def _closest_point_on_rect(
    px: float,
    py: float,
    left: float,
    bottom: float,
    right: float,
    top: float,
) -> tuple[float, float]:
    if left < px < right and bottom < py < top:
        candidates = (
            (px - left, (left, py)),
            (right - px, (right, py)),
            (py - bottom, (px, bottom)),
            (top - py, (px, top)),
        )
        return min(candidates, key=lambda item: item[0])[1]
    return (min(max(px, left), right), min(max(py, bottom), top))


def _speech_vertices(
    left: float,
    bottom: float,
    right: float,
    top: float,
    px: float,
    py: float,
    tail: float = 12.0,
) -> tuple[tuple[float, float], ...]:
    ax, ay = _closest_point_on_rect(px, py, left, bottom, right, top)
    edge_scores = {
        "bottom": abs(ay - bottom) + (0.0 if left - 0.5 <= ax <= right + 0.5 else 1e6),
        "top": abs(ay - top) + (0.0 if left - 0.5 <= ax <= right + 0.5 else 1e6),
        "left": abs(ax - left) + (0.0 if bottom - 0.5 <= ay <= top + 0.5 else 1e6),
        "right": abs(ax - right) + (0.0 if bottom - 0.5 <= ay <= top + 0.5 else 1e6),
    }
    edge = min(edge_scores, key=edge_scores.get)
    if edge == "bottom":
        a1 = (max(left, ax - tail), bottom)
        a2 = (min(right, ax + tail), bottom)
        return ((left, bottom), a1, (px, py), a2, (right, bottom), (right, top), (left, top))
    if edge == "right":
        a1 = (right, max(bottom, ay - tail))
        a2 = (right, min(top, ay + tail))
        return ((left, bottom), (right, bottom), a1, (px, py), a2, (right, top), (left, top))
    if edge == "top":
        a1 = (min(right, ax + tail), top)
        a2 = (max(left, ax - tail), top)
        return ((left, bottom), (right, bottom), (right, top), a1, (px, py), a2, (left, top))
    a1 = (left, min(top, ay + tail))
    a2 = (left, max(bottom, ay - tail))
    return ((left, bottom), (right, bottom), (right, top), (left, top), a1, (px, py), a2)


def _cloud_vertices(
    left: float,
    bottom: float,
    right: float,
    top: float,
    px: float,
    py: float,
) -> tuple[tuple[float, float], ...]:
    cx = (left + right) / 2.0
    cy = (bottom + top) / 2.0
    rx = max((right - left) / 2.0, 8.0) * 1.06
    ry = max((top - bottom) / 2.0, 8.0) * 1.18
    count = 28
    points = []
    for index in range(count):
        angle = index / count * 2.0 * math.pi
        bump = 1.0 + 0.16 * math.cos(index * 2.15)
        points.append((cx + rx * bump * math.cos(angle), cy + ry * bump * math.sin(angle)))
    nearest = min(range(count), key=lambda index: (points[index][0] - px) ** 2 + (points[index][1] - py) ** 2)
    points[nearest] = (px, py)
    return tuple(points)


def callout_layout(markup: MarkupAnnotation) -> CalloutLayout:
    """Pointer is (x0, y0). Box is either explicit or centred on (x1, y1)."""

    px, py = markup.x0, markup.y0
    if markup.box_width > 1 and markup.box_height > 1:
        left = markup.x1
        bottom = markup.y1
        right = left + markup.box_width
        top = bottom + markup.box_height
        if left - 2 <= px <= right + 2 and bottom - 2 <= py <= top + 2:
            px = left - 28.0
            py = bottom - 16.0
    else:
        cx, cy = markup.x1, markup.y1
        if abs(cx - px) < 8 and abs(cy - py) < 8:
            cx = px + 130.0
            cy = py + 48.0
        width = min(max(abs(cx - px) * 0.45, 120.0), 260.0)
        height = min(max(abs(cy - py) * 0.28, 40.0), 90.0)
        left = cx - width / 2.0
        right = cx + width / 2.0
        bottom = cy - height / 2.0
        top = cy + height / 2.0
    attach = _closest_point_on_rect(px, py, left, bottom, right, top)
    text = (markup.contents or "").strip() or "註解"
    vertices = (
        _cloud_vertices(left, bottom, right, top, px, py)
        if markup.kind == "cloud"
        else _speech_vertices(left, bottom, right, top, px, py)
    )
    path = comment_polyline(left, bottom, right, top, px, py) if markup.kind == "comment" else ()
    return CalloutLayout((px, py), left, bottom, right, top, attach, vertices, text, path)


CALLOUT_POINTER_GAP = 18.0
COMMENT_STUB = 16.0
COMMENT_ARM = 80.0
COMMENT_TIP = 12.0


def callout_box_from_pointer(
    pointer_x: float,
    pointer_y: float,
    rect_width: float,
    rect_height: float,
    gap: float = CALLOUT_POINTER_GAP,
) -> tuple[float, float]:
    """Return the text-box origin so the tail/arrow sits on the pointer."""

    return (pointer_x - max(float(rect_width), 1.0) / 2.0, pointer_y + gap)


def comment_box_from_pointer(
    pointer_x: float,
    pointer_y: float,
    rect_width: float,
    rect_height: float,
    arm: float = COMMENT_ARM,
    stub: float = COMMENT_STUB,
    tip: float = COMMENT_TIP,
) -> tuple[float, float]:
    """Place the note box above-right of the arrow so the cranked arm is visible."""

    width = max(float(rect_width), 1.0)
    box_bottom = pointer_y + stub + tip
    center_x = pointer_x + arm
    return (center_x - width / 2.0, box_bottom)


def comment_polyline(
    left: float,
    bottom: float,
    right: float,
    top: float,
    pointer_x: float,
    pointer_y: float,
    stub: float = COMMENT_STUB,
) -> tuple[tuple[float, float], tuple[float, float], tuple[float, float], tuple[float, float]]:
    """Box-bottom → drop → horizontal arm → arrow tip."""

    center_x = (left + right) / 2.0
    if pointer_y <= bottom - 4:
        knee_y = min(bottom - stub, (bottom + pointer_y) / 2.0)
        knee_y = max(pointer_y + 8.0, knee_y)
        attach = (center_x, bottom)
        return (attach, (center_x, knee_y), (pointer_x, knee_y), (pointer_x, pointer_y))
    if pointer_y >= top + 4:
        knee_y = max(top + stub, (top + pointer_y) / 2.0)
        knee_y = min(pointer_y - 8.0, knee_y)
        attach = (center_x, top)
        return (attach, (center_x, knee_y), (pointer_x, knee_y), (pointer_x, pointer_y))
    knee_y = bottom - stub
    return ((center_x, bottom), (center_x, knee_y), (pointer_x, knee_y), (pointer_x, pointer_y))


def _polyline_arrowhead(
    tip: tuple[float, float],
    previous: tuple[float, float],
    size: float = 8.0,
) -> tuple[tuple[float, float], tuple[float, float], tuple[float, float]]:
    dx = tip[0] - previous[0]
    dy = tip[1] - previous[1]
    length = math.hypot(dx, dy) or 1.0
    ux, uy = dx / length, dy / length
    base_x = tip[0] - ux * size
    base_y = tip[1] - uy * size
    px, py = -uy * size * 0.55, ux * size * 0.55
    return (tip, (base_x + px, base_y + py), (base_x - px, base_y - py))


def paint_callout_markup(draw, markup: MarkupAnnotation, to_image, draw_text: bool = True) -> None:
    """Draw a callout / speech / cloud shape onto a PIL overlay."""

    from PIL import ImageFont

    layout = callout_layout(markup)
    red, green, blue = (int(max(0.0, min(channel, 1.0)) * 255) for channel in markup.color_rgb)
    fill_rgb = resolve_annotation_fill(markup.color_rgb, markup.fill_rgb, markup.fill_none)
    fill = None
    if fill_rgb is not None:
        fill = (
            int(fill_rgb[0] * 255),
            int(fill_rgb[1] * 255),
            int(fill_rgb[2] * 255),
            230,
        )
    outline = (red, green, blue, 255)

    def image_point(x: float, y: float) -> tuple[float, float]:
        return to_image(x, y)

    p1 = image_point(layout.left, layout.top)
    p2 = image_point(layout.right, layout.bottom)
    box = (min(p1[0], p2[0]), min(p1[1], p2[1]), max(p1[0], p2[0]), max(p1[1], p2[1]))
    if markup.kind == "comment":
        if fill is None:
            draw.rectangle(box, outline=outline, width=2)
        else:
            draw.rectangle(box, fill=fill, outline=outline, width=2)
        path = [image_point(x, y) for x, y in layout.path]
        if len(path) >= 2:
            draw.line(path, fill=outline, width=2)
            head = _polyline_arrowhead(layout.path[-1], layout.path[-2])
            draw.polygon([image_point(x, y) for x, y in head], fill=outline)
    elif markup.kind == "callout":
        if fill is None:
            draw.rectangle(box, outline=outline, width=2)
        else:
            draw.rectangle(box, fill=fill, outline=outline, width=2)
        pointer = image_point(*layout.pointer)
        attach = image_point(*layout.attach)
        draw.line([attach, pointer], fill=outline, width=2)
        size = 7
        draw.polygon(
            [
                pointer,
                (pointer[0] - size, pointer[1] - size),
                (pointer[0] + size, pointer[1] - size),
            ],
            fill=outline,
        )
    else:
        if fill is None:
            draw.polygon([image_point(x, y) for x, y in layout.vertices], outline=outline)
        else:
            draw.polygon([image_point(x, y) for x, y in layout.vertices], fill=fill, outline=outline)
    if not draw_text:
        return
    try:
        font = ImageFont.load_default()
    except Exception:
        font = None
    draw.text((box[0] + 6, box[1] + 5), layout.text[:28], fill=(28, 28, 28, 255), font=font)


def _build_freetext_callout(markup: MarkupAnnotation) -> DictionaryObject:
    layout = callout_layout(markup)
    color = _color_array(markup.color_rgb)
    fill = resolve_annotation_fill(markup.color_rgb, markup.fill_rgb, markup.fill_none)
    annotation = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/FreeText"),
            NameObject("/IT"): NameObject("/FreeTextCallout"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(layout.left),
                    FloatObject(layout.bottom),
                    FloatObject(layout.right),
                    FloatObject(layout.top),
                ]
            ),
            NameObject("/Contents"): TextStringObject(layout.text),
            NameObject("/DA"): TextStringObject("/Helv 10 Tf 0 0 0 rg"),
            NameObject("/C"): color,
            NameObject("/CL"): ArrayObject(
                [
                    FloatObject(layout.pointer[0]),
                    FloatObject(layout.pointer[1]),
                    FloatObject(layout.attach[0]),
                    FloatObject(layout.attach[1]),
                ]
            ),
            NameObject("/LE"): NameObject("/OpenArrow"),
            NameObject("/BS"): DictionaryObject(
                {NameObject("/W"): NumberObject(1), NameObject("/S"): NameObject("/S")}
            ),
            NameObject("/F"): NumberObject(4),
        }
    )
    if fill is not None:
        annotation[NameObject("/IC")] = _color_array(fill)
    return annotation


def _build_shape_polygon(markup: MarkupAnnotation) -> DictionaryObject:
    layout = callout_layout(markup)
    vertices = ArrayObject()
    for x, y in layout.vertices:
        vertices.append(FloatObject(x))
        vertices.append(FloatObject(y))
    xs = [point[0] for point in layout.vertices]
    ys = [point[1] for point in layout.vertices]
    fill = resolve_annotation_fill(markup.color_rgb, markup.fill_rgb, markup.fill_none)
    annotation = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/Polygon"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(min(xs) - 2),
                    FloatObject(min(ys) - 2),
                    FloatObject(max(xs) + 2),
                    FloatObject(max(ys) + 2),
                ]
            ),
            NameObject("/Vertices"): vertices,
            NameObject("/C"): _color_array(markup.color_rgb),
            NameObject("/BS"): DictionaryObject({NameObject("/W"): NumberObject(1)}),
            NameObject("/Contents"): TextStringObject(layout.text),
            NameObject("/F"): NumberObject(4),
        }
    )
    if fill is not None:
        annotation[NameObject("/IC")] = _color_array(fill)
    if markup.kind == "cloud":
        annotation[NameObject("/BE")] = DictionaryObject(
            {NameObject("/S"): NameObject("/C"), NameObject("/I"): NumberObject(1)}
        )
    return annotation


def _build_box_freetext(markup: MarkupAnnotation) -> DictionaryObject:
    layout = callout_layout(markup)
    fill = resolve_annotation_fill(markup.color_rgb, markup.fill_rgb, markup.fill_none)
    annotation = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/FreeText"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(layout.left + 4),
                    FloatObject(layout.bottom + 4),
                    FloatObject(layout.right - 4),
                    FloatObject(layout.top - 4),
                ]
            ),
            NameObject("/Contents"): TextStringObject(layout.text),
            NameObject("/DA"): TextStringObject("/Helv 10 Tf 0 0 0 rg"),
            NameObject("/BS"): DictionaryObject({NameObject("/W"): NumberObject(0)}),
            NameObject("/F"): NumberObject(4),
        }
    )
    if fill is not None:
        annotation[NameObject("/C")] = _color_array(fill)
    return annotation


def _build_comment_square(markup: MarkupAnnotation) -> DictionaryObject:
    layout = callout_layout(markup)
    fill = resolve_annotation_fill(markup.color_rgb, markup.fill_rgb, markup.fill_none)
    annotation = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/Square"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(layout.left),
                    FloatObject(layout.bottom),
                    FloatObject(layout.right),
                    FloatObject(layout.top),
                ]
            ),
            NameObject("/C"): _color_array(markup.color_rgb),
            NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), FloatObject(1.5)]),
            NameObject("/F"): NumberObject(4),
        }
    )
    if fill is not None:
        annotation[NameObject("/IC")] = _color_array(fill)
    return annotation


def _build_comment_polyline(markup: MarkupAnnotation) -> DictionaryObject:
    layout = callout_layout(markup)
    points = layout.path or (layout.attach, layout.pointer)
    vertices = ArrayObject()
    xs: list[float] = []
    ys: list[float] = []
    for x, y in points:
        vertices.append(FloatObject(x))
        vertices.append(FloatObject(y))
        xs.append(x)
        ys.append(y)
    pad = 12.0
    return DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/PolyLine"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(min(xs) - pad),
                    FloatObject(min(ys) - pad),
                    FloatObject(max(xs) + pad),
                    FloatObject(max(ys) + pad),
                ]
            ),
            NameObject("/Vertices"): vertices,
            NameObject("/LE"): ArrayObject([NameObject("/None"), NameObject("/ClosedArrow")]),
            NameObject("/C"): _color_array(markup.color_rgb),
            NameObject("/IC"): _color_array(markup.color_rgb),
            NameObject("/BS"): DictionaryObject(
                {NameObject("/W"): FloatObject(1.5), NameObject("/S"): NameObject("/S")}
            ),
            NameObject("/F"): NumberObject(4),
        }
    )


def iter_markup_pdf_annotations(markup: MarkupAnnotation) -> list[DictionaryObject]:
    if markup.kind == "comment":
        return [_build_comment_square(markup), _build_comment_polyline(markup), _build_box_freetext(markup)]
    if markup.kind in {"speech", "cloud"}:
        return [_build_shape_polygon(markup), _build_box_freetext(markup)]
    return [build_markup_annotation(markup)]


def build_markup_annotation(markup: MarkupAnnotation) -> DictionaryObject:
    kind = markup.kind
    if kind not in MARKUP_SUBTYPES:
        raise ValueError(f"未知的標註類型：{kind}")
    if kind == "callout":
        return _build_freetext_callout(markup)
    if kind in {"speech", "cloud"}:
        return _build_shape_polygon(markup)
    color = _color_array(markup.color_rgb)

    if kind in {"highlight", "underline", "strikeout"}:
        left, bottom, right, top = _normalized_rect(markup.x0, markup.y0, markup.x1, markup.y1)
        quad = ArrayObject(
            [
                FloatObject(left),
                FloatObject(top),
                FloatObject(right),
                FloatObject(top),
                FloatObject(left),
                FloatObject(bottom),
                FloatObject(right),
                FloatObject(bottom),
            ]
        )
        annotation = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject(MARKUP_SUBTYPES[kind]),
                NameObject("/Rect"): ArrayObject(
                    [FloatObject(left), FloatObject(bottom), FloatObject(right), FloatObject(top)]
                ),
                NameObject("/QuadPoints"): quad,
                NameObject("/C"): color,
                NameObject("/F"): NumberObject(4),
            }
        )
        if markup.contents:
            annotation[NameObject("/Contents")] = TextStringObject(markup.contents)
        return annotation

    if kind in {"rect", "ellipse"}:
        left, bottom, right, top = _normalized_rect(markup.x0, markup.y0, markup.x1, markup.y1)
        annotation = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject(MARKUP_SUBTYPES[kind]),
                NameObject("/Rect"): ArrayObject(
                    [FloatObject(left), FloatObject(bottom), FloatObject(right), FloatObject(top)]
                ),
                NameObject("/C"): color,
                NameObject("/BS"): DictionaryObject(
                    {NameObject("/W"): NumberObject(2), NameObject("/S"): NameObject("/S")}
                ),
                NameObject("/F"): NumberObject(4),
            }
        )
        if markup.contents:
            annotation[NameObject("/Contents")] = TextStringObject(markup.contents)
        return annotation

    if kind in {"line", "arrow"}:
        left, bottom, right, top = _normalized_rect(markup.x0, markup.y0, markup.x1, markup.y1)
        padding = 6.0
        annotation = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/Line"),
                NameObject("/L"): ArrayObject(
                    [
                        FloatObject(markup.x0),
                        FloatObject(markup.y0),
                        FloatObject(markup.x1),
                        FloatObject(markup.y1),
                    ]
                ),
                NameObject("/Rect"): ArrayObject(
                    [
                        FloatObject(left - padding),
                        FloatObject(bottom - padding),
                        FloatObject(right + padding),
                        FloatObject(top + padding),
                    ]
                ),
                NameObject("/C"): color,
                NameObject("/BS"): DictionaryObject({NameObject("/W"): NumberObject(2)}),
                NameObject("/F"): NumberObject(4),
            }
        )
        if kind == "arrow":
            annotation[NameObject("/LE")] = ArrayObject([NameObject("/None"), NameObject("/OpenArrow")])
        if markup.contents:
            annotation[NameObject("/Contents")] = TextStringObject(markup.contents)
        return annotation

    # sticky note
    size = 18.0
    annotation = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Annot"),
            NameObject("/Subtype"): NameObject("/Text"),
            NameObject("/Rect"): ArrayObject(
                [
                    FloatObject(markup.x0),
                    FloatObject(markup.y0 - size),
                    FloatObject(markup.x0 + size),
                    FloatObject(markup.y0),
                ]
            ),
            NameObject("/Contents"): TextStringObject(markup.contents or "備註"),
            NameObject("/Name"): NameObject("/Comment"),
            NameObject("/Open"): BooleanObject(False),
            NameObject("/C"): color,
            NameObject("/F"): NumberObject(4),
        }
    )
    return annotation


def apply_markup_annotations(
    source: Path,
    target: Path,
    markups: list[tuple[int, MarkupAnnotation]],
    password: str = "",
) -> int:
    if not markups:
        raise ValueError("請先加入至少一個標註。")
    reader = open_reader(source, password)
    page_count = len(reader.pages)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    applied = 0
    for page_index, markup in markups:
        if page_index < 0 or page_index >= page_count:
            continue
        for annotation in iter_markup_pdf_annotations(markup):
            add_annotation_to_page(writer, writer.pages[page_index], annotation)
        applied += 1
    if applied == 0:
        raise ValueError("沒有可套用的標註。")
    write_pdf(writer, target)
    return applied


# --- page cropping -----------------------------------------------------------

def crop_pdf_pages(
    source: Path,
    target: Path,
    rect: tuple[float, float, float, float],
    pages_spec: str = "",
    password: str = "",
) -> int:
    """Crop selected pages to ``rect`` (left, bottom, right, top) in PDF points.

    The rectangle is the area to keep and is clamped to each page's media box.
    """

    left, bottom, right, top = rect
    if right <= left or top <= bottom:
        raise ValueError("裁切範圍無效，請重新框選或輸入正確數值。")

    reader = open_reader(source, password)
    page_count = len(reader.pages)
    selected = (
        set(parse_pages(pages_spec, page_count))
        if (pages_spec or "").strip()
        else set(range(page_count))
    )
    if not selected:
        raise ValueError("沒有要裁切的頁面。")

    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    cropped = 0
    for index, page in enumerate(writer.pages):
        if index not in selected:
            continue
        media = page.mediabox
        new_left = max(left, float(media.left))
        new_bottom = max(bottom, float(media.bottom))
        new_right = min(right, float(media.right))
        new_top = min(top, float(media.top))
        if new_right <= new_left or new_top <= new_bottom:
            continue
        box = RectangleObject([new_left, new_bottom, new_right, new_top])
        page.cropbox = box
        page.trimbox = box
        cropped += 1
    if cropped == 0:
        raise ValueError("裁切範圍與頁面沒有重疊，未輸出任何頁面。")
    write_pdf(writer, target)
    return cropped


# --- advanced split ----------------------------------------------------------

def parse_page_groups(spec: str, page_count: int) -> list[list[int]]:
    """Parse a spec where each comma-separated part becomes its own group.

    Example: ``"1-3,4-6,7"`` -> ``[[0,1,2], [3,4,5], [6]]``.
    """

    spec = (spec or "").strip()
    if not spec:
        raise ValueError("請輸入範圍，例如 1-3,4-6,7-9。")
    groups: list[list[int]] = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start_text, end_text = part.split("-", 1)
            start, end = int(start_text), int(end_text)
            if start > end:
                raise ValueError(f"頁碼範圍不正確：{part}")
            pages = list(range(start - 1, end))
        else:
            pages = [int(part) - 1]
        for page in pages:
            if page < 0 or page >= page_count:
                raise ValueError(f"頁碼超出範圍：{page + 1}，目前共有 {page_count} 頁。")
        groups.append(pages)
    if not groups:
        raise ValueError("沒有有效的拆分範圍。")
    return groups


def split_pdf_advanced(
    source: Path,
    target_zip: Path,
    mode: str,
    value: str,
    password: str = "",
) -> int:
    """Split a PDF into multiple files inside a ZIP.

    ``mode="every"`` groups every ``value`` pages into one file. ``mode="ranges"``
    treats ``value`` as a comma list where each part becomes one output file.
    Returns the number of output PDFs.
    """

    reader = open_reader(source, password)
    page_count = len(reader.pages)
    if page_count == 0:
        raise ValueError("PDF 沒有頁面可拆分。")

    if mode == "every":
        try:
            size = int(str(value).strip())
        except (TypeError, ValueError):
            raise ValueError("請輸入每檔頁數（正整數）。")
        if size < 1:
            raise ValueError("每檔頁數必須至少為 1。")
        groups = [list(range(start, min(start + size, page_count))) for start in range(0, page_count, size)]
    elif mode == "ranges":
        groups = parse_page_groups(value, page_count)
    else:
        raise ValueError(f"未知的拆分模式：{mode}")

    stem = safe_output_name(source.stem)
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        generated: list[Path] = []
        for part_index, pages in enumerate(groups, start=1):
            if not pages:
                continue
            writer = PdfWriter()
            for page_index in pages:
                writer.add_page(reader.pages[page_index])
            first = pages[0] + 1
            last = pages[-1] + 1
            label = f"{first}" if first == last else f"{first}-{last}"
            part = temp_path / safe_output_name(f"{part_index:03d}-{stem}-pages-{label}.pdf")
            write_pdf(writer, part)
            generated.append(part)
        if not generated:
            raise ValueError("沒有可輸出的拆分檔案。")
        with zipfile.ZipFile(target_zip, "w", zipfile.ZIP_DEFLATED) as archive:
            for item in generated:
                archive.write(item, item.name)
    return len(generated)


# --- Bates numbering ---------------------------------------------------------

BATES_POSITIONS = {
    "bottom-right": "右下",
    "bottom-left": "左下",
    "top-right": "右上",
    "top-left": "左上",
}


def add_bates_numbering(
    source: Path,
    target: Path,
    prefix: str = "",
    start: int = 1,
    digits: int = 6,
    suffix: str = "",
    position: str = "bottom-right",
    password: str = "",
) -> int:
    """Stamp sequential Bates numbers (e.g. ABC-000001) on every page."""

    if digits < 1 or digits > 12:
        raise ValueError("位數請介於 1 至 12。")
    if position not in BATES_POSITIONS:
        raise ValueError("不支援的 Bates 位置。")

    reader = open_reader(source, password)
    writer = PdfWriter()
    writer.append_pages_from_reader(reader)
    for index, page in enumerate(writer.pages):
        number = start + index
        text = f"{prefix}{number:0{digits}d}{suffix}"
        width = float(page.mediabox.width)
        height = float(page.mediabox.height)
        rect_width = max(120.0, len(text) * 8.0)
        margin = 24.0
        if position.endswith("right"):
            x = max(width - rect_width - margin, margin)
        else:
            x = margin
        if position.startswith("bottom"):
            y = margin
        else:
            y = max(height - margin - 18.0, margin)
        annotation = DictionaryObject(
            {
                NameObject("/Type"): NameObject("/Annot"),
                NameObject("/Subtype"): NameObject("/FreeText"),
                NameObject("/Rect"): ArrayObject(
                    [FloatObject(x), FloatObject(y), FloatObject(x + rect_width), FloatObject(y + 20)]
                ),
                NameObject("/Contents"): TextStringObject(text),
                NameObject("/DA"): TextStringObject("/Helv 10 Tf 0.1 0.1 0.1 rg"),
                NameObject("/Border"): ArrayObject([NumberObject(0), NumberObject(0), NumberObject(0)]),
                NameObject("/F"): NumberObject(4),
            }
        )
        add_annotation_to_page(writer, page, annotation)
    write_pdf(writer, target)
    return len(writer.pages)


# --- permission control ------------------------------------------------------

def encrypt_pdf_with_permissions(
    source: Path,
    target: Path,
    owner_password: str,
    allow_print: bool = True,
    allow_copy: bool = True,
    allow_modify: bool = True,
    user_password: str = "",
    password: str = "",
) -> None:
    """Encrypt with an owner password and restrict printing/copying/modifying.

    ``user_password`` (empty by default) lets recipients open the file without a
    password while still being bound by the permission flags.
    """

    if not owner_password:
        raise ValueError("請輸入擁有者密碼（用來鎖定權限）。")

    reader = open_reader(source, password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    permissions = UserAccessPermissions(0)
    if allow_print:
        permissions |= UserAccessPermissions.PRINT | UserAccessPermissions.PRINT_TO_REPRESENTATION
    if allow_copy:
        permissions |= (
            UserAccessPermissions.EXTRACT | UserAccessPermissions.EXTRACT_TEXT_AND_GRAPHICS
        )
    if allow_modify:
        permissions |= (
            UserAccessPermissions.MODIFY
            | UserAccessPermissions.ADD_OR_MODIFY
            | UserAccessPermissions.FILL_FORM_FIELDS
            | UserAccessPermissions.ASSEMBLE_DOC
        )

    writer.encrypt(
        user_password=user_password,
        owner_password=owner_password,
        permissions_flag=permissions,
    )
    write_pdf(writer, target)


# --- Adobe-Pro-like page / stamp / compare / bookmark tools ------------------

TEXT_MARKUP_KINDS = frozenset({"highlight", "underline", "strikeout"})


def _resolve_page_indexes(reader: PdfReader, pages_spec: str) -> list[int]:
    if (pages_spec or "").strip():
        return parse_pages(pages_spec, len(reader.pages))
    return list(range(len(reader.pages)))


def _pymupdf_search_rects(page, query: str, case_sensitive: bool = False) -> list:
    needle = (query or "").strip()
    if not needle:
        return []
    if case_sensitive:
        return page.search_for(needle)
    rects: list = []
    for x0, y0, x1, y1, word, *_rest in page.get_text("words"):
        if needle.lower() in word.lower():
            rects.append(fitz.Rect(x0, y0, x1, y1))
    if rects:
        return rects
    return page.search_for(needle)


def insert_pdf_pages(
    source: Path,
    insert_from: Path,
    target: Path,
    at_index: int,
    pages_spec: str = "",
    password: str = "",
    insert_password: str = "",
) -> int:
    """Insert pages from ``insert_from`` into ``source`` at ``at_index`` (0-based)."""

    source_reader = open_reader(source, password)
    insert_reader = open_reader(insert_from, insert_password)
    source_count = len(source_reader.pages)
    if at_index < 0 or at_index > source_count:
        raise ValueError(f"插入位置超出範圍：{at_index + 1}，目前共有 {source_count} 頁。")

    insert_indexes = _resolve_page_indexes(insert_reader, pages_spec)
    if not insert_indexes:
        raise ValueError("沒有可插入的頁面。")

    writer = PdfWriter()
    for index in range(at_index):
        writer.add_page(source_reader.pages[index])
    for index in insert_indexes:
        writer.add_page(insert_reader.pages[index])
    for index in range(at_index, source_count):
        writer.add_page(source_reader.pages[index])
    write_pdf(writer, target)
    return len(insert_indexes)


def replace_pdf_pages(
    source: Path,
    replacement: Path,
    target: Path,
    start_index: int,
    pages_spec: str = "",
    password: str = "",
    replacement_password: str = "",
) -> int:
    """Replace consecutive pages in ``source`` with pages from ``replacement``."""

    source_reader = open_reader(source, password)
    replacement_reader = open_reader(replacement, replacement_password)
    source_count = len(source_reader.pages)
    replacement_indexes = _resolve_page_indexes(replacement_reader, pages_spec)
    if not replacement_indexes:
        raise ValueError("沒有可替換的頁面。")

    replace_count = len(replacement_indexes)
    if start_index < 0 or start_index >= source_count:
        raise ValueError(f"替換起始位置超出範圍：{start_index + 1}，目前共有 {source_count} 頁。")
    if start_index + replace_count > source_count:
        raise ValueError(
            f"替換範圍超出文件頁數：需 {replace_count} 頁，從第 {start_index + 1} 頁起不足。"
        )

    writer = PdfWriter()
    for index in range(start_index):
        writer.add_page(source_reader.pages[index])
    for index in replacement_indexes:
        writer.add_page(replacement_reader.pages[index])
    for index in range(start_index + replace_count, source_count):
        writer.add_page(source_reader.pages[index])
    write_pdf(writer, target)
    return replace_count


def add_image_stamp(
    source: Path,
    target: Path,
    image_path: Path,
    page_index: int,
    x: float,
    y: float,
    width: float,
    height: float,
    password: str = "",
) -> None:
    """Place an image stamp/signature on a page using PyMuPDF."""

    ensure_pymupdf_available()
    image_path = Path(image_path)
    if not image_path.exists():
        raise ValueError(f"找不到圖片檔案：{image_path.name}")

    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        if page_index < 0 or page_index >= document.page_count:
            raise ValueError("頁碼超出範圍。")
        page = document[page_index]
        rect = fitz.Rect(x, y, x + width, y + height)
        page.insert_image(rect, filename=str(image_path))
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()


def add_signature_image(
    source: Path,
    target: Path,
    image_path: Path,
    page_index: int,
    x: float,
    y: float,
    width: float,
    height: float,
    password: str = "",
) -> None:
    """Alias for :func:`add_image_stamp` — place a signature image on a page."""

    add_image_stamp(source, target, image_path, page_index, x, y, width, height, password)


def add_text_stamp(
    source: Path,
    target: Path,
    text: str,
    page_index: int,
    x: float,
    y: float,
    width: float,
    height: float,
    color_rgb: tuple[float, float, float] = (0.8, 0.0, 0.0),
    password: str = "",
) -> None:
    """Draw a bordered text stamp on a PDF page using PyMuPDF."""

    ensure_pymupdf_available()
    stamp_text = (text or "").strip()
    if not stamp_text:
        raise ValueError("圖章文字不可為空。")

    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        if page_index < 0 or page_index >= document.page_count:
            raise ValueError("頁碼超出範圍。")
        page = document[page_index]
        rect = fitz.Rect(x, y, x + width, y + height)
        page.draw_rect(rect, color=color_rgb, width=2)
        fontsize = max(8.0, min(height * 0.55, width / max(len(stamp_text), 1) * 0.75))
        rc = page.insert_textbox(
            rect,
            stamp_text,
            fontsize=fontsize,
            color=color_rgb,
            align=fitz.TEXT_ALIGN_CENTER,
        )
        if rc < 0:
            page.insert_textbox(
                rect,
                stamp_text,
                fontsize=max(6.0, fontsize * 0.7),
                color=color_rgb,
                align=fitz.TEXT_ALIGN_CENTER,
            )
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()


def compress_pdf_advanced(
    source: Path,
    target: Path,
    image_dpi: int = 100,
    jpeg_quality: int = 60,
    password: str = "",
) -> tuple[int, int]:
    """Re-write a PDF with embedded images re-sampled via PyMuPDF."""

    ensure_pymupdf_available()
    old_size = source.stat().st_size if source.exists() else 0
    document = fitz.open(str(source))
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        document.rewrite_images(
            dpi_threshold=image_dpi,
            dpi_target=image_dpi,
            quality=jpeg_quality,
            lossy=True,
        )
        document.save(str(target), garbage=4, deflate=True, clean=True)
    finally:
        document.close()
    new_size = target.stat().st_size if target.exists() else 0
    return old_size, new_size


def compare_pdf_text(
    left: Path,
    right: Path,
    target_report: Path,
    left_password: str = "",
    right_password: str = "",
) -> int:
    """Write a UTF-8 page-by-page text comparison report; return differing page count."""

    left_reader = open_reader(left, left_password)
    right_reader = open_reader(right, right_password)
    left_count = len(left_reader.pages)
    right_count = len(right_reader.pages)
    max_pages = max(left_count, right_count)
    diff_count = 0
    lines = [
        "Victor PDF Tools Box - Text Compare",
        f"Left: {left.name}",
        f"Right: {right.name}",
        "",
    ]

    for index in range(max_pages):
        left_text = (
            (left_reader.pages[index].extract_text() or "").strip() if index < left_count else ""
        )
        right_text = (
            (right_reader.pages[index].extract_text() or "").strip()
            if index < right_count
            else ""
        )
        if left_text != right_text:
            diff_count += 1
            lines.append(f"--- Page {index + 1} DIFF ---")
            lines.append("[Left]")
            lines.append(left_text or "(empty)")
            lines.append("[Right]")
            lines.append(right_text or "(empty)")
            lines.append("")

    if diff_count == 0:
        lines.append("All pages match.")
    target_report.write_text("\n".join(lines), encoding="utf-8")
    return diff_count


def split_pdf_by_bookmarks(source: Path, target_zip: Path, password: str = "") -> int:
    """Split a PDF into one file per top-level bookmark range, packed in a ZIP."""

    reader = open_reader(source, password)
    page_count = len(reader.pages)
    items = extract_outline(source, password)
    top_level = [item for item in items if item.level == 0]
    if not top_level:
        raise ValueError("PDF 沒有書籤，無法依書籤拆分。")

    sorted_items = sorted(top_level, key=lambda item: item.page_index)
    ranges: list[tuple[str, int, int]] = []
    for index, item in enumerate(sorted_items):
        start = item.page_index
        if index + 1 < len(sorted_items):
            end = sorted_items[index + 1].page_index - 1
        else:
            end = page_count - 1
        if end < start:
            end = start
        ranges.append((item.title, start, end))

    stem = safe_output_name(source.stem)
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        generated: list[Path] = []
        for part_index, (title, start, end) in enumerate(ranges, start=1):
            writer = PdfWriter()
            for page_index in range(start, end + 1):
                writer.add_page(reader.pages[page_index])
            part = temp_path / safe_output_name(f"{part_index:03d}-{stem}-{title}.pdf")
            write_pdf(writer, part)
            generated.append(part)
        with zipfile.ZipFile(target_zip, "w", zipfile.ZIP_DEFLATED) as archive:
            for item in generated:
                archive.write(item, item.name)
    return len(generated)


def apply_text_markups_for_query(
    source: Path,
    target: Path,
    query: str,
    kind: str = "highlight",
    color_rgb: tuple[float, float, float] = (1.0, 0.92, 0.23),
    password: str = "",
    case_sensitive: bool = False,
) -> int:
    """Search for ``query`` and add highlight/underline/strikeout annotations."""

    if not (query or "").strip():
        raise ValueError("請輸入要標註的搜尋文字。")
    if kind not in TEXT_MARKUP_KINDS:
        raise ValueError(f"不支援的標註類型：{kind}，請使用 highlight、underline 或 strikeout。")
    ensure_pymupdf_available()

    annot_adders = {
        "highlight": lambda page, rect: page.add_highlight_annot(rect),
        "underline": lambda page, rect: page.add_underline_annot(rect),
        "strikeout": lambda page, rect: page.add_strikeout_annot(rect),
    }
    add_annot = annot_adders[kind]

    document = fitz.open(str(source))
    match_count = 0
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        for page in document:
            for rect in _pymupdf_search_rects(page, query, case_sensitive):
                annotation = add_annot(page, rect)
                if annotation is not None:
                    annotation.set_colors(stroke=color_rgb)
                    annotation.update()
                    match_count += 1
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()
    return match_count


def flatten_form_fields(source: Path, target: Path, password: str = "") -> int:
    """Flatten PDF form widgets into page content when present."""

    ensure_pymupdf_available()
    document = fitz.open(str(source))
    field_count = 0
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        for page in document:
            for _widget in page.widgets() or []:
                field_count += 1
        if field_count:
            document.bake(widgets=True, annots=False)
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()
    return field_count


def _form_widget_type_name(field_type: int) -> str:
    if fitz is not None:
        mapping = {
            int(getattr(fitz, "PDF_WIDGET_TYPE_BUTTON", 1)): "button",
            int(getattr(fitz, "PDF_WIDGET_TYPE_CHECKBOX", 2)): "checkbox",
            int(getattr(fitz, "PDF_WIDGET_TYPE_RADIOBUTTON", 5)): "radio",
            int(getattr(fitz, "PDF_WIDGET_TYPE_TEXT", 7)): "text",
            int(getattr(fitz, "PDF_WIDGET_TYPE_LISTBOX", 4)): "listbox",
            int(getattr(fitz, "PDF_WIDGET_TYPE_COMBOBOX", 3)): "combobox",
            int(getattr(fitz, "PDF_WIDGET_TYPE_SIGNATURE", 6)): "signature",
        }
        return mapping.get(int(field_type or 0), "unknown")
    return FORM_WIDGET_TYPE_NAMES.get(int(field_type or 0), "unknown")


def _open_pymupdf_document(source: Path, password: str = ""):
    ensure_pymupdf_available()
    document = fitz.open(str(source))
    if document.is_encrypted:
        if not password or document.authenticate(password) == 0:
            document.close()
            raise ValueError(f"{source.name} 已加密，請輸入密碼。")
    return document


def list_form_fields(source: Path, password: str = "") -> list[FormField]:
    """Return AcroForm widgets in page order."""

    document = _open_pymupdf_document(source, password)
    fields: list[FormField] = []
    try:
        for page_index, page in enumerate(document):
            for widget in page.widgets() or []:
                choices = tuple(str(item) for item in (getattr(widget, "choice_values", None) or []))
                fields.append(
                    FormField(
                        name=str(widget.field_name or ""),
                        page_index=page_index,
                        field_type=_form_widget_type_name(widget.field_type),
                        value="" if widget.field_value is None else str(widget.field_value),
                        choices=choices,
                    )
                )
    finally:
        document.close()
    return fields


def fill_form_fields(
    source: Path,
    target: Path,
    values: dict[str, str],
    password: str = "",
) -> int:
    """Write ``values`` keyed by field name into widgets, then save."""

    if not values:
        raise ValueError("請輸入要填寫的欄位值。")
    document = _open_pymupdf_document(source, password)
    updated = 0
    try:
        for page in document:
            for widget in page.widgets() or []:
                name = str(widget.field_name or "")
                if not name or name not in values:
                    continue
                if _form_widget_type_name(widget.field_type) == "signature":
                    continue
                widget.field_value = values[name]
                widget.update()
                updated += 1
        if updated == 0:
            raise ValueError("沒有符合的表單欄位可填寫。")
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()
    return updated


def _pad_rgb_image(image: Image.Image, size: tuple[int, int], fill=(255, 255, 255)) -> Image.Image:
    rgb = image.convert("RGB")
    if rgb.size == size:
        return rgb
    canvas = Image.new("RGB", size, fill)
    canvas.paste(rgb, (0, 0))
    return canvas


def compare_pdf_visual(
    left: Path,
    right: Path,
    target: Path,
    left_password: str = "",
    right_password: str = "",
    dpi: int = 72,
    threshold: float = 8.0,
) -> int:
    """Render both PDFs, pixel-diff each page, and write a side-by-side PDF.

    Returns the number of pages that differ (including missing pages).
    """

    from PIL import ImageChops, ImageDraw, ImageStat

    if not PDF_RENDER_AVAILABLE or pdfium is None:
        raise ValueError("視覺比對需要 pypdfium2 預覽元件。")
    if dpi < 36 or dpi > 200:
        raise ValueError("視覺比對解析度請介於 36 至 200 DPI。")

    left_pages = render_pdf_page_images(left, left_password, dpi=dpi)
    right_pages = render_pdf_page_images(right, right_password, dpi=dpi)
    max_pages = max(len(left_pages), len(right_pages))
    if max_pages == 0:
        raise ValueError("沒有可比對的頁面。")

    composed: list[Image.Image] = []
    diff_count = 0
    for index in range(max_pages):
        left_img = left_pages[index][1] if index < len(left_pages) else None
        right_img = right_pages[index][1] if index < len(right_pages) else None
        if left_img is None or right_img is None:
            diff_count += 1
            is_diff = True
            ref = left_img or right_img
            width, height = ref.size
            left_pad = (
                left_img.convert("RGB")
                if left_img is not None
                else Image.new("RGB", (width, height), (230, 230, 230))
            )
            right_pad = (
                right_img.convert("RGB")
                if right_img is not None
                else Image.new("RGB", (width, height), (230, 230, 230))
            )
            highlighted = right_pad
        else:
            width = max(left_img.width, right_img.width)
            height = max(left_img.height, right_img.height)
            left_pad = _pad_rgb_image(left_img, (width, height))
            right_pad = _pad_rgb_image(right_img, (width, height))
            diff = ImageChops.difference(left_pad, right_pad)
            stats = ImageStat.Stat(diff)
            mean_val = sum(stats.mean) / max(len(stats.mean), 1)
            is_diff = mean_val > threshold
            if is_diff:
                diff_count += 1
                mask = diff.convert("L").point(lambda pixel: 170 if pixel > max(int(threshold), 4) else 0)
                red = Image.new("RGB", (width, height), (220, 40, 40))
                highlighted = Image.composite(red, right_pad, mask)
            else:
                highlighted = right_pad

        gap = 16
        header = 36
        canvas = Image.new("RGB", (width * 2 + gap + 24, height + header + 16), (36, 40, 44))
        canvas.paste(left_pad, (8, header))
        canvas.paste(highlighted, (16 + width, header))
        draw = ImageDraw.Draw(canvas)
        status = "DIFF" if is_diff else "MATCH"
        color = (255, 180, 80) if is_diff else (140, 220, 160)
        draw.text((8, 8), f"Page {index + 1}  {status}    Left | Right", fill=color)
        composed.append(canvas)

    rgb_pages = [image.convert("RGB") for image in composed]
    rgb_pages[0].save(target, save_all=True, append_images=rgb_pages[1:], resolution=float(dpi))
    return diff_count


def secure_redact_query(
    source: Path,
    target: Path,
    query: str,
    password: str = "",
    case_sensitive: bool = False,
) -> int:
    """Search for ``query`` and permanently redact all matches via PyMuPDF."""

    if not (query or "").strip():
        raise ValueError("請輸入要遮蔽的搜尋文字。")
    ensure_pymupdf_available()

    document = fitz.open(str(source))
    match_count = 0
    try:
        if document.is_encrypted:
            if not password or document.authenticate(password) == 0:
                raise ValueError(f"{source.name} 已加密，請輸入密碼。")
        for page in document:
            rects = _pymupdf_search_rects(page, query, case_sensitive)
            for rect in rects:
                page.add_redact_annot(rect, fill=(0, 0, 0))
                match_count += 1
            if rects:
                page.apply_redactions()
        if match_count == 0:
            raise ValueError(f"找不到要遮蔽的文字：{query.strip()}")
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()
    return match_count


def parse_size_bytes(value: str, default: int | None = None) -> int:
    text = (value or "").strip().lower().replace(" ", "")
    if not text:
        if default is not None:
            return default
        raise ValueError("請輸入檔案大小，例如 10MB。")
    match = re.fullmatch(r"(\d+(?:\.\d+)?)(kb|mb|gb|b|k|m|g)?", text)
    if not match:
        raise ValueError("檔案大小格式無效，例如 10MB、25MB 或 512KB。")
    amount = float(match.group(1))
    unit = match.group(2) or "mb"
    factor = {
        "b": 1,
        "k": 1024,
        "kb": 1024,
        "m": 1024 * 1024,
        "mb": 1024 * 1024,
        "g": 1024 * 1024 * 1024,
        "gb": 1024 * 1024 * 1024,
    }[unit]
    size = int(amount * factor)
    if size < 1:
        raise ValueError("檔案大小必須大於 0。")
    return size


def _copy_pdf_file(source: Path, target: Path) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    if source.resolve() == target.resolve():
        return
    shutil.copy2(source, target)


def _osd_rotate_degrees(image: Image.Image) -> int:
    if not OCR_AVAILABLE or pytesseract is None:
        return 0
    try:
        osd = pytesseract.image_to_osd(image)
    except Exception:
        return 0
    match = re.search(r"Rotate:\s*(\d+)", str(osd or ""))
    if not match:
        return 0
    return int(match.group(1)) % 360


def _rotate_image_clockwise(image: Image.Image, degrees: int) -> Image.Image:
    turn = int(degrees) % 360
    if turn == 90:
        return image.transpose(Image.ROTATE_270)
    if turn == 180:
        return image.transpose(Image.ROTATE_180)
    if turn == 270:
        return image.transpose(Image.ROTATE_90)
    return image


def _upright_scanned_image(image: Image.Image) -> Image.Image:
    degrees = _osd_rotate_degrees(image)
    if not degrees:
        return image
    return _rotate_image_clockwise(image, degrees)


def _binary_row_variance(image: Image.Image) -> float:
    gray = image.convert("L")
    width, height = gray.size
    pixels = gray.tobytes()
    totals: list[int] = []
    for row in range(0, height, 2):
        start = row * width
        totals.append(sum(1 for pixel in pixels[start : start + width] if pixel < 128))
    if len(totals) < 2:
        return 0.0
    mean = sum(totals) / len(totals)
    return sum((value - mean) ** 2 for value in totals) / len(totals)


def estimate_skew_degrees(image: Image.Image) -> float:
    gray = image.convert("L")
    longest = max(gray.size)
    if longest > 720:
        scale = 720 / float(longest)
        resample = getattr(getattr(Image, "Resampling", Image), "BILINEAR", Image.BILINEAR)
        gray = gray.resize(
            (max(1, int(gray.width * scale)), max(1, int(gray.height * scale))),
            resample,
        )
    binary = gray.point(lambda pixel: 0 if pixel < 180 else 255)
    best_angle = 0.0
    best_score = -1.0
    for tenths in range(-70, 71, 5):
        angle = tenths / 10.0
        rotated = binary.rotate(angle, resample=Image.NEAREST, expand=False, fillcolor=255)
        score = _binary_row_variance(rotated)
        if score > best_score:
            best_score = score
            best_angle = angle
    return best_angle


def deskew_image(image: Image.Image, max_degrees: float = 8.0) -> Image.Image:
    angle = estimate_skew_degrees(image)
    if abs(angle) < 0.35 or abs(angle) > max_degrees:
        return image
    return image.rotate(angle, resample=getattr(Image, "BICUBIC", Image.BILINEAR), expand=True, fillcolor="white")


def prepare_scanned_page_image(image: Image.Image) -> Image.Image:
    from PIL import ImageFilter

    upright = _upright_scanned_image(image)
    deskewed = deskew_image(upright)
    cleaned = deskewed.filter(ImageFilter.MedianFilter(size=3))
    if deskewed is not upright and upright is not image:
        upright.close()
    if cleaned is not deskewed and deskewed is not image:
        deskewed.close()
    return cleaned


def cleanup_scanned_pdf(
    source: Path,
    target: Path,
    password: str = "",
    language: str = "auto",
    pages_spec: str = "",
    dpi: int = 300,
    progress=None,
) -> int:
    """Deskew, auto-orient, denoise, then rebuild a searchable PDF."""

    ensure_ocr_available()
    pages, tess_lang, tess_config = _ocr_iter_pages(source, password, pages_spec, language, dpi, progress)
    writer = PdfWriter()
    count = 0
    for page_index, image in pages:
        count += 1
        if progress:
            progress(count - 1, 0, f"掃描件整理第 {page_index + 1} 頁，請稍候…")
        cleaned = prepare_scanned_page_image(image)
        prepared = _prepare_ocr_image(cleaned)
        try:
            pdf_bytes = pytesseract.image_to_pdf_or_hocr(
                prepared,
                extension="pdf",
                lang=tess_lang,
                config=tess_config,
            )
        finally:
            if prepared is not cleaned:
                prepared.close()
            if cleaned is not image:
                cleaned.close()
            image.close()
        reader = PdfReader(io.BytesIO(pdf_bytes))
        for page in reader.pages:
            writer.add_page(page)
    if count == 0:
        raise ValueError("沒有可整理的頁面。")
    if progress:
        progress(count, count, f"掃描件整理完成 {count} 頁")
    write_pdf(writer, target)
    return count


def compress_pdf_to_size(
    source: Path,
    target: Path,
    max_bytes: int,
    password: str = "",
) -> tuple[int, int, str]:
    """Compress until the file is at or under ``max_bytes`` when possible."""

    if max_bytes < 20_000:
        raise ValueError("目標大小請至少 20 KB。")
    old_size = source.stat().st_size
    if old_size <= max_bytes:
        _copy_pdf_file(source, target)
        return old_size, target.stat().st_size, "already-small"

    attempts: list[tuple[str, tuple[int, int] | None]] = [
        ("basic", None),
        ("advanced-150-70", (150, 70)),
        ("advanced-120-55", (120, 55)),
        ("advanced-100-45", (100, 45)),
        ("advanced-72-35", (72, 35)),
        ("advanced-72-22", (72, 22)),
    ]
    best_method = "original"
    best_size = old_size
    best_data = source.read_bytes()
    with tempfile.TemporaryDirectory() as temp_dir:
        for method, params in attempts:
            candidate = Path(temp_dir) / f"{method}.pdf"
            try:
                if params is None:
                    compress_pdf(source, candidate, password)
                else:
                    compress_pdf_advanced(
                        source,
                        candidate,
                        image_dpi=params[0],
                        jpeg_quality=params[1],
                        password=password,
                    )
            except Exception:
                continue
            if not candidate.is_file():
                continue
            size = candidate.stat().st_size
            if size < best_size:
                best_size = size
                best_method = method
                best_data = candidate.read_bytes()
            if size <= max_bytes:
                break
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_bytes(best_data)
    return old_size, best_size, best_method


def sanitize_pdf_for_external(source: Path, target: Path, password: str = "") -> dict[str, int]:
    """Remove JavaScript, embedded files, annotations, and metadata for external send."""

    document = _open_pymupdf_document(source, password)
    annotation_count = 0
    try:
        for page in document:
            annotation_count += len(list(page.annots() or []))
        embedded_count = int(document.embfile_count() or 0)
        scrub = getattr(document, "scrub", None)
        if callable(scrub):
            scrub(
                attached_files=True,
                clean_pages=True,
                embedded_files=True,
                hidden_text=True,
                javascript=True,
                metadata=True,
                redactions=True,
                redact_images=0,
                remove_links=False,
                reset_fields=True,
                reset_responses=True,
                thumbnails=True,
                xml_metadata=True,
            )
        else:
            for page in document:
                for annot in list(page.annots() or []):
                    page.delete_annot(annot)
            while document.embfile_count():
                document.embfile_del(0)
            document.set_metadata({})
            document.del_xml_metadata()
        for page in document:
            for annot in list(page.annots() or []):
                try:
                    page.delete_annot(annot)
                except Exception:
                    continue
        document.save(str(target), garbage=4, deflate=True, clean=True)
    finally:
        document.close()
    return {"annotations": annotation_count, "embedded_files": embedded_count}


def flatten_or_strip_annotations(
    source: Path,
    target: Path,
    mode: str = "flatten",
    password: str = "",
) -> int:
    normalized = (mode or "flatten").strip().lower()
    if normalized not in {"flatten", "strip"}:
        raise ValueError("註解處理請選擇壓平或清除。")
    document = _open_pymupdf_document(source, password)
    count = 0
    try:
        for page in document:
            count += len(list(page.annots() or []))
        if normalized == "flatten" and count:
            document.bake(annots=True, widgets=False)
        elif normalized == "strip":
            for page in document:
                for annot in list(page.annots() or []):
                    page.delete_annot(annot)
        document.save(str(target), garbage=4, deflate=True)
    finally:
        document.close()
    return count


def _pptx_available() -> bool:
    try:
        import pptx  # noqa: F401

        return True
    except Exception:
        return False


def _pdf_to_pptx_via_images(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    dpi: int = 150,
    progress=None,
) -> tuple[int, str]:
    if not _pptx_available():
        raise ValueError("PDF 轉 PowerPoint 需要 Python 套件 python-pptx。")
    from pptx import Presentation
    from pptx.util import Emu

    pages = list(iter_pdf_page_images(source, password, pages_spec, dpi))
    if not pages:
        raise ValueError("沒有可轉換的頁面。")
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    first_image = pages[0][1]
    presentation.slide_width = Emu(int(first_image.width / float(dpi) * 914400))
    presentation.slide_height = Emu(int(first_image.height / float(dpi) * 914400))
    with tempfile.TemporaryDirectory() as temp_dir:
        for index, (page_index, image) in enumerate(pages, start=1):
            if progress:
                progress(index - 1, len(pages), f"正在匯出第 {page_index + 1} 頁到投影片…")
            slide = presentation.slides.add_slide(blank)
            image_path = Path(temp_dir) / f"page-{page_index + 1}.png"
            image.convert("RGB").save(image_path, format="PNG")
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            image.close()
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))
    if progress:
        progress(len(pages), len(pages), f"已轉成 {len(pages)} 張投影片")
    return len(pages), "images"


def pdf_to_pptx(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    dpi: int = 150,
    progress=None,
) -> tuple[int, str]:
    unlocked = _write_unlocked_pdf_copy(source, password)
    try:
        if not (pages_spec or "").strip():
            ok, method = _pdf_to_office_via_libreoffice(unlocked, target, "pptx")
            if ok:
                page_count = len(open_reader(source, password).pages)
                return page_count, method
        return _pdf_to_pptx_via_images(source, target, password, pages_spec, dpi, progress)
    finally:
        if unlocked != source:
            unlocked.unlink(missing_ok=True)


def signature_coverage_is_intact(file_size: int, byte_range: list[int]) -> bool:
    if len(byte_range) < 4:
        return False
    start1, length1, start2, length2 = (int(item) for item in byte_range[:4])
    if start1 != 0 or length1 < 0 or length2 < 0 or start2 < start1 + length1:
        return False
    return start2 + length2 == int(file_size)


def _pdf_xref_value_text(value) -> str:
    if isinstance(value, tuple) and len(value) == 2:
        return str(value[1] or "")
    return str(value or "")


def _signature_info_from_xref(
    document,
    xref: int,
    raw: bytes,
    page_index: int,
    field_name: str,
) -> SignatureInfo:
    byte_range_text = _pdf_xref_value_text(document.xref_get_key(xref, "ByteRange"))
    contents_text = _pdf_xref_value_text(document.xref_get_key(xref, "Contents"))
    signer = _pdf_xref_value_text(document.xref_get_key(xref, "Name")).strip(" /")
    signed_at = _pdf_xref_value_text(document.xref_get_key(xref, "M")).strip(" /")
    numbers = [int(item) for item in re.findall(r"-?\d+", byte_range_text)]
    has_contents = bool(re.search(r"[0-9A-Fa-f]{16,}", contents_text))
    if not numbers:
        return SignatureInfo(
            field_name=field_name,
            page_index=page_index,
            signer=signer,
            signed_at=signed_at,
            has_contents=has_contents,
            intact=False,
            detail="尚未簽署，或沒有 ByteRange。",
        )
    intact = has_contents and signature_coverage_is_intact(len(raw), numbers)
    if intact:
        detail = "簽署範圍覆蓋整個檔案，簽完後未再附加內容。"
    elif has_contents:
        detail = "有簽章內容，但檔案在簽署範圍之後可能被改過或追加過。"
    else:
        detail = "簽章欄位不完整。"
    return SignatureInfo(
        field_name=field_name,
        page_index=page_index,
        signer=signer,
        signed_at=signed_at,
        has_contents=has_contents,
        intact=intact,
        detail=detail,
    )


def inspect_pdf_signatures(source: Path, password: str = "") -> list[SignatureInfo]:
    document = _open_pymupdf_document(source, password)
    raw = source.read_bytes()
    infos: list[SignatureInfo] = []
    seen: set[int] = set()
    try:
        for page_index, page in enumerate(document):
            for widget in page.widgets() or []:
                if _form_widget_type_name(widget.field_type) != "signature":
                    continue
                xref = int(getattr(widget, "xref", 0) or 0)
                if xref in seen:
                    continue
                seen.add(xref)
                infos.append(
                    _signature_info_from_xref(
                        document,
                        xref,
                        raw,
                        page_index,
                        str(widget.field_name or f"Signature{len(infos) + 1}"),
                    )
                )
        for xref in range(1, document.xref_length()):
            if xref in seen:
                continue
            type_text = _pdf_xref_value_text(document.xref_get_key(xref, "Type"))
            if type_text != "/Sig":
                continue
            seen.add(xref)
            infos.append(
                _signature_info_from_xref(
                    document,
                    xref,
                    raw,
                    -1,
                    f"Signature{len(infos) + 1}",
                )
            )
    finally:
        document.close()
    return infos


def verify_pdf_signatures(source: Path, target: Path, password: str = "") -> int:
    infos = inspect_pdf_signatures(source, password)
    lines = [
        "Victor PDF Tools Box - 數位簽章檢查",
        f"檔案：{source.name}",
        f"簽章數：{len(infos)}",
        "",
    ]
    if not infos:
        lines.append("這份 PDF 沒有憑證數位簽章（影像簽名／圖章不算）。")
    for index, item in enumerate(infos, start=1):
        page_label = "文件層" if item.page_index < 0 else str(item.page_index + 1)
        lines.extend(
            [
                f"--- 簽章 {index} ---",
                f"欄位：{item.field_name or '（未命名）'}",
                f"頁碼：{page_label}",
                f"簽署人：{item.signer or '（未標示）'}",
                f"時間：{item.signed_at or '（未標示）'}",
                f"有簽章內容：{'是' if item.has_contents else '否'}",
                f"簽署後未被追加修改：{'是' if item.intact else '否'}",
                f"說明：{item.detail}",
                "",
            ]
        )
    target.write_text("\n".join(lines), encoding="utf-8")
    return len(infos)


def _unique_attachment_name(name: str, used: set[str]) -> str:
    stem = Path(name).name.strip() or "attachment.bin"
    candidate = safe_output_name(stem)
    if candidate not in used:
        used.add(candidate)
        return candidate
    base = Path(candidate).stem
    suffix = Path(candidate).suffix or ".bin"
    index = 2
    while True:
        renamed = safe_output_name(f"{base}-{index}{suffix}")
        if renamed not in used:
            used.add(renamed)
            return renamed
        index += 1


def extract_pdf_attachments(source: Path, target_zip: Path, password: str = "") -> int:
    document = _open_pymupdf_document(source, password)
    extracted = 0
    used: set[str] = set()
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            folder = Path(temp_dir)
            for index in range(int(document.embfile_count() or 0)):
                info = document.embfile_info(index) or {}
                name = str(info.get("filename") or info.get("name") or f"embedded-{index + 1}.bin")
                data = document.embfile_get(index)
                if not data:
                    continue
                path = folder / _unique_attachment_name(name, used)
                path.write_bytes(bytes(data))
                extracted += 1
            for page_index, page in enumerate(document):
                for annot in list(page.annots() or []):
                    kind = annot.type[1] if isinstance(annot.type, tuple) else str(annot.type)
                    if kind != "FileAttachment":
                        continue
                    payload = None
                    getter = getattr(annot, "get_file", None)
                    if callable(getter):
                        payload = getter()
                    data = b""
                    name = f"attachment-p{page_index + 1}.bin"
                    if isinstance(payload, dict):
                        data = bytes(payload.get("file") or payload.get("content") or b"")
                        name = str(payload.get("filename") or name)
                    elif isinstance(payload, (bytes, bytearray)):
                        data = bytes(payload)
                    if not data:
                        continue
                    path = folder / _unique_attachment_name(name, used)
                    path.write_bytes(data)
                    extracted += 1
            if extracted == 0:
                raise ValueError("這份 PDF 沒有可抽出的內嵌附件。")
            with zipfile.ZipFile(target_zip, "w", zipfile.ZIP_DEFLATED) as archive:
                for item in sorted(folder.iterdir()):
                    archive.write(item, item.name)
    finally:
        document.close()
    return extracted


def _extracted_image_extension(ext: str) -> str:
    normalized = (ext or "bin").lower().lstrip(".")
    if normalized in {"jpeg", "jpg", "jpe"}:
        return "jpg"
    if normalized in {"png", "jp2", "jpx", "j2k", "tif", "tiff", "bmp", "gif", "webp"}:
        return "jp2" if normalized in {"jpx", "j2k"} else normalized
    if normalized in {"jb2", "jbig2"}:
        return "jb2"
    return normalized or "bin"


def _compose_image_with_smask(image_bytes: bytes, document, smask_xref: int) -> Image.Image | None:
    try:
        mask_info = document.extract_image(smask_xref)
        mask_bytes = bytes((mask_info or {}).get("image") or b"")
        if not mask_bytes:
            return None
        with Image.open(io.BytesIO(image_bytes)) as image, Image.open(io.BytesIO(mask_bytes)) as mask:
            rgba = image.convert("RGBA")
            alpha = mask.convert("L")
            if alpha.size != rgba.size:
                alpha = alpha.resize(rgba.size, Image.Resampling.BILINEAR)
            rgba.putalpha(alpha)
            return rgba.copy()
    except Exception:
        return None


def _extracted_image_payload(document, xref: int) -> tuple[Image.Image, bytes, str] | None:
    try:
        info = document.extract_image(xref)
    except Exception:
        return None
    data = bytes((info or {}).get("image") or b"")
    if not data:
        return None
    smask = int((info or {}).get("smask") or 0)
    ext = _extracted_image_extension(str((info or {}).get("ext") or ""))
    if smask:
        composed = _compose_image_with_smask(data, document, smask)
        if composed is not None:
            buffer = io.BytesIO()
            composed.save(buffer, format="PNG")
            return composed, buffer.getvalue(), "png"
    try:
        image = Image.open(io.BytesIO(data)).convert("RGB")
        image.load()
    except Exception:
        return None
    return image, data, ext


def _image_fingerprint(image: Image.Image) -> bytes:
    return image.convert("L").resize((12, 12), Image.Resampling.BILINEAR).tobytes()


def _luminance_stats(image: Image.Image) -> tuple[float, float, int]:
    sample = image.convert("RGB").resize((48, 48), Image.Resampling.BILINEAR)
    pixels = list(sample.getdata())
    lum = [0.299 * r + 0.587 * g + 0.114 * b for r, g, b in pixels]
    mean = sum(lum) / len(lum)
    variance = sum((value - mean) ** 2 for value in lum) / len(lum)
    stddev = variance ** 0.5
    edges = 0.0
    width = 48
    for y in range(47):
        for x in range(47):
            index = y * width + x
            here = lum[index]
            edges += abs(here - lum[index + 1]) + abs(here - lum[index + width])
    edge_mean = edges / (47 * 47 * 2)
    quantized = {(r // 16, g // 16, b // 16) for r, g, b in pixels}
    return stddev, edge_mean, len(quantized)


def _displayed_page_fraction(
    rects: list[tuple[float, float, float, float]] | None,
    page_size: tuple[float, float] | None,
) -> float:
    if not rects or not page_size:
        return 1.0
    page_w, page_h = page_size
    page_area = page_w * page_h
    if page_area <= 0:
        return 1.0
    area = 0.0
    for x0, y0, x1, y1 in rects:
        area += abs(x1 - x0) * abs(y1 - y0)
    return area / page_area


def embedded_image_is_useful(
    image: Image.Image,
    *,
    quality: str = "main",
    rects: list[tuple[float, float, float, float]] | None = None,
    page_size: tuple[float, float] | None = None,
) -> bool:
    width, height = image.size
    if width < 8 or height < 8:
        return False
    if quality != "main":
        return True
    min_side = min(width, height)
    max_side = max(width, height)
    if min_side < 48 or width * height < 80 * 80:
        return False
    if max_side / min_side > 10:
        return False
    fraction = _displayed_page_fraction(rects, page_size)
    if fraction < 0.012 and min_side < 140:
        return False
    stddev, edge_mean, unique_bins = _luminance_stats(image)
    if unique_bins <= 3 or (stddev < 12 and unique_bins <= 6):
        return False
    if edge_mean < 6.5 and unique_bins <= 28:
        return False
    if fraction > 0.55 and edge_mean < 10 and unique_bins <= 20:
        return False
    return True


def _save_extracted_pdf_image(document, xref: int, dest_without_suffix: Path) -> Path | None:
    payload = _extracted_image_payload(document, xref)
    if payload is None:
        return None
    _image, data, ext = payload
    path = dest_without_suffix.with_suffix(f".{ext}")
    path.write_bytes(data)
    return path


def _iter_embedded_image_xrefs(document, page_indexes: list[int]):
    seen_xrefs: set[int] = set()
    for page_index in page_indexes:
        page = document[page_index]
        images = list(page.get_images(full=True) or [])
        smask_xrefs = {int(item[1]) for item in images if len(item) > 1 and item[1]}
        for item in images:
            xref = int(item[0])
            if xref in seen_xrefs or xref in smask_xrefs:
                continue
            width = int(item[2]) if len(item) > 2 else 0
            height = int(item[3]) if len(item) > 3 else 0
            if width < 8 or height < 8:
                continue
            seen_xrefs.add(xref)
            yield page_index, xref, width, height


def _embedded_image_rects(page, xref: int) -> list[tuple[float, float, float, float]]:
    rects: list[tuple[float, float, float, float]] = []
    try:
        for rect in page.get_image_rects(xref) or []:
            rects.append((float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)))
    except Exception:
        pass
    return rects


def list_pdf_embedded_images(
    source: Path,
    password: str = "",
    pages_spec: str = "",
    quality: str = "main",
) -> list[dict]:
    page_indexes = selected_page_indexes(source, password, pages_spec)
    if not page_indexes:
        return []
    document = _open_pymupdf_document(source, password)
    found: list[dict] = []
    fingerprints: set[bytes] = set()
    try:
        for page_index, xref, width, height in _iter_embedded_image_xrefs(document, page_indexes):
            page = document[page_index]
            rects = _embedded_image_rects(page, xref)
            payload = _extracted_image_payload(document, xref)
            if payload is None:
                continue
            image, _data, _ext = payload
            page_size = (float(page.rect.width), float(page.rect.height))
            if not embedded_image_is_useful(image, quality=quality, rects=rects, page_size=page_size):
                continue
            fingerprint = _image_fingerprint(image)
            if fingerprint in fingerprints:
                continue
            fingerprints.add(fingerprint)
            found.append(
                {
                    "page_index": page_index,
                    "xref": xref,
                    "width": width,
                    "height": height,
                    "rects": rects,
                }
            )
    finally:
        document.close()
    return found


def extract_pdf_images(
    source: Path,
    target: Path,
    password: str = "",
    pages_spec: str = "",
    quality: str = "main",
) -> int:
    """Save embedded PDF image XObjects (not full-page rasterization)."""

    page_indexes = selected_page_indexes(source, password, pages_spec)
    if not page_indexes:
        raise ValueError("沒有可處理的頁面。")
    document = _open_pymupdf_document(source, password)
    extracted = 0
    used: set[str] = set()
    fingerprints: set[bytes] = set()
    stem = _WINDOWS_FORBIDDEN_NAME.sub("_", source.stem).strip(" .") or "output"
    page_serial: dict[int, int] = {}
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            folder = Path(temp_dir)
            for page_index, xref, _width, _height in _iter_embedded_image_xrefs(document, page_indexes):
                page = document[page_index]
                rects = _embedded_image_rects(page, xref)
                payload = _extracted_image_payload(document, xref)
                if payload is None:
                    continue
                image, data, ext = payload
                page_size = (float(page.rect.width), float(page.rect.height))
                if not embedded_image_is_useful(image, quality=quality, rects=rects, page_size=page_size):
                    continue
                fingerprint = _image_fingerprint(image)
                if fingerprint in fingerprints:
                    continue
                fingerprints.add(fingerprint)
                page_serial[page_index] = page_serial.get(page_index, 0) + 1
                dest = folder / f"{stem}-p{page_index + 1:04d}-{page_serial[page_index]:02d}.{ext}"
                candidate = dest.name
                if candidate in used:
                    candidate = _unique_attachment_name(candidate, used)
                else:
                    used.add(candidate)
                dest = folder / candidate
                dest.write_bytes(data)
                extracted += 1
            if extracted == 0:
                if quality == "main":
                    raise ValueError(
                        "沒有符合的主要圖片（已略過小圖示、細線和漸層背景）。"
                        "可改選「全部內嵌圖」，或這頁圖是向量／掃描件請用「PDF 轉圖片」。"
                    )
                raise ValueError(
                    "這份 PDF 沒有可抽出的內嵌圖片。若圖是整頁掃描或向量繪製，請改用「PDF 轉圖片」。"
                )
            if target.suffix.lower() == ".zip":
                with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as archive:
                    for item in sorted(folder.iterdir()):
                        archive.write(item, item.name)
            else:
                target.mkdir(parents=True, exist_ok=True)
                for item in sorted(folder.iterdir()):
                    shutil.copy2(item, target / item.name)
    finally:
        document.close()
    return extracted


def _pdf_page_group_bytes(reader: PdfReader, page_indexes: list[int]) -> bytes:
    writer = PdfWriter()
    for index in page_indexes:
        writer.add_page(reader.pages[index])
    buffer = io.BytesIO()
    writer.write(buffer)
    return buffer.getvalue()


def split_pdf_by_size(
    source: Path,
    target_zip: Path,
    max_bytes: int,
    password: str = "",
) -> int:
    if max_bytes < 1:
        raise ValueError("每份大小必須大於 0。")
    reader = open_reader(source, password)
    page_count = len(reader.pages)
    if page_count == 0:
        raise ValueError("PDF 沒有頁面可拆分。")
    groups: list[list[int]] = []
    current: list[int] = []
    for index in range(page_count):
        trial = current + [index]
        data = _pdf_page_group_bytes(reader, trial)
        if current and len(data) > max_bytes:
            groups.append(current)
            current = [index]
        else:
            current = trial
    if current:
        groups.append(current)
    stem = safe_output_name(source.stem)
    with tempfile.TemporaryDirectory() as temp_dir:
        generated: list[Path] = []
        for part_index, pages in enumerate(groups, start=1):
            part = Path(temp_dir) / f"{stem}-part-{part_index:03d}.pdf"
            part.write_bytes(_pdf_page_group_bytes(reader, pages))
            generated.append(part)
        with zipfile.ZipFile(target_zip, "w", zipfile.ZIP_DEFLATED) as archive:
            for item in generated:
                archive.write(item, item.name)
    return len(generated)
